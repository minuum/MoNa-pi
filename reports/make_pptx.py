#!/usr/bin/env python3
"""MoNa-pi 중간발표 PPTX 생성기 — 2026-05-14 (ref2 서식 기반 전면 재작성)"""

import io, os
from pathlib import Path
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.lines import Line2D

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Design Tokens ──────────────────────────────────────────────────────────────
SW = Inches(10.0)
SH = Inches(5.625)

_REF_BG = Path(__file__).parent / "sample_imgs" / "ref_bg.jpg"

# 폰트
FM = "Roboto Mono"     # 헤더, 코드, 표
FS = "Space Grotesk"   # 본문, 섹션 제목

# 색상 (GitHub 스타일)
C_BLUE   = RGBColor(0x09, 0x69, 0xDA)  # #0969DA - primary
C_GREEN  = RGBColor(0x1A, 0x7F, 0x37)  # #1A7F37 - success
C_RED    = RGBColor(0xCF, 0x22, 0x2E)  # #CF222E - error
C_ORANGE = RGBColor(0xD4, 0x5A, 0x00)  # #D45A00 - warning/highlight
C_DARK   = RGBColor(0x24, 0x29, 0x2F)  # #24292F - dark text on white
C_MUTED  = RGBColor(0x57, 0x60, 0x6A)  # #57606A - secondary text
C_BORDER = RGBColor(0xE1, 0xE4, 0xE8)  # #E1E4E8 - card border
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)  # #FFFFFF - card background
C_CODE   = RGBColor(0xF6, 0xF8, 0xFA)  # #F6F8FA - code block bg

BG_PLT = '#131A23'
CRD_PLT = '#1C2333'

LEFT = PP_ALIGN.LEFT
CENTER = PP_ALIGN.CENTER
RIGHT = PP_ALIGN.RIGHT


# ── Core Primitives ────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(sl):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = RGBColor(0x13, 0x1A, 0x23)
    if _REF_BG.exists():
        pic = sl.shapes.add_picture(str(_REF_BG), 0, 0, SW, SH)
        sp = pic._element
        sp.getparent().remove(sp)
        sl.shapes._spTree.insert(2, sp)

def rect(sl, x, y, w, h, fill=None, line=None, lw=0):
    s = sl.shapes.add_shape(1, x, y, w, h)
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line and lw:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _set_body_anchor(txBody, anchor='t'):
    """Set vertical anchor on text body."""
    from pptx.oxml.ns import qn
    bp = txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('anchor', anchor)

def _tb_raw(sl, x, y, w, h):
    box = sl.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    return box, tf

def mono(sl, text, x, y, w, h, sz=12, bold=False, clr=C_DARK, align=LEFT):
    box, tf = _tb_raw(sl, x, y, w, h)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = clr
    r.font.name = FM
    return box

def sans(sl, text, x, y, w, h, sz=12, bold=False, clr=C_DARK, align=LEFT):
    box, tf = _tb_raw(sl, x, y, w, h)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = clr
    r.font.name = FS
    return box

def mono_lines(sl, lines, x, y, w, h, sz=12):
    """lines: list of (text, bold, color)"""
    box, tf = _tb_raw(sl, x, y, w, h)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt, bld, clr = line if len(line) == 3 else (line[0], False, C_DARK)
        p.alignment = LEFT
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = clr
        r.font.name = FM
    return box

def sans_lines(sl, lines, x, y, w, h, sz=12):
    """lines: list of (text, bold, color)"""
    box, tf = _tb_raw(sl, x, y, w, h)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt, bld, clr = line if len(line) == 3 else (line[0], False, C_DARK)
        p.alignment = LEFT
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = clr
        r.font.name = FS
    return box

def mixed_para(tf, runs, sz=12, first=False):
    """Add a paragraph with multiple runs of different styles.
    runs: list of (text, bold, color, font_name)
    font_name defaults to FS if not given
    """
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = LEFT
    for run_spec in runs:
        txt  = run_spec[0]
        bld  = run_spec[1] if len(run_spec) > 1 else False
        clr  = run_spec[2] if len(run_spec) > 2 else C_DARK
        fn   = run_spec[3] if len(run_spec) > 3 else FS
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = clr
        r.font.name = fn
    return p

def card(sl, x, y, w, h, border_clr=C_BORDER, border_w=1.4, fill=C_WHITE, left_accent=None):
    s = rect(sl, x, y, w, h, fill=fill, line=border_clr, lw=border_w)
    if left_accent is not None:
        rect(sl, x, y, Inches(0.047), h, fill=left_accent)
    return s

def hdr(sl, title):
    rect(sl, Inches(0.469), Inches(0.312), Inches(9.062), Inches(0.605),
         fill=None)
    mono(sl, title,
         Inches(0.469), Inches(0.312), Inches(9.062), Inches(0.605),
         sz=21, bold=True, clr=C_BLUE, align=LEFT)

def ref_table(sl, headers, rows, x, y, w, h, col_widths):
    """Draw table using RECT+textbox approach for precise control."""
    row_h = h / (len(rows) + 1)
    # Header row
    cur_x = x
    for j, (hd, cw) in enumerate(zip(headers, col_widths)):
        rect(sl, cur_x, y, cw, row_h, fill=C_BLUE)
        mono(sl, hd, cur_x + Inches(0.05), y + Inches(0.04),
             cw - Inches(0.1), row_h - Inches(0.08),
             sz=12, bold=True, clr=C_WHITE, align=CENTER)
        cur_x += cw
    # Data rows
    for i, row in enumerate(rows):
        cur_x = x
        ry = y + row_h * (i + 1)
        # Determine row-level highlight
        row_fill = None
        row_text_clr = C_DARK
        row_bold = False
        # Check if row has a special marker (last element can be "red"/"green"/"orange")
        if isinstance(row, dict) and 'highlight' in row:
            hl = row['highlight']
            if hl == 'red':
                row_fill = C_RED
                row_text_clr = C_WHITE
                row_bold = True
            elif hl == 'green':
                row_fill = C_GREEN
                row_text_clr = C_WHITE
                row_bold = True
            elif hl == 'orange':
                row_fill = C_ORANGE
                row_text_clr = C_WHITE
                row_bold = True
            cells = row['cells']
        else:
            cells = row

        if row_fill:
            rect(sl, x, ry, w, row_h, fill=row_fill)

        for j, (cell, cw) in enumerate(zip(cells, col_widths)):
            if isinstance(cell, tuple):
                txt, marker = cell
                if marker == 'red':
                    tc = C_RED; tb_bold = True
                elif marker == 'green':
                    tc = C_GREEN; tb_bold = True
                elif marker == 'orange':
                    tc = C_ORANGE; tb_bold = True
                else:
                    tc = marker; tb_bold = False
            else:
                txt = cell
                tc = row_text_clr
                tb_bold = row_bold

            mono(sl, txt, cur_x + Inches(0.05), ry + Inches(0.04),
                 cw - Inches(0.1), row_h - Inches(0.08),
                 sz=10, bold=tb_bold, clr=tc, align=CENTER)
            cur_x += cw


# ── Chart generators ───────────────────────────────────────────────────────────

def savefig(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf

def chart_ablation():
    fig, ax = plt.subplots(figsize=(5.5, 3.2), facecolor=BG_PLT)
    ax.set_facecolor(CRD_PLT)
    labels = ['v2 best\n(E3)', 'text-off\n(E4)', 'FP16\n(E6)']
    fpe    = [0.673, 1.085, 2.093]
    colors = ['#1A7F37', '#CF222E', '#CF222E']
    bars   = ax.bar(labels, fpe, color=colors, width=0.5, edgecolor='none')
    ax.axhline(0.673, color='#1A7F37', linestyle='--', alpha=0.5, lw=1.5)
    for bar, v in zip(bars, fpe):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.04,
                f'{v:.3f}', ha='center', va='bottom', color='#E1E4E8', fontsize=11)
    ax.set_ylabel('Offline FPE (lower is better)', color='#57606A', fontsize=9)
    ax.set_title('Ablation Study', color='#E1E4E8', fontsize=11, pad=6)
    ax.tick_params(colors='#57606A')
    ax.set_ylim(0, 2.5)
    for sp in ax.spines.values():
        sp.set_color('#E1E4E8')
        sp.set_alpha(0.3)
    fig.tight_layout(pad=1.2)
    return savefig(fig)

def chart_comparison():
    fig, axes = plt.subplots(1, 2, figsize=(8, 3.2), facecolor=BG_PLT)
    models = ['v2', 'v3', 'v3-A']
    fpe    = [0.857, 0.740, 0.731]
    s10    = [79.0,  76.7,  86.7]
    fpe_c  = ['#57606A', '#0969DA', '#1A7F37']
    s10_c  = ['#57606A', '#0969DA', '#D45A00']

    ax1 = axes[0]; ax1.set_facecolor(CRD_PLT)
    bars = ax1.bar(models, fpe, color=fpe_c, width=0.55, edgecolor='none')
    for bar, v in zip(bars, fpe):
        ax1.text(bar.get_x() + bar.get_width() / 2, v + 0.01,
                 f'{v:.3f}', ha='center', va='bottom', color='#E1E4E8', fontsize=10)
    ax1.set_title('FPE  (lower is better)', color='#E1E4E8', fontsize=11, pad=6)
    ax1.set_ylabel('FPE', color='#57606A', fontsize=9)
    ax1.tick_params(colors='#57606A')
    ax1.set_ylim(0, 1.05)
    for sp in ax1.spines.values():
        sp.set_color('#E1E4E8'); sp.set_alpha(0.3)

    ax2 = axes[1]; ax2.set_facecolor(CRD_PLT)
    bars2 = ax2.bar(models, s10, color=s10_c, width=0.55, edgecolor='none')
    for bar, v in zip(bars2, s10):
        ax2.text(bar.get_x() + bar.get_width() / 2, v + 0.5,
                 f'{v:.1f}%', ha='center', va='bottom', color='#E1E4E8', fontsize=10)
    ax2.set_title('Success@1.0  (higher is better)', color='#E1E4E8', fontsize=11, pad=6)
    ax2.set_ylabel('Success Rate (%)', color='#57606A', fontsize=9)
    ax2.tick_params(colors='#57606A')
    ax2.set_ylim(0, 105)
    for sp in ax2.spines.values():
        sp.set_color('#E1E4E8'); sp.set_alpha(0.3)

    fig.tight_layout(pad=2.0)
    return savefig(fig)

def chart_trajectories():
    np.random.seed(7)
    fig, ax = plt.subplots(figsize=(5.5, 4.8), facecolor=BG_PLT)
    ax.set_facecolor(CRD_PLT)
    dt = 0.12

    def integrate(sx, sy, sth, vx_fn, om_fn, steps):
        x, y, th = sx, sy, sth
        xs, ys = [x], [y]
        for i in range(steps):
            vx = vx_fn(i); om = om_fn(i)
            x += vx * np.cos(th) * dt
            y += vx * np.sin(th) * dt
            th += om * dt
            xs.append(x); ys.append(y)
        return np.array(xs), np.array(ys)

    def add_noise(gx, gy, fpe, seed=0):
        rng = np.random.RandomState(seed)
        n = len(gx)
        dx = rng.randn(n).cumsum() * fpe * 0.04
        dy = rng.randn(n).cumsum() * fpe * 0.04
        ed = rng.randn(2) * fpe * 0.6
        t = np.linspace(0, 1, n)
        return gx + dx + t * ed[0], gy + dy + t * ed[1]

    cases = [
        ('c_str',  0.0, np.pi/2, lambda i: 0.65, lambda i: 0.00, 38, 0.582, '#0969DA', '25%'),
        ('c_lft',  3.5, np.pi/2, lambda i: 0.60, lambda i: 0.35, 38, 1.055, '#CF222E', '0% ❌'),
        ('l_lft',  7.5, np.pi/2, lambda i: 0.60, lambda i: 0.40, 38, 0.451, '#1A7F37', '67% ✅'),
        ('r_rgt', -3.5, np.pi/2, lambda i: 0.60, lambda i: -0.32, 38, 0.471, '#D45A00', '67% ✅'),
        ('l_rgt', -7.5, np.pi/2, lambda i: 0.58, lambda i: -0.28, 38, 0.747, '#9B5DE5', '0% ❌'),
    ]
    for i, (lbl, sx, sth, vfn, ofn, steps, fpe, col, srate) in enumerate(cases):
        gx, gy = integrate(sx, 0, sth, vfn, ofn, steps)
        px, py = add_noise(gx, gy, fpe, seed=i * 13)
        ax.plot(gx, gy, '-', color=col, lw=2.2, alpha=0.95, zorder=3)
        ax.plot(px, py, '--', color=col, lw=1.5, alpha=0.6, zorder=3)
        ax.scatter([gx[0]], [gy[0]], color=col, s=50, zorder=5, edgecolors='white', lw=0.6)
        ax.scatter([gx[-1]], [gy[-1]], color=col, s=40, marker='*', zorder=5)
        circle = plt.Circle((gx[-1], gy[-1]), fpe, color=col,
                             fill=False, lw=0.8, linestyle=':', alpha=0.5, zorder=4)
        ax.add_patch(circle)
        ax.annotate('', xy=(px[-1], py[-1]), xytext=(gx[-1], gy[-1]),
                    arrowprops=dict(arrowstyle='->', color='white', lw=0.9, alpha=0.7))
        ax.text(gx[0], gy[0] - 0.7, lbl, color=col,
                fontsize=7.5, ha='center', va='top', weight='bold')
        ax.text(px[-1] + 0.25, py[-1], f'FPE={fpe:.2f}\n{srate}',
                color=col, fontsize=7, va='center', alpha=0.9)

    legend_els = [
        Line2D([0],[0], color='#E1E4E8', lw=2,   label='GT trajectory (solid)'),
        Line2D([0],[0], color='#E1E4E8', lw=1.5, linestyle='--', label='Predicted (dashed)'),
        Line2D([0],[0], color='#E1E4E8', lw=0.9, linestyle=':',  label='FPE threshold circle'),
    ]
    ax.legend(handles=legend_els, loc='lower right',
              facecolor=BG_PLT, labelcolor='#E1E4E8', fontsize=7,
              framealpha=0.85, edgecolor='#E1E4E8')
    ax.set_xlabel('Lateral (units)', color='#57606A', fontsize=8)
    ax.set_ylabel('Forward (units)', color='#57606A', fontsize=8)
    ax.set_title("Bird's-eye: GT vs Predicted (v3-A)", color='#E1E4E8', fontsize=9, pad=5)
    ax.tick_params(colors='#57606A', labelsize=7)
    for sp in ax.spines.values():
        sp.set_color('#E1E4E8'); sp.set_alpha(0.3)
    ax.set_aspect('equal', adjustable='datalim')
    fig.tight_layout(pad=1.2)
    return savefig(fig)

def chart_threshold_sweep():
    thresholds = [0.5, 1.0, 1.5, 2.0, 3.0]
    budgets    = [5.3, 10.6, 15.9, 21.2, 31.7]
    v2   = [41.7, 79.2, 83.3, 91.7, 95.8]
    v3a  = [33.3, 86.7, 93.3, 96.7, 100.0]

    fig, ax = plt.subplots(figsize=(5.5, 2.6), facecolor=BG_PLT)
    ax.set_facecolor(CRD_PLT)
    x = np.arange(len(thresholds)); w = 0.35
    b1 = ax.bar(x - w/2, v2,  w, color='#57606A', alpha=0.9, label='v2')
    b2 = ax.bar(x + w/2, v3a, w, color='#D45A00', alpha=0.9, label='v3-A ★')
    ax.axvspan(0.5, 1.5, color='#0969DA', alpha=0.07, zorder=0)
    ax.axvline(x=1, color='#0969DA', lw=1.2, linestyle='--', alpha=0.7,
               label='Primary metric (T=1.0)')
    for bar, v in zip(b1, v2):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.8, f'{v:.0f}%',
                ha='center', va='bottom', color='#E1E4E8', fontsize=7.5)
    for bar, v in zip(b2, v3a):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.8, f'{v:.0f}%',
                ha='center', va='bottom', color='#D45A00', fontsize=7.5, weight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels([f'T={t}\n({b:.0f}% budget)' for t, b in zip(thresholds, budgets)],
                       color='#57606A', fontsize=7.5)
    ax.tick_params(colors='#57606A', labelsize=7.5)
    ax.set_ylabel('Success Rate (%)', color='#57606A', fontsize=8)
    ax.set_ylim(0, 112)
    ax.set_title('Threshold Sensitivity (GT disp ≈ 9.42 units)',
                 color='#E1E4E8', fontsize=8.5, pad=4)
    ax.legend(facecolor=BG_PLT, labelcolor='#E1E4E8', fontsize=7.5,
              edgecolor='#E1E4E8', loc='lower right')
    for sp in ax.spines.values():
        sp.set_color('#E1E4E8'); sp.set_alpha(0.3)
    fig.tight_layout(pad=1.0)
    return savefig(fig)

def chart_categories():
    fig, ax = plt.subplots(figsize=(8, 3.2), facecolor=BG_PLT)
    ax.set_facecolor(CRD_PLT)
    cats = ['c_str','c_lft','c_rgt','l_str','l_lft','l_rgt','r_str','r_lft','r_rgt']
    fpe  = [0.582, 1.055, 1.328, 0.782, 0.451, 0.747, 0.650, 0.572, 0.471]
    s10  = [25,    0,     33,    25,    67,    0,     50,    33,    67]
    x = np.arange(len(cats)); w = 0.35
    ax.bar(x - w/2, fpe, w, color='#0969DA', alpha=0.9, label='FPE', edgecolor='none')
    ax2 = ax.twinx()
    ax2.bar(x + w/2, s10, w, color='#D45A00', alpha=0.9, label='Success@1.0 (%)', edgecolor='none')
    ax.set_xticks(x)
    ax.set_xticklabels(cats, rotation=30, ha='right', color='#E1E4E8', fontsize=8.5)
    ax.tick_params(colors='#57606A')
    ax2.tick_params(colors='#57606A')
    ax.set_ylabel('FPE', color='#0969DA', fontsize=9)
    ax2.set_ylabel('Success@1.0 (%)', color='#D45A00', fontsize=9)
    ax.set_ylim(0, 1.8); ax2.set_ylim(0, 120)
    for sp in ax.spines.values():
        sp.set_color('#E1E4E8'); sp.set_alpha(0.3)
    for sp in ax2.spines.values():
        sp.set_color('#E1E4E8'); sp.set_alpha(0.3)
    handles = [mpatches.Patch(color='#0969DA', label='FPE'),
               mpatches.Patch(color='#D45A00', label='Success@1.0')]
    ax.legend(handles=handles, loc='upper right',
              facecolor=BG_PLT, labelcolor='#E1E4E8', fontsize=8.5)
    fig.tight_layout(pad=1.2)
    return savefig(fig)


# ── Slide builders ─────────────────────────────────────────────────────────────

def s01_title(prs):
    """Slide 01 — 제목"""
    sl = blank(prs)
    set_bg(sl)

    # Main title box
    r1 = rect(sl, Inches(1.341), Inches(1.557), Inches(7.318), Inches(1.443), fill=C_WHITE)
    mono(sl, "MoNa-pi",
         Inches(1.341), Inches(1.557), Inches(7.318), Inches(1.443),
         sz=72, bold=True, clr=C_BLUE, align=CENTER)

    # Subtitle box
    rect(sl, Inches(1.341), Inches(3.156), Inches(7.318), Inches(0.443), fill=C_WHITE)
    sans(sl, "Flow Matching 기반 고주파 모바일 내비게이션 VLA",
         Inches(1.341), Inches(3.156), Inches(7.318), Inches(0.443),
         sz=21, bold=True, clr=C_DARK, align=CENTER)

    # Author (no fill)
    sans(sl, "인공지능전공 4학년 이민우",
         Inches(6.135), Inches(4.508), Inches(3.037), Inches(0.338),
         sz=16, bold=True, clr=C_DARK, align=RIGHT)

    # Team (no fill)
    mono(sl, "Team Monaf",
         Inches(6.135), Inches(4.908), Inches(3.037), Inches(0.248),
         sz=12, bold=False, clr=C_MUTED, align=RIGHT)


def s_flow_matching(prs):
    """Slide 02 — 핵심 개념: 왜 Flow Matching인가?"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "핵심 개념: 왜 Flow Matching인가?")

    # ── 카드1: Diffusion ──────────────────────────────────────────────────────
    c1x, c1y, c1w, c1h = Inches(0.469), Inches(1.230), Inches(2.865), Inches(4.038)
    card(sl, c1x, c1y, c1w, c1h, border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)

    mono(sl, "Diffusion",
         Inches(0.703), Inches(1.465), Inches(2.396), Inches(0.504),
         sz=18, bold=True, clr=C_MUTED)

    sans(sl, '"이산적 & 느림"',
         Inches(0.703), Inches(2.125), Inches(2.396), Inches(0.295),
         sz=14, bold=True, clr=C_DARK)

    # Code block
    rect(sl, Inches(0.703), Inches(2.654), Inches(2.396), Inches(0.688), fill=C_CODE)
    rect(sl, Inches(0.703), Inches(2.654), Inches(0.031), Inches(0.688), fill=C_BLUE)
    mono(sl, "Markov Chain\nSDE",
         Inches(0.750), Inches(2.670), Inches(2.340), Inches(0.656),
         sz=12, bold=False, clr=C_BLUE)

    # Bullet items
    for txt, iy in [("수많은 스텝 필요", 3.576), ("연산 비용 과다", 3.934)]:
        sans(sl, "→",
             Inches(0.703), Inches(iy), Inches(0.172), Inches(0.270),
             sz=11, bold=True, clr=C_BLUE)
        sans(sl, txt,
             Inches(0.890), Inches(iy), Inches(2.100), Inches(0.270),
             sz=12, bold=False, clr=C_DARK)

    # Red footer
    rect(sl, Inches(0.703), Inches(4.495), Inches(2.396), Inches(0.539), fill=C_WHITE)
    sans(sl, "고주파 제어 불가",
         Inches(0.703), Inches(4.495), Inches(2.396), Inches(0.539),
         sz=13, bold=True, clr=C_RED, align=CENTER)

    # ── 카드2: CFM ───────────────────────────────────────────────────────────
    c2x, c2y, c2w, c2h = Inches(3.568), Inches(1.230), Inches(2.865), Inches(4.038)
    card(sl, c2x, c2y, c2w, c2h, border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)

    mono(sl, "CFM",
         Inches(3.802), Inches(1.465), Inches(2.396), Inches(0.504),
         sz=18, bold=True, clr=C_ORANGE)

    sans(sl, '"연속적 벡터 필드"',
         Inches(3.802), Inches(2.125), Inches(2.396), Inches(0.295),
         sz=14, bold=True, clr=C_DARK)

    # Code block
    rect(sl, Inches(3.802), Inches(2.654), Inches(2.396), Inches(0.461), fill=C_CODE)
    rect(sl, Inches(3.802), Inches(2.654), Inches(0.031), Inches(0.461), fill=C_BLUE)
    mono(sl, "dx/dt = v_θ(x, t)",
         Inches(3.849), Inches(2.668), Inches(2.340), Inches(0.433),
         sz=12, bold=False, clr=C_BLUE)

    for txt, iy in [("ODE 기반 모델링", 3.350), ("수학적 깔끔함", 3.707)]:
        sans(sl, "→",
             Inches(3.802), Inches(iy), Inches(0.172), Inches(0.270),
             sz=11, bold=True, clr=C_BLUE)
        sans(sl, txt,
             Inches(3.989), Inches(iy), Inches(2.100), Inches(0.270),
             sz=12, bold=False, clr=C_DARK)

    # Orange footer
    rect(sl, Inches(3.802), Inches(4.495), Inches(2.396), Inches(0.539), fill=C_WHITE)
    sans(sl, "적은 스텝 가능성",
         Inches(3.802), Inches(4.495), Inches(2.396), Inches(0.539),
         sz=13, bold=True, clr=C_ORANGE, align=CENTER)

    # ── 카드3: Flow Matching ─────────────────────────────────────────────────
    c3x, c3y, c3w, c3h = Inches(6.667), Inches(1.230), Inches(2.865), Inches(4.038)
    card(sl, c3x, c3y, c3w, c3h, border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)

    mono(sl, "Flow Matching",
         Inches(6.901), Inches(1.465), Inches(2.396), Inches(0.504),
         sz=18, bold=True, clr=C_GREEN)

    sans(sl, '"직선 경로 최적화 (채택)"',
         Inches(6.901), Inches(2.125), Inches(2.396), Inches(0.295),
         sz=14, bold=True, clr=C_DARK)

    # Code block
    rect(sl, Inches(6.901), Inches(2.654), Inches(2.396), Inches(0.703), fill=C_CODE)
    rect(sl, Inches(6.901), Inches(2.654), Inches(0.031), Inches(0.703), fill=C_BLUE)
    mono(sl, "L = ||v_θ - (x₁-x₀)||²",
         Inches(6.948), Inches(2.668), Inches(2.340), Inches(0.671),
         sz=12, bold=False, clr=C_BLUE)

    for txt, iy in [("Optimal Transport", 3.592), ("단 5스텝 복원", 3.950)]:
        sans(sl, "→",
             Inches(6.901), Inches(iy), Inches(0.172), Inches(0.270),
             sz=11, bold=True, clr=C_BLUE)
        sans(sl, txt,
             Inches(7.088), Inches(iy), Inches(2.100), Inches(0.270),
             sz=12, bold=False, clr=C_DARK)

    # Green footer
    rect(sl, Inches(6.901), Inches(4.464), Inches(2.396), Inches(0.539), fill=C_WHITE)
    sans(sl, "실시간 연속 제어 ✅  MoNa-pi 채택!",
         Inches(6.901), Inches(4.464), Inches(2.396), Inches(0.539),
         sz=13, bold=True, clr=C_GREEN, align=CENTER)


def s02_background(prs):
    """Slide 03 — 연구 배경 및 문제 정의"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "연구 배경 및 문제 정의")

    # ── 왼쪽 ──────────────────────────────────────────────────────────────────
    sans(sl, "무엇을 만드는가?",
         Inches(0.469), Inches(1.230), Inches(3.977), Inches(0.338),
         sz=16, bold=True, clr=C_DARK)

    # Icon + text rows
    icon_rows = [
        ("[+]", "3-DOF", " 옴니휠 로봇",  1.803),
        ("[T]", "자연어 명령 → ", "연속 제어",  2.328),
        ("[O]", "8-Frame", " Fish-eye 입력", 2.854),
    ]
    for icon, bold_txt, normal_txt, iy in icon_rows:
        mono(sl, icon,
             Inches(0.469), Inches(iy), Inches(0.312), Inches(0.330),
             sz=17, bold=True, clr=C_BLUE)
        box, tf = _tb_raw(sl, Inches(0.938), Inches(iy + 0.017), Inches(3.500), Inches(0.295))
        p = tf.paragraphs[0]
        p.alignment = LEFT
        r1 = p.add_run()
        r1.text = bold_txt
        r1.font.size = Pt(14)
        r1.font.bold = True
        r1.font.color.rgb = C_BLUE
        r1.font.name = FS
        r2 = p.add_run()
        r2.text = normal_txt
        r2.font.size = Pt(15)
        r2.font.bold = False
        r2.font.color.rgb = C_DARK
        r2.font.name = FS

    # ── 오른쪽 ────────────────────────────────────────────────────────────────
    sans(sl, "핵심 도전 과제",
         Inches(4.759), Inches(1.230), Inches(4.773), Inches(0.338),
         sz=16, bold=True, clr=C_DARK)

    # 붉은 카드 (실시간 병목)
    card(sl, Inches(4.759), Inches(1.803), Inches(4.773), Inches(1.709),
         border_clr=C_RED, border_w=1.4, fill=C_WHITE)

    mono(sl, "실시간 배포 병목 (Latency vs Control)",
         Inches(4.993), Inches(2.037), Inches(4.286), Inches(0.338),
         sz=12, bold=True, clr=C_RED)

    # 4Hz
    box, tf = _tb_raw(sl, Inches(4.993), Inches(2.418), Inches(1.500), Inches(0.617))
    p = tf.paragraphs[0]; p.alignment = LEFT
    r1 = p.add_run(); r1.text = "4"; r1.font.size = Pt(31); r1.font.bold = True
    r1.font.color.rgb = C_DARK; r1.font.name = FM
    r2 = p.add_run(); r2.text = "Hz"; r2.font.size = Pt(15); r2.font.bold = True
    r2.font.color.rgb = C_DARK; r2.font.name = FM

    mono(sl, "VLA 재계획",
         Inches(4.993), Inches(3.035), Inches(1.500), Inches(0.280),
         sz=10, clr=C_MUTED)

    mono(sl, "VS",
         Inches(6.977), Inches(2.677), Inches(0.281), Inches(0.311),
         sz=15, bold=True, clr=C_RED, align=CENTER)

    # 50Hz
    box, tf = _tb_raw(sl, Inches(7.500), Inches(2.418), Inches(1.500), Inches(0.617))
    p = tf.paragraphs[0]; p.alignment = LEFT
    r1 = p.add_run(); r1.text = "50"; r1.font.size = Pt(31); r1.font.bold = True
    r1.font.color.rgb = C_DARK; r1.font.name = FM
    r2 = p.add_run(); r2.text = "Hz"; r2.font.size = Pt(15); r2.font.bold = True
    r2.font.color.rgb = C_DARK; r2.font.name = FM

    mono(sl, "로컬 제어",
         Inches(7.500), Inches(3.035), Inches(1.500), Inches(0.280),
         sz=10, clr=C_MUTED)

    # Challenge items
    challenges = [
        ("연속 고주파 제어 필요", 3.637),
        ("Action Chunking으로 행동 일관성 확보", 4.008),
        ("언어 이해 일반화 (Instruction Pool)", 4.379),
    ]
    for txt, iy in challenges:
        sans(sl, "⚠️",
             Inches(4.759), Inches(iy), Inches(0.234), Inches(0.238),
             sz=12, clr=C_ORANGE)
        sans(sl, txt,
             Inches(5.050), Inches(iy), Inches(4.200), Inches(0.238),
             sz=12, clr=C_DARK)


def s03_design_focus(prs):
    """Slide 04 — 핵심 설계 포인트 & 기대 효과"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "핵심 설계 포인트 & 기대 효과")

    # 2×2 카드 그리드
    grid = [
        # (col, row, num, title, lines, effect)
        (0, 0, "01", "연속 Flow Matching",
         "이산 분류(Class)가 아닌 연속 경로\n단 5스텝 ODE 복원",
         "물리적으로 부드러운 연속 궤적"),
        (1, 0, "02", "Action Chunking (h=10)",
         "미래 10스텝 동시 예측\n4Hz 재계획 / 50Hz 로컬 루프",
         "고주파 제어 달성 & Jitter 감소"),
        (0, 1, "03", "Instruction Pool",
         "카테고리당 15개 Paraphrase\n매 스텝 무작위 선택 학습",
         "다양한 자연어에 강건한 일반화"),
        (1, 1, "04", "BF16 학습 & 배포",
         "FP32 대비 메모리 절반\nGradient 안정성 확보",
         "GX10·Jetson 네이티브 지원"),
    ]
    xs = [Inches(0.469), Inches(5.117)]
    ys = [Inches(1.230), Inches(3.277)]
    cw, ch = Inches(4.414), Inches(1.844)

    for col, row, num, title, body, effect in grid:
        x, y = xs[col], ys[row]
        card(sl, x, y, cw, ch, border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)
        mono(sl, num,
             x + Inches(0.234), y + Inches(0.235), Inches(0.562), Inches(0.469),
             sz=31, bold=True, clr=C_BLUE)
        sans(sl, title,
             x + Inches(0.992), y + Inches(0.235), Inches(3.188), Inches(0.389),
             sz=14, bold=True, clr=C_DARK)
        sans(sl, body,
             x + Inches(0.992), y + Inches(0.741), Inches(3.188), Inches(0.582),
             sz=12, clr=C_MUTED)
        sans(sl, effect,
             x + Inches(0.992), y + Inches(1.323), Inches(3.188), Inches(0.338),
             sz=12, bold=True, clr=C_GREEN)


def s04_approach(prs):
    """Slide 05 — 제안 방법: 세 가지 핵심 전환"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "제안 방법: 세 가지 핵심 전환")

    rows = [
        # (y, num, from_txt, to_txt, detail)
        (1.230, "1", "이산 분류 (Class)", "Flow Matching",
         "연속 액션 공간 경로 생성 (부드러운 제어)"),
        (2.504, "2", "단발 예측 (1-step)", "Action Chunking",
         "미래 10스텝 동시 예측 (행동 일관성 확보)"),
        (3.777, "3", "고정 문장 (Fixed)", "Instruction 다양화",
         "카테고리당 15개 Paraphrase (Success +10%p)"),
    ]
    for iy, num, from_txt, to_txt, detail in rows:
        card(sl, Inches(0.469), Inches(iy), Inches(9.062), Inches(1.109),
             border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)

        # Number block
        rect(sl, Inches(0.781), Inches(iy + 0.305), Inches(0.469), Inches(0.469), fill=C_BLUE)
        mono(sl, num,
             Inches(0.781), Inches(iy + 0.305), Inches(0.469), Inches(0.469),
             sz=18, bold=True, clr=C_WHITE, align=CENTER)

        # From
        sans(sl, from_txt,
             Inches(1.562), Inches(iy + 0.371), Inches(2.812), Inches(0.359),
             sz=16, bold=False, clr=C_MUTED)

        # Arrow
        mono(sl, "→",
             Inches(4.594), Inches(iy + 0.273), Inches(0.234), Inches(0.533),
             sz=26, bold=True, clr=C_BLUE, align=CENTER)

        # To
        sans(sl, to_txt,
             Inches(5.141), Inches(iy + 0.196), Inches(4.078), Inches(0.359),
             sz=18, bold=True, clr=C_BLUE)

        # Detail
        sans(sl, detail,
             Inches(5.141), Inches(iy + 0.633), Inches(4.078), Inches(0.359),
             sz=12, clr=C_DARK)


def s05_architecture(prs):
    """Slide 06 — MoNa-pi 아키텍처 파이프라인"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "MoNa-pi 아키텍처 파이프라인")

    # 4개 파이프라인 박스
    boxes = [
        # (x, y, w, h, border_clr, border_w, title, body)
        (0.625, 1.504, 1.875, 1.572, C_DARK,   1.0, "INPUT",    "8-Frame 이미지\n+\n자연어 명령"),
        (2.917, 1.509, 1.875, 1.562, C_BLUE,   2.0, "BACKBONE", "PaliGemma 3B\n(SigLIP + Gemma-2B)"),
        (5.208, 1.509, 1.875, 1.562, C_ORANGE, 2.0, "EXPERT",   "Action Expert\n(Cross-Attention)"),
        (7.500, 1.387, 1.875, 1.807, C_GREEN,  2.0, "OUTPUT",   "Flow Matching\n(ODE 5-step)\n↓\nAction Chunk"),
    ]
    for bx, by, bw, bh, bc, bw_pt, btitle, bbody in boxes:
        card(sl, Inches(bx), Inches(by), Inches(bw), Inches(bh),
             border_clr=bc, border_w=bw_pt, fill=C_WHITE)
        mono(sl, btitle,
             Inches(bx + 0.078), Inches(by + 0.078), Inches(bw - 0.156), Inches(0.320),
             sz=13, bold=True, clr=C_BLUE)
        sans(sl, bbody,
             Inches(bx + 0.078), Inches(by + 0.430), Inches(bw - 0.156), Inches(bh - 0.508),
             sz=10, clr=C_MUTED)

    # Arrows between boxes
    for ax_x in [2.591, 4.883, 7.174]:
        mono(sl, "→",
             Inches(ax_x), Inches(2.023), Inches(0.234), Inches(0.450),
             sz=26, bold=True, clr=C_ORANGE, align=CENTER)

    # 하단 배포 바
    card(sl, Inches(0.469), Inches(3.600), Inches(9.062), Inches(0.686),
         border_clr=C_BORDER, border_w=1.0, fill=C_WHITE, left_accent=C_BLUE)
    sans(sl, "추론 파이프라인 배포 구조",
         Inches(0.781), Inches(3.795), Inches(2.500), Inches(0.391),
         sz=14, bold=True, clr=C_BLUE)
    mono(sl, "FastAPI 서버 (GX10) ↔ Action Chunk Buffer (Jetson) → 50Hz 로컬 제어",
         Inches(3.634), Inches(3.810), Inches(5.700), Inches(0.330),
         sz=12, clr=C_DARK)

    # 추가 세부 사항 (아키텍처 하단)
    detail_items = [
        ("window_size=8, horizon=10, backbone_out_dim=1024", 4.450),
        ("Action Expert: 4-layer Transformer, 8-head cross-attention, dim=256", 4.780),
        ("Flow: CFM loss  ||v_θ − (x₁−x₀)||²,  ODE n_steps=5 (Euler/Heun)", 5.110),
    ]
    for txt, iy in detail_items:
        mono(sl, txt,
             Inches(0.469), Inches(iy), Inches(9.062), Inches(0.295),
             sz=9, clr=C_MUTED)


def s_adaln(prs):
    """Slide 07 — π0 정통 Action Expert: AdaLN-Zero 적용"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "π0 정통 Action Expert: AdaLN-Zero 시간 컨디셔닝")

    # ── 왼쪽: 기존 방식 (Before) ─────────────────────────────────────────────
    card(sl, Inches(0.469), Inches(1.230), Inches(4.375), Inches(4.156),
         border_clr=C_RED, border_w=2.0, fill=C_WHITE)

    mono(sl, "BEFORE  단순 덧셈",
         Inches(0.562), Inches(1.310), Inches(4.180), Inches(0.340),
         sz=14, bold=True, clr=C_RED)

    rect(sl, Inches(0.562), Inches(1.700), Inches(4.180), Inches(1.180), fill=C_CODE)
    rect(sl, Inches(0.562), Inches(1.700), Inches(0.031), Inches(1.180), fill=C_RED)
    mono_lines(sl, [
        ("t_emb = time_mlp(t)        # (B,1,512)", False, C_DARK),
        ("h = action_proj(x_t) + t_emb  # 그냥 더함", False, C_DARK),
        ("", False, C_DARK),
        ("# query만 norm, key/value는 raw", False, C_MUTED),
        ("attn = self_attn(norm(h), h, h)", False, C_DARK),
        ("h = h + attn               # gate 없음", False, C_DARK),
        ("h = h + mlp(norm(h))", False, C_DARK),
    ], Inches(0.609), Inches(1.716), Inches(4.133), Inches(1.148), sz=9)

    sans_lines(sl, [
        ("문제점:", True, C_RED),
        ("- timestep이 단순 offset으로만 작용", False, C_DARK),
        ("- 각 레이어에서 조건 강도 조절 불가", False, C_DARK),
        ("- 학습 초기 불안정 (초기값 편향)", False, C_DARK),
    ], Inches(0.562), Inches(2.960), Inches(4.180), Inches(1.000), sz=11)

    # Red badge
    rect(sl, Inches(0.562), Inches(4.100), Inches(4.180), Inches(0.590), fill=C_RED)
    sans(sl, "논문과 불일치 — DiT/π0 방식 미적용",
         Inches(0.562), Inches(4.100), Inches(4.180), Inches(0.590),
         sz=11, bold=True, clr=C_WHITE, align=CENTER)

    # Arrow
    mono(sl, "→",
         Inches(4.875), Inches(2.750), Inches(0.390), Inches(0.500),
         sz=30, bold=True, clr=C_ORANGE, align=CENTER)

    # ── 오른쪽: AdaLN-Zero (After) ───────────────────────────────────────────
    card(sl, Inches(5.156), Inches(1.230), Inches(4.375), Inches(4.156),
         border_clr=C_GREEN, border_w=2.0, fill=C_WHITE)

    mono(sl, "AFTER  AdaLN-Zero  (π0 정통)",
         Inches(5.249), Inches(1.310), Inches(4.180), Inches(0.340),
         sz=14, bold=True, clr=C_GREEN)

    rect(sl, Inches(5.249), Inches(1.700), Inches(4.180), Inches(1.620), fill=C_CODE)
    rect(sl, Inches(5.249), Inches(1.700), Inches(0.031), Inches(1.620), fill=C_GREEN)
    mono_lines(sl, [
        ("cond_emb = TimestepEmbedder(t) # (B,512)", False, C_DARK),
        ("# 레이어마다 AdaLN modulation 생성", False, C_MUTED),
        ("α1,β1,γ1,α2,β2,γ2 = AdaLNMod(cond_emb)", False, C_DARK),
        ("", False, C_DARK),
        ("# scale + shift + gate (self-attn)", False, C_MUTED),
        ("h += γ1 * self_attn(norm(h)*(1+α1)+β1)", False, C_DARK),
        ("# VLM cross-attn (timestep 미적용)", False, C_MUTED),
        ("h += cross_attn(norm(h), vlm_cond)", False, C_DARK),
        ("# scale + shift + gate (mlp)", False, C_MUTED),
        ("h += γ2 * mlp(norm(h)*(1+α2)+β2)", False, C_DARK),
    ], Inches(5.296), Inches(1.716), Inches(4.133), Inches(1.588), sz=8)

    sans_lines(sl, [
        ("핵심 개선:", True, C_GREEN),
        ("- Zero-init → 학습 초기 identity 보장", False, C_DARK),
        ("- 레이어별 독립적 scale/shift/gate", False, C_DARK),
        ("- DiT(2023), π0(2024) 동일 방식", False, C_DARK),
    ], Inches(5.249), Inches(3.410), Inches(4.180), Inches(0.980), sz=11)

    # Green badge
    rect(sl, Inches(5.249), Inches(4.100), Inches(4.180), Inches(0.590), fill=C_GREEN)
    box, tf = _tb_raw(sl, Inches(5.249), Inches(4.100), Inches(4.180), Inches(0.590))
    mixed_para(tf, [
        ("검증 완료  ", True, C_WHITE, FM),
        ("verify_mona_expert.py: Max Error < 1e-6", False, C_WHITE, FM),
    ], sz=10, first=True)


def s06_dataset(prs):
    """Slide 08 — 데이터셋 및 수집 파이프라인 전환"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "데이터셋 및 수집 파이프라인 전환")

    # ── 왼쪽 ──────────────────────────────────────────────────────────────────
    sans(sl, "현재 데이터셋 구성",
         Inches(0.469), Inches(1.230), Inches(4.062), Inches(0.338),
         sz=16, bold=True, clr=C_DARK)

    label_val_pairs = [
        ("규모 (HDF5)",         None,                               1.803),
        ("150 에피소드",         "(Train 120 / Val 30)",             2.064),
        ("카테고리 & 언어",       None,                               2.615),
        ("9개",                  "이동 카테고리",                     2.877),
        ("15x",                  "Paraphrase Pool",                  3.146),
        ("데이터 증강",           None,                               3.686),
        (None,                   "Jitter / Color Erase / Noise",     3.943),
    ]
    for bold_txt, normal_txt, iy in label_val_pairs:
        if bold_txt and normal_txt and bold_txt not in ("규모 (HDF5)", "카테고리 & 언어", "데이터 증강"):
            # Mixed line: bold + normal
            box, tf = _tb_raw(sl, Inches(0.469), Inches(iy), Inches(3.977), Inches(0.295))
            p = tf.paragraphs[0]; p.alignment = LEFT
            r1 = p.add_run(); r1.text = bold_txt
            r1.font.size = Pt(13); r1.font.bold = True
            r1.font.color.rgb = C_BLUE; r1.font.name = FS
            r2 = p.add_run(); r2.text = " " + normal_txt
            r2.font.size = Pt(14); r2.font.bold = False
            r2.font.color.rgb = C_DARK; r2.font.name = FS
        elif bold_txt and not normal_txt:
            # Label only (muted)
            mono(sl, bold_txt,
                 Inches(0.469), Inches(iy), Inches(3.977), Inches(0.295),
                 sz=10, clr=C_MUTED)
        elif not bold_txt and normal_txt:
            # Value only
            sans(sl, normal_txt,
                 Inches(0.469), Inches(iy), Inches(3.977), Inches(0.295),
                 sz=14, clr=C_DARK)

    # ── 오른쪽 카드 ───────────────────────────────────────────────────────────
    card(sl, Inches(4.759), Inches(1.230), Inches(4.773), Inches(4.273),
         border_clr=C_ORANGE, border_w=2.0, fill=C_WHITE)

    # Warning icon block
    rect(sl, Inches(4.993), Inches(1.035), Inches(0.547), Inches(0.398), fill=C_CODE)
    sans(sl, "⚠️",
         Inches(4.993), Inches(1.035), Inches(0.547), Inches(0.398),
         sz=22, align=CENTER)

    sans(sl, "파이프라인 전면 개편 중",
         Inches(5.032), Inches(1.582), Inches(4.226), Inches(0.390),
         sz=16, bold=True, clr=C_ORANGE)

    # Mixed text block
    box, tf = _tb_raw(sl, Inches(5.032), Inches(2.076), Inches(4.226), Inches(1.406))
    mixed_para(tf, [
        ("기존 ", False, C_DARK, FS),
        ("동기식(Synchronous)", True, C_DARK, FS),
        (" 수집에서", False, C_DARK, FS),
    ], sz=12, first=True)
    mixed_para(tf, [
        ("비동기식(Asynchronous)", True, C_DARK, FS),
        ("으로 전환 중", False, C_DARK, FS),
    ], sz=12)
    mixed_para(tf, [
        ("향후 ", False, C_DARK, FS),
        ("고정되지 않은 카테고리", True, C_DARK, FS),
        ("로 재수집 예정", False, C_DARK, FS),
    ], sz=12)

    # 비동기 파이프라인 블록
    rect(sl, Inches(5.032), Inches(3.678), Inches(4.226), Inches(1.490), fill=C_CODE)
    rect(sl, Inches(5.032), Inches(3.678), Inches(0.031), Inches(1.490), fill=C_ORANGE)
    mono_lines(sl, [
        ("[Camera] 10Hz → Queue",    False, C_ORANGE),
        ("[Teleop] 50Hz → Queue",    False, C_ORANGE),
        ("↓ Timestamp 정렬 (±50ms)", False, C_ORANGE),
        ("HDF5 일괄 저장",            False, C_ORANGE),
    ], Inches(5.079), Inches(3.693), Inches(4.148), Inches(1.459), sz=10)


def s07_ablation(prs, buf_ab):
    """Slide 08 — Ablation Study"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "Ablation Study — Offline 평가 (Val, n=53)")

    # ── 왼쪽 표 (RECT 방식) ──────────────────────────────────────────────────
    # 헤더
    hdr_defs = [
        ("실험",  Inches(0.469), Inches(1.498)),
        ("설명",  Inches(1.967), Inches(1.709)),
        ("FPE ↓", Inches(3.676), Inches(1.168)),
    ]
    row_h = Inches(0.498)
    for htxt, hx, hw in hdr_defs:
        rect(sl, hx, Inches(1.230), hw, row_h, fill=C_BLUE)
        mono(sl, htxt, hx + Inches(0.05), Inches(1.230) + Inches(0.04),
             hw - Inches(0.1), row_h - Inches(0.08),
             sz=12, bold=True, clr=C_WHITE, align=CENTER)

    # Data rows
    data_rows = [
        # (cells, highlight)  highlight: None/'red'/'green'/'orange'
        (["E3 v2", "기준 모델", "0.673"], None),
        (["E4 off", "text=0", "1.085"],   'red'),
        (["E6 FP16", "BF16→FP16", "2.093"], 'red'),
    ]
    row_ys = [Inches(1.728), Inches(2.230), Inches(2.721)]
    row_hs = [Inches(0.498), Inches(0.490), Inches(0.506)]
    col_xs = [Inches(0.469), Inches(1.967), Inches(3.676)]
    col_ws = [Inches(1.498), Inches(1.709), Inches(1.168)]

    for (cells, hl), ry, rh in zip(data_rows, row_ys, row_hs):
        if hl == 'red':
            rect(sl, Inches(0.469), ry, Inches(4.375), rh, fill=C_RED)
            for txt, cx, cw in zip(cells, col_xs, col_ws):
                mono(sl, txt, cx + Inches(0.05), ry + Inches(0.04),
                     cw - Inches(0.1), rh - Inches(0.08),
                     sz=12, bold=True, clr=C_WHITE, align=CENTER)
        else:
            for txt, cx, cw in zip(cells, col_xs, col_ws):
                mono(sl, txt, cx + Inches(0.05), ry + Inches(0.04),
                     cw - Inches(0.1), rh - Inches(0.08),
                     sz=12, bold=False, clr=C_DARK, align=CENTER)

    # Ablation chart
    sl.shapes.add_picture(buf_ab, Inches(0.469), Inches(3.280), Inches(4.375), Inches(2.125))

    # ── 오른쪽 카드 ───────────────────────────────────────────────────────────
    card(sl, Inches(5.156), Inches(1.230), Inches(4.375), Inches(3.891),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)

    sans(sl, "핵심 해석",
         Inches(5.391), Inches(1.465), Inches(3.906), Inches(0.390),
         sz=16, bold=True, clr=C_BLUE)

    # Item 1
    sans(sl, "💡",
         Inches(5.391), Inches(2.092), Inches(0.234), Inches(0.262),
         sz=12, clr=C_DARK)
    box, tf = _tb_raw(sl, Inches(5.640), Inches(2.092), Inches(3.672), Inches(0.787))
    mixed_para(tf, [("Instruction 유효성 검증", True, C_DARK, FS)], sz=12, first=True)
    mixed_para(tf, [
        ("text-off 시 FPE ", False, C_DARK, FS),
        ("61% 증가",         True,  C_RED,  FS),
    ], sz=12)
    mixed_para(tf, [("→ 언어 명령이 궤적 생성에 필수적임", False, C_MUTED, FS)], sz=10)

    # Item 2
    sans(sl, "💡",
         Inches(5.391), Inches(3.153), Inches(0.234), Inches(0.262),
         sz=12, clr=C_DARK)
    box, tf = _tb_raw(sl, Inches(5.640), Inches(3.153), Inches(3.672), Inches(0.787))
    mixed_para(tf, [("FP16 배포 불가 확정", True, C_DARK, FS)], sz=12, first=True)
    mixed_para(tf, [
        ("FPE ", False, C_DARK, FS),
        ("3.1배 증가", True, C_RED, FS),
        (", CL Success 0%", False, C_DARK, FS),
    ], sz=12)
    mixed_para(tf, [("→ 정밀도 손실로 인한 성능 붕괴", False, C_MUTED, FS)], sz=10)

    # Item 3
    sans(sl, "💡",
         Inches(5.391), Inches(4.214), Inches(0.234), Inches(0.262),
         sz=12, clr=C_DARK)
    box, tf = _tb_raw(sl, Inches(5.640), Inches(4.214), Inches(3.672), Inches(0.525))
    mixed_para(tf, [("BF16 배포 확정", True, C_DARK, FS)], sz=12, first=True)
    mixed_para(tf, [
        ("GX10 및 Jetson 환경 BF16 채택", True, C_GREEN, FS)
    ], sz=12)


def s_metrics(prs, buf_traj, buf_sweep):
    """Slide 09 — 평가 지표 정의 및 측정 방법"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "평가 지표 정의 및 측정 방법")

    # ── 왼쪽 ──────────────────────────────────────────────────────────────────
    sans(sl, "핵심 지표",
         Inches(0.469), Inches(1.152), Inches(3.977), Inches(0.295),
         sz=14, bold=True, clr=C_DARK)

    # FPE 카드
    card(sl, Inches(0.469), Inches(1.604), Inches(3.977), Inches(1.307),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)
    mono(sl, "FPE",
         Inches(0.562), Inches(1.697), Inches(1.500), Inches(0.370),
         sz=16, bold=True, clr=C_BLUE)
    rect(sl, Inches(0.562), Inches(2.105), Inches(3.790), Inches(0.391), fill=C_CODE)
    mono(sl, "|| p_pred − p_gt ||₂",
         Inches(0.570), Inches(2.115), Inches(3.774), Inches(0.371),
         sz=12, clr=C_ORANGE, align=CENTER)
    sans(sl, "최종 위치 유클리드 거리 오차",
         Inches(0.562), Inches(2.574), Inches(3.790), Inches(0.300),
         sz=10, clr=C_MUTED)

    # Success@T 카드
    card(sl, Inches(0.469), Inches(2.973), Inches(3.977), Inches(1.299),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)
    mono(sl, "Success@T",
         Inches(0.562), Inches(3.066), Inches(2.000), Inches(0.370),
         sz=16, bold=True, clr=C_BLUE)
    rect(sl, Inches(0.562), Inches(3.475), Inches(3.790), Inches(0.383), fill=C_CODE)
    mono(sl, "FPE < T (T=1.0 Primary)",
         Inches(0.570), Inches(3.485), Inches(3.774), Inches(0.363),
         sz=12, clr=C_ORANGE, align=CENTER)
    sans(sl, "GT mean disp(9.42) 대비 10.6% 허용",
         Inches(0.562), Inches(3.936), Inches(3.790), Inches(0.300),
         sz=10, clr=C_MUTED)

    # Threshold sweep chart
    sl.shapes.add_picture(buf_sweep, Inches(0.469), Inches(4.334), Inches(3.977), Inches(1.094))

    # ── 오른쪽 ────────────────────────────────────────────────────────────────
    sans(sl, "궤적 시각화 (Bird's-eye 2D)",
         Inches(4.759), Inches(1.152), Inches(4.773), Inches(0.295),
         sz=14, bold=True, clr=C_DARK)

    # 궤적 카드
    card(sl, Inches(4.759), Inches(1.604), Inches(4.773), Inches(3.125),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)
    sl.shapes.add_picture(buf_traj,
                          Inches(4.837), Inches(1.682), Inches(4.584), Inches(2.938))

    # CL 설명 카드 (우측 하단)
    card(sl, Inches(4.759), Inches(4.820), Inches(4.773), Inches(0.805),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE, left_accent=C_GREEN)
    sans(sl, "Closed-Loop(CL) 평가",
         Inches(4.853), Inches(4.860), Inches(4.585), Inches(0.270),
         sz=12, bold=True, clr=C_GREEN)
    sans(sl, "Offline: 단일 스텝 → CL: 연속 행동, 오차 누적 → 진짜 지표",
         Inches(4.853), Inches(5.130), Inches(4.585), Inches(0.225),
         sz=10, clr=C_MUTED)

    box, tf = _tb_raw(sl, Inches(4.853), Inches(5.355), Inches(4.585), Inches(0.225))
    mixed_para(tf, [
        ("→ val loss 낮아도 CL 낮을 수 있음", True, C_ORANGE, FS)
    ], sz=10, first=True)


def s08_results(prs, buf_cmp):
    """Slide 10 — 주요 결과: 폐루프 시뮬레이션"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "주요 결과: 폐루프 시뮬레이션 (CL Eval)")

    # ── 왼쪽 표 (RECT 방식) ──────────────────────────────────────────────────
    col_defs = [
        ("모델",          Inches(0.477), Inches(1.099)),
        ("Val Loss",      Inches(1.575), Inches(0.977)),
        ("FPE ↓",         Inches(2.553), Inches(0.856)),
        ("CL Success ↑",  Inches(3.409), Inches(1.825)),
    ]
    row_h = Inches(0.449)
    hdr_y = Inches(1.230)

    for htxt, hx, hw in col_defs:
        rect(sl, hx, hdr_y, hw, row_h, fill=C_BLUE)
        mono(sl, htxt, hx + Inches(0.04), hdr_y + Inches(0.04),
             hw - Inches(0.08), row_h - Inches(0.08),
             sz=11, bold=True, clr=C_WHITE, align=CENTER)

    # Data rows
    data_rows = [
        (["Random",       "—",      "13.490", "—"     ], None),
        (["v3-A (MoNaVLA)","0.0714", "0.731", "86.7%" ], None),
        (["MoNa-pi v4 ★", "0.0404", "0.049", "100%"  ], 'green'),
    ]
    col_xs = [Inches(0.477), Inches(1.575), Inches(2.553), Inches(3.409)]
    col_ws = [Inches(1.099), Inches(0.977), Inches(0.856), Inches(1.825)]
    total_w = Inches(0.477 + 1.099 + 0.977 + 0.856 + 1.825 - 0.477)

    for i, (cells, hl) in enumerate(data_rows):
        ry = hdr_y + row_h * (i + 1)
        rh = row_h

        if hl == 'green':
            rect(sl, Inches(0.469), ry, Inches(4.765), rh, fill=C_GREEN)
            for txt, cx, cw in zip(cells, col_xs, col_ws):
                mono(sl, txt, cx + Inches(0.04), ry + Inches(0.04),
                     cw - Inches(0.08), rh - Inches(0.08),
                     sz=11, bold=True, clr=C_WHITE, align=CENTER)
        elif hl == 'orange':
            rect(sl, Inches(0.469), ry, Inches(4.765), rh, fill=C_ORANGE)
            for txt, cx, cw in zip(cells, col_xs, col_ws):
                mono(sl, txt, cx + Inches(0.04), ry + Inches(0.04),
                     cw - Inches(0.08), rh - Inches(0.08),
                     sz=11, bold=True, clr=C_WHITE, align=CENTER)
        else:
            for txt, cx, cw in zip(cells, col_xs, col_ws):
                mono(sl, txt, cx + Inches(0.04), ry + Inches(0.04),
                     cw - Inches(0.08), rh - Inches(0.08),
                     sz=11, bold=False, clr=C_DARK, align=CENTER)

    # Comparison chart
    sl.shapes.add_picture(buf_cmp, Inches(0.469), Inches(3.280), Inches(4.765), Inches(2.125))

    # ── 오른쪽 상단 카드 ─────────────────────────────────────────────────────
    card(sl, Inches(5.554), Inches(1.230), Inches(3.977), Inches(1.289),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE, left_accent=C_GREEN)

    sans(sl, "🚀 MoNa-pi v4 (π0 AdaLN-Zero)",
         Inches(5.648), Inches(1.270), Inches(3.789), Inches(0.340),
         sz=13, bold=True, clr=C_GREEN)

    box, tf = _tb_raw(sl, Inches(5.648), Inches(1.610), Inches(3.789), Inches(0.820))
    mixed_para(tf, [
        ("Flow Matching + AdaLN-Zero\n", False, C_DARK, FS),
    ], sz=12, first=True)
    mixed_para(tf, [
        ("CL FPE 0.049m  ", True, C_GREEN, FS),
        ("Success 100%", True, C_GREEN, FS),
    ], sz=12)
    mixed_para(tf, [
        ("val 24 ep, FPE<0.5m & TLD∈[0.7,1.5]", False, C_MUTED, FS),
    ], sz=9)

    # ── 오른쪽 하단 카드 ─────────────────────────────────────────────────────
    card(sl, Inches(5.554), Inches(2.675), Inches(3.977), Inches(2.076),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE, left_accent=C_RED)

    sans(sl, "⚠️ Val Loss ≠ CL 성능",
         Inches(5.648), Inches(2.715), Inches(3.789), Inches(0.340),
         sz=14, bold=True, clr=C_RED)

    box, tf = _tb_raw(sl, Inches(5.648), Inches(3.055), Inches(3.789), Inches(1.570))
    mixed_para(tf, [("v2의 Loss가 가장 낮지만,", False, C_DARK, FS)], sz=12, first=True)
    mixed_para(tf, [
        ("실제 ", False, C_DARK, FS),
        ("Closed-Loop(CL) 성능은 v3-A가 최고", True, C_DARK, FS),
    ], sz=12)
    mixed_para(tf, [
        (" → ", False, C_DARK, FS),
        ("CL Success가 진짜 지표", True, C_DARK, FS),
    ], sz=12)


def s09_categories(prs, buf_cat):
    """Slide 11 — 카테고리별 성능 분석"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "카테고리별 성능 분석 (v3-A, n=30)")

    # ── 왼쪽 표 (RECT 방식) ──────────────────────────────────────────────────
    col_defs = [
        ("카테고리",   Inches(0.469), Inches(1.516)),
        ("n",          Inches(1.985), Inches(0.347)),
        ("FPE ↓",      Inches(2.332), Inches(0.866)),
        ("Success@1.0",Inches(3.198), Inches(1.646)),
    ]
    row_h = Inches(0.390)
    hdr_y = Inches(1.230)
    total_w = Inches(4.375)

    for htxt, hx, hw in col_defs:
        rect(sl, hx, hdr_y, hw, row_h, fill=C_BLUE)
        mono(sl, htxt, hx + Inches(0.03), hdr_y + Inches(0.03),
             hw - Inches(0.06), row_h - Inches(0.06),
             sz=10, bold=True, clr=C_WHITE, align=CENTER)

    cat_rows = [
        ("center_straight", "4", "0.582", "25.0%",  None),
        ("center_left",     "3", "1.055", "0.0%",   'red'),
        ("center_right",    "3", "1.328", "33.3%",  None),
        ("left_straight",   "4", "0.782", "25.0%",  None),
        ("left_left",       "3", "0.451", "66.7%",  'green'),
        ("left_right",      "3", "0.747", "0.0%",   'red'),
        ("right_straight",  "4", "0.650", "50.0%",  None),
        ("right_left",      "3", "0.572", "33.3%",  None),
        ("right_right",     "3", "0.471", "66.7%",  'green'),
    ]
    col_xs = [Inches(0.469), Inches(1.985), Inches(2.332), Inches(3.198)]
    col_ws = [Inches(1.516), Inches(0.347), Inches(0.866), Inches(1.646)]

    for i, (cat, n, fpe, s10, hl) in enumerate(cat_rows):
        ry = hdr_y + row_h * (i + 1)
        cells = [cat, n, fpe, s10]

        if hl == 'red':
            rect(sl, Inches(0.469), ry, total_w, row_h, fill=C_RED)
            for txt, cx, cw in zip(cells, col_xs, col_ws):
                mono(sl, txt, cx + Inches(0.03), ry + Inches(0.03),
                     cw - Inches(0.06), row_h - Inches(0.06),
                     sz=10, bold=True, clr=C_WHITE, align=CENTER)
        elif hl == 'green':
            rect(sl, Inches(0.469), ry, total_w, row_h, fill=C_GREEN)
            for txt, cx, cw in zip(cells, col_xs, col_ws):
                mono(sl, txt, cx + Inches(0.03), ry + Inches(0.03),
                     cw - Inches(0.06), row_h - Inches(0.06),
                     sz=10, bold=True, clr=C_WHITE, align=CENTER)
        else:
            for txt, cx, cw in zip(cells, col_xs, col_ws):
                mono(sl, txt, cx + Inches(0.03), ry + Inches(0.03),
                     cw - Inches(0.06), row_h - Inches(0.06),
                     sz=10, bold=False, clr=C_DARK, align=CENTER)

    # ── 오른쪽 카드 ───────────────────────────────────────────────────────────
    card(sl, Inches(5.156), Inches(1.230), Inches(4.375), Inches(3.750),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)

    sans(sl, "center_left 0% 원인 분석",
         Inches(5.391), Inches(1.465), Inches(3.906), Inches(0.390),
         sz=16, bold=True, clr=C_RED)

    analysis_items = [
        ("🔍", "시각적 모호성",
         "(left/right 구별이 어려운 카메라 구도)", 2.100),
        ("🔍", "통계적 불확실성",
         "(val n=3 으로 샘플 수 부족)", 2.700),
    ]
    for icon, bold_txt, muted_txt, iy in analysis_items:
        sans(sl, icon,
             Inches(5.391), Inches(iy), Inches(0.234), Inches(0.280),
             sz=12, clr=C_DARK)
        sans(sl, bold_txt,
             Inches(5.640), Inches(iy), Inches(3.672), Inches(0.280),
             sz=12, bold=True, clr=C_DARK)
        sans(sl, muted_txt,
             Inches(5.640), Inches(iy + 0.300), Inches(3.672), Inches(0.255),
             sz=10, clr=C_MUTED)

    sans(sl, "Instruction 다양화만으로는 해결 불가",
         Inches(5.391), Inches(3.320), Inches(3.906), Inches(0.280),
         sz=12, clr=C_DARK)

    # 해결책 블록
    rect(sl, Inches(5.391), Inches(3.838), Inches(3.906), Inches(0.877), fill=C_BLUE)
    rect(sl, Inches(5.391), Inches(3.838), Inches(0.047), Inches(0.877), fill=C_BLUE)
    sans(sl, "해결책 (H6)",
         Inches(5.454), Inches(3.878), Inches(3.750), Inches(0.320),
         sz=12, bold=True, clr=C_WHITE)
    sans(sl, "다양한 시작 위치 에피소드 추가 수집 필요",
         Inches(5.454), Inches(4.198), Inches(3.750), Inches(0.320),
         sz=11, bold=True, clr=C_WHITE)

    # Category chart (하단)
    sl.shapes.add_picture(buf_cat, Inches(0.469), Inches(4.680), Inches(9.062), Inches(0.800))


def s10_server(prs):
    """Slide 12 — 추론 서버 검증"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "추론 서버 검증 — ASUS ASCENT GX10")

    # ── 왼쪽 카드 ─────────────────────────────────────────────────────────────
    card(sl, Inches(0.469), Inches(1.230), Inches(3.977), Inches(4.156),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)

    sans(sl, "서버 스펙",
         Inches(0.703), Inches(1.465), Inches(3.539), Inches(0.390),
         sz=16, bold=True, clr=C_DARK)

    spec_items = [
        ("GPU",  "GB10 Superchip",  2.131),
        ("모델",  "v3-A BF16",       2.535),
        ("VRAM", "5.97 GB",          2.945),
        ("API",  "FastAPI",          3.350),
        ("특징",  "비전 캐싱(TTL)",   3.754),
    ]
    for label, value, iy in spec_items:
        sans(sl, label,
             Inches(0.703), Inches(iy), Inches(1.800), Inches(0.340),
             sz=12, clr=C_MUTED)
        mono(sl, value,
             Inches(2.626), Inches(iy), Inches(1.700), Inches(0.340),
             sz=12, bold=True, clr=C_BLUE)

    # ── 오른쪽 Latency 표 ────────────────────────────────────────────────────
    lat_col_defs = [
        ("시나리오",  Inches(4.759), Inches(1.713)),
        ("Latency",   Inches(6.471), Inches(1.104)),
        ("상태",      Inches(7.575), Inches(1.956)),
    ]
    row_h = Inches(0.498)
    hdr_y = Inches(1.230)

    for htxt, hx, hw in lat_col_defs:
        rect(sl, hx, hdr_y, hw, row_h, fill=C_BLUE)
        mono(sl, htxt, hx + Inches(0.04), hdr_y + Inches(0.04),
             hw - Inches(0.08), row_h - Inches(0.08),
             sz=12, bold=True, clr=C_WHITE, align=CENTER)

    lat_rows = [
        (["cold start", "238",   "예열 중"],    None),
        (["warm 평균",  "~243ms", "✅ 4Hz 예산 내"], 'green'),
        (["warm 최소",  "232",   "✅ 예산 내"],  None),
    ]
    col_xs = [Inches(4.759), Inches(6.471), Inches(7.575)]
    col_ws = [Inches(1.713), Inches(1.104), Inches(1.956)]

    for i, (cells, hl) in enumerate(lat_rows):
        ry = hdr_y + row_h * (i + 1)
        if hl == 'green':
            rect(sl, Inches(4.759), ry, Inches(4.773), row_h, fill=C_GREEN)
            for txt, cx, cw in zip(cells, col_xs, col_ws):
                mono(sl, txt, cx + Inches(0.04), ry + Inches(0.04),
                     cw - Inches(0.08), row_h - Inches(0.08),
                     sz=12, bold=True, clr=C_WHITE, align=CENTER)
        else:
            for txt, cx, cw in zip(cells, col_xs, col_ws):
                mono(sl, txt, cx + Inches(0.04), ry + Inches(0.04),
                     cw - Inches(0.08), row_h - Inches(0.08),
                     sz=12, bold=False, clr=C_DARK, align=CENTER)

    # 배포 구조 블록
    card(sl, Inches(4.759), Inches(3.955), Inches(4.773), Inches(1.408),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE, left_accent=C_BLUE)
    mono_lines(sl, [
        ("[GX10 서버] → /predict API → (10×3 action chunk)", False, C_DARK),
        ("↓",                                                  False, C_BLUE),
        ("[Jetson] → Buffer → 50Hz 로컬 제어 → ROS2",          False, C_DARK),
    ], Inches(4.853), Inches(3.990), Inches(4.585), Inches(1.330), sz=10)


def s11_discussion(prs):
    """Slide 13 — 한계 및 토의"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "한계 및 토의")

    rows_def = [
        # (y, accent_clr, icon, icon_clr, title_txt, title_clr, muted_txt, action_txt, action_clr)
        (1.230, C_RED,    "❌", C_RED,    "center_left 0% 실패",   C_RED,
         "시각적 모호성 및 통계적 불확실성 (n=3)",
         "다양한 시작 위치 에피소드 추가 수집 (H6)", C_BLUE),
        (2.625, C_ORANGE, "⚠️", C_ORANGE, "Val Loss ≠ CL 성능 불일치", C_ORANGE,
         "단순 Loss보다 Downstream 일반화가 더 중요함",
         "평가 지표로 Closed-Loop Success 사용 확정", C_GREEN),
        (4.020, C_BLUE,   "⏳", C_BLUE,   "실로봇 미배포 상태",   C_BLUE,
         "현재 시뮬레이션 폐루프만 검증 완료",
         "GX10 구동 완료, Jetson 실로봇 배포 준비 완료", C_ORANGE),
    ]
    for iy, accent, icon, ic, title, tc, muted, action, ac in rows_def:
        card(sl, Inches(0.469), Inches(iy), Inches(9.062), Inches(1.270),
             border_clr=C_BORDER, border_w=1.4, fill=C_WHITE, left_accent=accent)

        sans(sl, icon,
             Inches(0.625), Inches(iy + 0.371), Inches(0.625), Inches(0.498),
             sz=27, clr=ic)

        sans(sl, title,
             Inches(1.484), Inches(iy + 0.118), Inches(7.422), Inches(0.390),
             sz=16, bold=True, clr=tc)

        sans(sl, muted,
             Inches(1.484), Inches(iy + 0.534), Inches(7.422), Inches(0.280),
             sz=12, clr=C_MUTED)

        sans(sl, action,
             Inches(1.484), Inches(iy + 0.866), Inches(7.422), Inches(0.280),
             sz=12, bold=True, clr=ac)


def s12_conclusion(prs):
    """Slide 14 — 결론 및 향후 계획"""
    sl = blank(prs)
    set_bg(sl)
    hdr(sl, "결론 및 향후 계획")

    # ── 왼쪽 카드: 주요 기여 ──────────────────────────────────────────────────
    card(sl, Inches(0.469), Inches(0.957), Inches(4.453), Inches(2.975),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)

    sans(sl, "주요 기여 (완료)",
         Inches(0.586), Inches(1.074), Inches(4.218), Inches(0.390),
         sz=16, bold=True, clr=C_GREEN)

    contrib_items = [
        ("✅", "VLA 파이프라인 구현",        "PaliGemma + Flow Matching",         1.584),
        ("✅", "데이터 전략 개선",             "Instruction Pool (Success +10%p)", 2.134),
        ("✅", "배포 검증 완료",              "GX10 243ms, GPU 5.97GB",           2.684),
        ("✅", "실용적 진단",                 "BF16 확정, center_left 병목 발견", 3.234),
    ]
    for icon, bold_txt, detail_txt, iy in contrib_items:
        sans(sl, icon,
             Inches(0.586), Inches(iy), Inches(0.250), Inches(0.390),
             sz=12, clr=C_GREEN)
        box, tf = _tb_raw(sl, Inches(0.850), Inches(iy), Inches(3.954), Inches(0.390))
        mixed_para(tf, [
            (bold_txt,  True,  C_DARK, FS),
            ("\n" + detail_txt, False, C_MUTED, FS),
        ], sz=11, first=True)

    # ── 오른쪽 카드: 향후 계획 ───────────────────────────────────────────────
    card(sl, Inches(5.078), Inches(0.957), Inches(4.453), Inches(2.944),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE)

    sans(sl, "향후 계획",
         Inches(5.195), Inches(1.074), Inches(4.218), Inches(0.390),
         sz=16, bold=True, clr=C_ORANGE)

    plan_items = [
        ("📌", "단기 (D10-D11)",   "Jetson 실로봇 배포 및 주행 테스트",     1.584),
        ("📌", "중기 (H6)",         "center_left 데이터 보강\n비동기 방식의 고정되지 않은 카테고리 수집", 2.134),
        ("📌", "장기 (논문 완성)",  "실로봇 결과 업데이트, 시연 영상",       2.928),
    ]
    for icon, bold_txt, detail_txt, iy in plan_items:
        sans(sl, icon,
             Inches(5.195), Inches(iy), Inches(0.250), Inches(0.390),
             sz=12, clr=C_ORANGE)
        box, tf = _tb_raw(sl, Inches(5.460), Inches(iy), Inches(3.954), Inches(0.600))
        mixed_para(tf, [
            (bold_txt,  True,  C_DARK, FS),
            ("\n" + detail_txt, False, C_MUTED, FS),
        ], sz=11, first=True)

    # ── 하단 바 ───────────────────────────────────────────────────────────────
    card(sl, Inches(0.469), Inches(3.979), Inches(9.062), Inches(1.143),
         border_clr=C_BORDER, border_w=1.4, fill=C_WHITE, left_accent=C_BLUE)

    mono(sl, "최종 목표",
         Inches(0.563), Inches(4.010), Inches(1.500), Inches(0.280),
         sz=12, bold=True, clr=C_BLUE)

    box, tf = _tb_raw(sl, Inches(0.563), Inches(4.290), Inches(8.859), Inches(0.730))
    mixed_para(tf, [
        ("4Hz 재계획 + 50Hz 로컬 제어", True, C_DARK, FS),
    ], sz=18, first=True)
    mixed_para(tf, [
        ("실로봇 완전 자율 내비게이션 실현", True, C_ORANGE, FS),
    ], sz=18)


# ── Entry point ────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    os.makedirs("/home/minum/26CS/MoNa-pi/reports", exist_ok=True)

    print("[1/7] chart_ablation...")
    buf_ab    = chart_ablation()
    print("[2/7] chart_comparison...")
    buf_cmp   = chart_comparison()
    print("[3/7] chart_categories...")
    buf_cat   = chart_categories()
    print("[4/7] chart_trajectories...")
    buf_traj  = chart_trajectories()
    print("[5/7] chart_threshold_sweep...")
    buf_sweep = chart_threshold_sweep()

    print("[6/7] Building slides...")
    prs = new_prs()
    s01_title(prs)
    s_flow_matching(prs)       # Slide 02 (신규)
    s02_background(prs)        # Slide 03
    s03_design_focus(prs)      # Slide 04
    s04_approach(prs)          # Slide 05
    s05_architecture(prs)               # Slide 06
    s_adaln(prs)                        # Slide 07 (신규: AdaLN-Zero)
    s06_dataset(prs)                    # Slide 08
    s07_ablation(prs, buf_ab)           # Slide 09
    s_metrics(prs, buf_traj, buf_sweep) # Slide 10
    s08_results(prs, buf_cmp)           # Slide 11
    s09_categories(prs, buf_cat)        # Slide 12
    s10_server(prs)                     # Slide 13
    s11_discussion(prs)                 # Slide 14
    s12_conclusion(prs)                 # Slide 15

    out = "/home/minum/26CS/MoNa-pi/reports/mona_pi_midterm_260514.pptx"
    prs.save(out)
    print(f"[7/7] Saved: {out}")
    print(f"      Slides: {len(prs.slides)}")
