#!/usr/bin/env python3
"""MoNa-pi 중간발표 PPTX v2 — 2026-05-28
Claude Dark Design: 다크 카드 + 블록 다이어그램 + 실제 이미지
"""

import io, os
from pathlib import Path
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch
from matplotlib.lines import Line2D

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Slide dimensions (16:9) ────────────────────────────────────────────────────
SW = Inches(10.0)
SH = Inches(5.625)

SAMPLE_DIR = Path(__file__).parent / "sample_imgs"

# ── Typography ──────────────────────────────────────────────────────────────────
FM = "Roboto Mono"
FS = "Space Grotesk"

# ── Color Palette (Claude Dark Design) ─────────────────────────────────────────
BG_SLIDE = RGBColor(0x0D, 0x11, 0x1A)
C_CARD   = RGBColor(0x14, 0x1E, 0x2E)
C_CARD2  = RGBColor(0x1B, 0x28, 0x3C)
C_TEXT   = RGBColor(0xE2, 0xE8, 0xF0)
C_MUTED  = RGBColor(0x94, 0xA3, 0xB8)
C_CYAN   = RGBColor(0x22, 0xD3, 0xEE)
C_GREEN  = RGBColor(0x4A, 0xDE, 0x80)
C_RED    = RGBColor(0xF8, 0x71, 0x71)
C_ORANGE = RGBColor(0xFB, 0x92, 0x3C)
C_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BORDER = RGBColor(0x2D, 0x3E, 0x52)

# matplotlib equivalents
BG_PLT    = '#0D111A'
CARD_PLT  = '#141E2E'
CYAN_PLT  = '#22D3EE'
GREEN_PLT = '#4ADE80'
RED_PLT   = '#F87171'
ORANGE_PLT= '#FB923C'
PURPLE_PLT= '#A78BFA'
TEXT_PLT  = '#E2E8F0'
MUTED_PLT = '#94A3B8'
BORDER_PLT= '#2D3E52'

LEFT   = PP_ALIGN.LEFT
CENTER = PP_ALIGN.CENTER
RIGHT  = PP_ALIGN.RIGHT


# ── Core Primitives ────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(sl):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = BG_SLIDE

def rect(sl, x, y, w, h, fill=None, line=None, lw=0):
    s = sl.shapes.add_shape(1, x, y, w, h)
    if fill is not None:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line and lw:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb_raw(sl, x, y, w, h, wrap=True):
    box = sl.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    return box, tf

def mono(sl, text, x, y, w, h, sz=12, bold=False, clr=None, align=LEFT, wrap=True):
    if clr is None: clr = C_TEXT
    box, tf = _tb_raw(sl, x, y, w, h, wrap=wrap)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = clr
    r.font.name = FM
    return box

def sans(sl, text, x, y, w, h, sz=12, bold=False, clr=None, align=LEFT, wrap=True):
    if clr is None: clr = C_TEXT
    box, tf = _tb_raw(sl, x, y, w, h, wrap=wrap)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.color.rgb = clr
    r.font.name = FS
    return box

def mono_lines(sl, lines, x, y, w, h, sz=12):
    box, tf = _tb_raw(sl, x, y, w, h)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt = line[0]; bld = line[1] if len(line) > 1 else False
        clr = line[2] if len(line) > 2 else C_TEXT
        p.alignment = LEFT
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = clr
        r.font.name = FM
    return box

def sans_lines(sl, lines, x, y, w, h, sz=12):
    box, tf = _tb_raw(sl, x, y, w, h)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        txt = line[0]; bld = line[1] if len(line) > 1 else False
        clr = line[2] if len(line) > 2 else C_TEXT
        p.alignment = LEFT
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = clr
        r.font.name = FS
    return box

def mixed_para(tf, runs, sz=12, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = LEFT
    for spec in runs:
        txt = spec[0]
        bld = spec[1] if len(spec) > 1 else False
        clr = spec[2] if len(spec) > 2 else C_TEXT
        fn  = spec[3] if len(spec) > 3 else FS
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = clr
        r.font.name = fn
    return p

def card(sl, x, y, w, h, border_clr=None, border_w=1.4, fill=None, left_accent=None):
    if fill is None: fill = C_CARD
    if border_clr is None: border_clr = C_BORDER
    s = rect(sl, x, y, w, h, fill=fill, line=border_clr, lw=border_w)
    if left_accent is not None:
        rect(sl, x, y, Inches(0.047), h, fill=left_accent)
    return s

def hdr(sl, title, subtitle=None):
    # Cyan accent bar (left)
    rect(sl, Inches(0.469), Inches(0.240), Inches(0.047), Inches(0.680), fill=C_CYAN)
    mono(sl, title,
         Inches(0.563), Inches(0.312), Inches(9.062), Inches(0.530),
         sz=20, bold=True, clr=C_CYAN, align=LEFT)
    if subtitle:
        sans(sl, subtitle,
             Inches(0.563), Inches(0.842), Inches(9.062), Inches(0.270),
             sz=11, clr=C_MUTED)
    # Bottom separator line
    rect(sl, Inches(0.469), Inches(1.020), Inches(9.062), Inches(0.031), fill=C_BORDER)


# ── Chart generators ───────────────────────────────────────────────────────────

def savefig(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def chart_architecture():
    """Proper block diagram of MoNa-pi architecture"""
    fig = plt.figure(figsize=(6.2, 5.0), facecolor=BG_PLT)
    ax = fig.add_axes([0.01, 0.02, 0.98, 0.96])
    ax.set_facecolor(BG_PLT)
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 9.5)
    ax.axis('off')

    def fbox(x, y, w, h, title, sub='', ec=CYAN_PLT, alpha_fill=0.12):
        bg = FancyBboxPatch((x, y), w, h, boxstyle="round,pad=0.12",
                             facecolor=ec, alpha=alpha_fill,
                             edgecolor=ec, linewidth=1.8, zorder=2)
        ax.add_patch(bg)
        ty = y + h * (0.65 if sub else 0.5)
        ax.text(x + w/2, ty, title, ha='center', va='center',
                color=ec, fontsize=8.5, fontweight='bold',
                fontfamily='monospace', zorder=3)
        if sub:
            ax.text(x + w/2, y + h * 0.22, sub, ha='center', va='center',
                    color=MUTED_PLT, fontsize=6.5,
                    fontfamily='monospace', zorder=3)

    def arr(x1, y1, x2, y2, c=CYAN_PLT, style='->'):
        ax.annotate('', xy=(x2, y2), xytext=(x1, y1),
                    arrowprops=dict(arrowstyle=style, color=c,
                                   lw=1.6, connectionstyle='arc3,rad=0.0'),
                    zorder=5)

    def arr_elbow(x1, y1, xm, y2, c=CYAN_PLT):
        ax.plot([x1, x1, xm], [y1, y2, y2], color=c, lw=1.6, zorder=4)
        ax.annotate('', xy=(xm, y2), xytext=(xm - 0.01, y2),
                    arrowprops=dict(arrowstyle='->', color=c, lw=1.6), zorder=5)

    # ── Row 1: Inputs ────────────────────────────────────────────────────────
    fbox(0.1, 8.1, 4.3, 1.1, '8-Frame Fish-eye Images', '(B, 8, H, W, 3)', ec=TEXT_PLT, alpha_fill=0.07)
    fbox(5.3, 8.1, 4.3, 1.1, 'Natural Language Command', 'instruction_pool (15 paraphrases/cat)', ec=TEXT_PLT, alpha_fill=0.07)

    # ── Row 2: Encoders ──────────────────────────────────────────────────────
    fbox(0.1, 6.6, 2.0, 1.1, 'SigLIP SO400M', '14px patches', ec=CYAN_PLT)
    fbox(2.3, 6.6, 2.1, 1.1, 'Linear Proj', 'patch → token', ec=CYAN_PLT)
    fbox(5.3, 6.6, 4.3, 1.1, 'Gemma Tokenizer + Embed', 'BPE → dense (2048)', ec=CYAN_PLT)

    # ── Row 3: LM Backbone ───────────────────────────────────────────────────
    fbox(0.1, 5.0, 9.5, 1.2, 'Gemma-2B Language Model', 'Full BF16 fine-tuning  →  (B, 64, 2048) context features', ec=CYAN_PLT, alpha_fill=0.18)

    # ── Row 4: Action Expert ─────────────────────────────────────────────────
    fbox(0.1, 3.4, 6.8, 1.2, 'Action Expert (Cross-Attn Transformer)', '4-layer · 8-head · dim=256  |  AdaLN-Zero time conditioning', ec=ORANGE_PLT, alpha_fill=0.15)
    fbox(7.2, 3.4, 2.4, 1.2, 'Noisy x_t + t', 'x₀~N(0,I), t~Uniform(0,1)', ec=PURPLE_PLT, alpha_fill=0.15)

    # ── Row 5: Flow Matching ─────────────────────────────────────────────────
    fbox(0.1, 1.8, 9.5, 1.2, 'Flow Matching ODE  (Euler/Heun  n=5)',
         'L = ||v_θ − (x₁−x₀)||²   ·   x₁ = GT action chunk', ec=GREEN_PLT, alpha_fill=0.15)

    # ── Output label ─────────────────────────────────────────────────────────
    ax.text(4.85, 1.25, 'Action Chunk (B, 10, 3)  →  4Hz replay · 50Hz local control',
            ha='center', va='center', color=GREEN_PLT,
            fontsize=8.5, fontweight='bold', fontfamily='monospace')

    # ── Arrows ────────────────────────────────────────────────────────────────
    # Images → SigLIP
    arr(2.25, 8.1, 1.1, 7.7)
    # SigLIP → Proj
    arr(2.1, 7.15, 2.3, 7.15)
    # Proj → LM
    arr(3.35, 6.6, 3.35, 6.2)
    # Text → Tokenizer
    arr(7.45, 8.1, 7.45, 7.7)
    # Tokenizer → LM
    arr(7.45, 6.6, 7.45, 6.2)
    # LM → Action Expert
    arr(4.85, 5.0, 4.85, 4.6)
    # Noisy → Action Expert
    arr(7.2, 4.0, 6.9, 4.0)
    # Action Expert → ODE
    arr(4.85, 3.4, 4.85, 3.0)
    # ODE → Output
    arr(4.85, 1.8, 4.85, 1.4)

    fig.tight_layout(pad=0.1)
    return savefig(fig)


def chart_flow_matching_concept():
    """Flow Matching: noise → action, ODE trajectory visualization"""
    fig, axes = plt.subplots(1, 3, figsize=(8.5, 3.0), facecolor=BG_PLT)
    titles = ['① Source  x₀~N(0,I)', '② Flow v_θ = (x₁−x₀)', '③ Target  x₁ = GT chunk']
    colors_main = [PURPLE_PLT, CYAN_PLT, GREEN_PLT]
    np.random.seed(42)

    for idx, (ax, title, col) in enumerate(zip(axes, titles, colors_main)):
        ax.set_facecolor(CARD_PLT)
        ax.set_xlim(-3, 3); ax.set_ylim(-3, 3)
        ax.set_aspect('equal'); ax.axis('off')
        ax.set_title(title, color=col, fontsize=9, fontweight='bold', pad=6)
        for sp in ax.spines.values():
            sp.set_color(BORDER_PLT)

        if idx == 0:
            pts = np.random.randn(120, 2)
            ax.scatter(pts[:, 0], pts[:, 1], c=PURPLE_PLT, s=12, alpha=0.6, edgecolors='none')
            circle = plt.Circle((0, 0), 1.5, fill=False, color=PURPLE_PLT, lw=1.2, linestyle='--', alpha=0.5)
            ax.add_patch(circle)
            ax.text(0, 0, 'N(0,I)', ha='center', va='center', color=PURPLE_PLT,
                    fontsize=10, fontweight='bold', fontfamily='monospace')

        elif idx == 1:
            # Show straight-line flows
            src = np.random.randn(18, 2)
            tgt_mean = np.array([1.0, 0.5])
            tgt = tgt_mean + np.random.randn(18, 2) * 0.4
            for s, t in zip(src, tgt):
                ax.annotate('', xy=t, xytext=s,
                            arrowprops=dict(arrowstyle='->', color=CYAN_PLT, lw=0.9, alpha=0.7))
            ax.scatter(src[:, 0], src[:, 1], c=PURPLE_PLT, s=15, alpha=0.5, edgecolors='none')
            ax.scatter(tgt[:, 0], tgt[:, 1], c=GREEN_PLT, s=15, alpha=0.5, edgecolors='none')
            ax.text(-2.5, 2.5, 'v_θ', color=CYAN_PLT, fontsize=13, fontweight='bold', fontfamily='monospace')

        else:
            tgt_mean = np.array([1.0, 0.5])
            pts = tgt_mean + np.random.randn(120, 2) * 0.35
            ax.scatter(pts[:, 0], pts[:, 1], c=GREEN_PLT, s=12, alpha=0.7, edgecolors='none')
            circle = plt.Circle(tgt_mean, 0.6, fill=False, color=GREEN_PLT, lw=1.5, linestyle='-', alpha=0.7)
            ax.add_patch(circle)
            ax.text(tgt_mean[0], tgt_mean[1], 'action', ha='center', va='center',
                    color=GREEN_PLT, fontsize=9, fontweight='bold', fontfamily='monospace')

    fig.suptitle('Conditional Flow Matching: Gaussian Noise → Action Chunk (5-step ODE)',
                 color=TEXT_PLT, fontsize=9, y=0.02)
    fig.tight_layout(pad=0.8)
    return savefig(fig)


def chart_ablation():
    fig, ax = plt.subplots(figsize=(5.5, 3.2), facecolor=BG_PLT)
    ax.set_facecolor(CARD_PLT)
    labels = ['v2 best\n(E3 Baseline)', 'text-off\n(E4 No Lang)', 'FP16\n(E6 dtype)']
    fpe    = [0.673, 1.085, 2.093]
    colors = [GREEN_PLT, RED_PLT, RED_PLT]
    bars   = ax.bar(labels, fpe, color=colors, width=0.5, edgecolor='none')
    ax.axhline(0.673, color=GREEN_PLT, linestyle='--', alpha=0.5, lw=1.5)
    for bar, v in zip(bars, fpe):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.04,
                f'{v:.3f}', ha='center', va='bottom', color=TEXT_PLT, fontsize=11)
    ax.set_ylabel('Offline FPE  (↓ lower is better)', color=MUTED_PLT, fontsize=9)
    ax.set_title('Ablation Study', color=TEXT_PLT, fontsize=11, pad=6)
    ax.tick_params(colors=MUTED_PLT)
    ax.set_ylim(0, 2.6)
    for sp in ax.spines.values(): sp.set_color(BORDER_PLT); sp.set_alpha(0.5)
    # Annotations
    ax.annotate('+61%', xy=(1, 1.085), xytext=(1.4, 1.5),
                color=RED_PLT, fontsize=10, fontweight='bold',
                arrowprops=dict(arrowstyle='->', color=RED_PLT, lw=1.2))
    ax.annotate('+211%', xy=(2, 2.093), xytext=(1.8, 2.4),
                color=RED_PLT, fontsize=10, fontweight='bold',
                arrowprops=dict(arrowstyle='->', color=RED_PLT, lw=1.2))
    fig.tight_layout(pad=1.2)
    return savefig(fig)


def chart_comparison():
    fig, axes = plt.subplots(1, 2, figsize=(8, 3.2), facecolor=BG_PLT)
    models = ['v2', 'v3', 'v3-A ★']
    fpe    = [0.857, 0.740, 0.731]
    s10    = [79.0,  76.7,  86.7]
    fpe_c  = [MUTED_PLT, CYAN_PLT, GREEN_PLT]
    s10_c  = [MUTED_PLT, CYAN_PLT, ORANGE_PLT]

    ax1 = axes[0]; ax1.set_facecolor(CARD_PLT)
    bars = ax1.bar(models, fpe, color=fpe_c, width=0.55, edgecolor='none')
    for bar, v in zip(bars, fpe):
        ax1.text(bar.get_x() + bar.get_width() / 2, v + 0.01,
                 f'{v:.3f}', ha='center', va='bottom', color=TEXT_PLT, fontsize=10)
    ax1.set_title('FPE  (↓ lower is better)', color=TEXT_PLT, fontsize=11, pad=6)
    ax1.set_ylabel('FPE', color=MUTED_PLT, fontsize=9)
    ax1.tick_params(colors=MUTED_PLT); ax1.set_ylim(0, 1.05)
    for sp in ax1.spines.values(): sp.set_color(BORDER_PLT)

    ax2 = axes[1]; ax2.set_facecolor(CARD_PLT)
    bars2 = ax2.bar(models, s10, color=s10_c, width=0.55, edgecolor='none')
    for bar, v in zip(bars2, s10):
        ax2.text(bar.get_x() + bar.get_width() / 2, v + 0.5,
                 f'{v:.1f}%', ha='center', va='bottom', color=TEXT_PLT, fontsize=10)
    ax2.set_title('Success@1.0  (↑ higher is better)', color=TEXT_PLT, fontsize=11, pad=6)
    ax2.set_ylabel('Success Rate (%)', color=MUTED_PLT, fontsize=9)
    ax2.tick_params(colors=MUTED_PLT); ax2.set_ylim(0, 105)
    for sp in ax2.spines.values(): sp.set_color(BORDER_PLT)
    # Star annotation on best
    ax2.annotate('★ Best', xy=(2, 86.7), xytext=(1.5, 95),
                 color=ORANGE_PLT, fontsize=10, fontweight='bold',
                 arrowprops=dict(arrowstyle='->', color=ORANGE_PLT, lw=1.2))

    fig.tight_layout(pad=2.0)
    return savefig(fig)


def chart_trajectories():
    np.random.seed(7)
    fig, ax = plt.subplots(figsize=(5.5, 4.8), facecolor=BG_PLT)
    ax.set_facecolor(CARD_PLT)
    dt = 0.12

    def integrate(sx, sy, sth, vx_fn, om_fn, steps):
        x, y, th = sx, sy, sth
        xs, ys = [x], [y]
        for i in range(steps):
            x += vx_fn(i) * np.cos(th) * dt
            y += vx_fn(i) * np.sin(th) * dt
            th += om_fn(i) * dt
            xs.append(x); ys.append(y)
        return np.array(xs), np.array(ys)

    def add_noise(gx, gy, fpe, seed=0):
        rng = np.random.RandomState(seed)
        n = len(gx)
        dx = rng.randn(n).cumsum() * fpe * 0.04
        dy = rng.randn(n).cumsum() * fpe * 0.04
        ed = rng.randn(2) * fpe * 0.6
        t = np.linspace(0, 1, n)
        return gx + dx + t * ed[0], gy + dy + t * ed[1]

    cases = [
        ('c_str',  0.0, np.pi/2, lambda i: 0.65, lambda i: 0.00, 38, 0.582, CYAN_PLT,   '25%'),
        ('c_lft',  3.5, np.pi/2, lambda i: 0.60, lambda i: 0.35, 38, 1.055, RED_PLT,    '0% ❌'),
        ('l_lft',  7.5, np.pi/2, lambda i: 0.60, lambda i: 0.40, 38, 0.451, GREEN_PLT,  '67% ✅'),
        ('r_rgt', -3.5, np.pi/2, lambda i: 0.60, lambda i:-0.32, 38, 0.471, ORANGE_PLT, '67% ✅'),
        ('l_rgt', -7.5, np.pi/2, lambda i: 0.58, lambda i:-0.28, 38, 0.747, PURPLE_PLT, '0% ❌'),
    ]
    for i, (lbl, sx, sth, vfn, ofn, steps, fpe, col, sr) in enumerate(cases):
        gx, gy = integrate(sx, 0, sth, vfn, ofn, steps)
        px, py = add_noise(gx, gy, fpe, seed=i * 13)
        ax.plot(gx, gy, '-', color=col, lw=2.2, alpha=0.95, zorder=3)
        ax.plot(px, py, '--', color=col, lw=1.5, alpha=0.6, zorder=3)
        ax.scatter([gx[0]], [gy[0]], color=col, s=50, zorder=5, edgecolors='white', lw=0.6)
        ax.scatter([gx[-1]], [gy[-1]], color=col, s=40, marker='*', zorder=5)
        circle = plt.Circle((gx[-1], gy[-1]), fpe, color=col,
                             fill=False, lw=0.8, linestyle=':', alpha=0.5, zorder=4)
        ax.add_patch(circle)
        ax.annotate('', xy=(px[-1], py[-1]), xytext=(gx[-1], gy[-1]),
                    arrowprops=dict(arrowstyle='->', color='white', lw=0.9, alpha=0.7))
        ax.text(gx[0], gy[0] - 0.7, lbl, color=col,
                fontsize=7.5, ha='center', va='top', weight='bold')
        ax.text(px[-1] + 0.25, py[-1], f'FPE={fpe:.2f}\n{sr}',
                color=col, fontsize=7, va='center', alpha=0.9)

    legend_els = [
        Line2D([0],[0], color=TEXT_PLT, lw=2,   label='GT trajectory'),
        Line2D([0],[0], color=TEXT_PLT, lw=1.5, linestyle='--', label='Predicted'),
        Line2D([0],[0], color=TEXT_PLT, lw=0.9, linestyle=':',  label='FPE circle'),
    ]
    ax.legend(handles=legend_els, loc='lower right',
              facecolor=BG_PLT, labelcolor=TEXT_PLT, fontsize=7,
              framealpha=0.9, edgecolor=BORDER_PLT)
    ax.set_xlabel('Lateral', color=MUTED_PLT, fontsize=8)
    ax.set_ylabel('Forward', color=MUTED_PLT, fontsize=8)
    ax.set_title("Bird's-eye: GT vs Predicted (v3-A)", color=TEXT_PLT, fontsize=9, pad=5)
    ax.tick_params(colors=MUTED_PLT, labelsize=7)
    for sp in ax.spines.values(): sp.set_color(BORDER_PLT); sp.set_alpha(0.5)
    ax.set_aspect('equal', adjustable='datalim')
    fig.tight_layout(pad=1.2)
    return savefig(fig)


def chart_threshold_sweep():
    thresholds = [0.5, 1.0, 1.5, 2.0, 3.0]
    budgets    = [5.3, 10.6, 15.9, 21.2, 31.7]
    v2   = [41.7, 79.2, 83.3, 91.7, 95.8]
    v3a  = [33.3, 86.7, 93.3, 96.7, 100.0]

    fig, ax = plt.subplots(figsize=(5.5, 2.6), facecolor=BG_PLT)
    ax.set_facecolor(CARD_PLT)
    x = np.arange(len(thresholds)); w = 0.35
    b1 = ax.bar(x - w/2, v2,  w, color=MUTED_PLT,  alpha=0.9, label='v2')
    b2 = ax.bar(x + w/2, v3a, w, color=ORANGE_PLT, alpha=0.9, label='v3-A ★')
    ax.axvspan(0.5, 1.5, color=CYAN_PLT, alpha=0.06, zorder=0)
    ax.axvline(x=1, color=CYAN_PLT, lw=1.2, linestyle='--', alpha=0.7, label='Primary (T=1.0)')
    for bar, v in zip(b1, v2):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.8, f'{v:.0f}%',
                ha='center', va='bottom', color=TEXT_PLT, fontsize=7.5)
    for bar, v in zip(b2, v3a):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.8, f'{v:.0f}%',
                ha='center', va='bottom', color=ORANGE_PLT, fontsize=7.5, weight='bold')
    ax.set_xticks(x)
    ax.set_xticklabels([f'T={t}\n({b:.0f}%)' for t, b in zip(thresholds, budgets)],
                       color=MUTED_PLT, fontsize=7.5)
    ax.tick_params(colors=MUTED_PLT, labelsize=7.5)
    ax.set_ylabel('Success Rate (%)', color=MUTED_PLT, fontsize=8)
    ax.set_ylim(0, 112)
    ax.set_title('Threshold Sensitivity (GT disp ≈ 9.42)', color=TEXT_PLT, fontsize=8.5, pad=4)
    ax.legend(facecolor=BG_PLT, labelcolor=TEXT_PLT, fontsize=7.5, edgecolor=BORDER_PLT)
    for sp in ax.spines.values(): sp.set_color(BORDER_PLT); sp.set_alpha(0.5)
    fig.tight_layout(pad=1.0)
    return savefig(fig)


def chart_categories():
    fig, ax = plt.subplots(figsize=(8.5, 3.2), facecolor=BG_PLT)
    ax.set_facecolor(CARD_PLT)
    cats = ['c_str','c_lft','c_rgt','l_str','l_lft','l_rgt','r_str','r_lft','r_rgt']
    fpe  = [0.582, 1.055, 1.328, 0.782, 0.451, 0.747, 0.650, 0.572, 0.471]
    s10  = [25,    0,     33,    25,    67,    0,     50,    33,    67]
    bar_colors = [RED_PLT if s == 0 else GREEN_PLT if s >= 60 else CYAN_PLT for s in s10]
    x = np.arange(len(cats)); w = 0.35
    ax.bar(x - w/2, fpe, w, color=CYAN_PLT, alpha=0.8, label='FPE ↓', edgecolor='none')
    ax2 = ax.twinx()
    ax2.bar(x + w/2, s10, w, color=bar_colors, alpha=0.85, label='Success@1.0 ↑', edgecolor='none')
    ax.set_xticks(x)
    ax.set_xticklabels(cats, rotation=30, ha='right', color=TEXT_PLT, fontsize=8.5)
    ax.tick_params(colors=MUTED_PLT); ax2.tick_params(colors=MUTED_PLT)
    ax.set_ylabel('FPE', color=CYAN_PLT, fontsize=9)
    ax2.set_ylabel('Success@1.0 (%)', color=ORANGE_PLT, fontsize=9)
    ax.set_ylim(0, 1.8); ax2.set_ylim(0, 120)
    for sp in ax.spines.values(): sp.set_color(BORDER_PLT); sp.set_alpha(0.4)
    for sp in ax2.spines.values(): sp.set_color(BORDER_PLT); sp.set_alpha(0.4)
    handles = [mpatches.Patch(color=CYAN_PLT, label='FPE'),
               mpatches.Patch(color=ORANGE_PLT, label='Success@1.0')]
    ax.legend(handles=handles, loc='upper right',
              facecolor=BG_PLT, labelcolor=TEXT_PLT, fontsize=8.5, edgecolor=BORDER_PLT)
    # Highlight failures
    for i, (s, cat) in enumerate(zip(s10, cats)):
        if s == 0:
            ax2.text(i + w/2, 3, '0%\n❌', ha='center', va='bottom',
                     color=RED_PLT, fontsize=8, fontweight='bold')
    fig.tight_layout(pad=1.2)
    return savefig(fig)


# ── Slide builders ─────────────────────────────────────────────────────────────

def s01_title(prs):
    sl = blank(prs); set_bg(sl)

    # Decorative background bars
    rect(sl, 0, 0, SW, Inches(0.18), fill=C_CYAN)
    rect(sl, 0, SH - Inches(0.18), SW, Inches(0.18), fill=C_CYAN)
    # Subtle side accent
    rect(sl, 0, 0, Inches(0.08), SH, fill=C_CYAN)

    # Large title
    mono(sl, "MoNa-π",
         Inches(1.2), Inches(1.0), Inches(7.5), Inches(1.5),
         sz=82, bold=True, clr=C_CYAN, align=CENTER)

    # Subtitle block
    card(sl, Inches(1.0), Inches(2.7), Inches(7.8), Inches(0.65),
         border_clr=C_CYAN, border_w=1.5, fill=C_CARD)
    sans(sl, "Flow Matching 기반 고주파 모바일 내비게이션 VLA",
         Inches(1.1), Inches(2.78), Inches(7.6), Inches(0.50),
         sz=18, bold=True, clr=C_TEXT, align=CENTER)

    # Separator line
    rect(sl, Inches(3.5), Inches(3.55), Inches(2.8), Inches(0.03), fill=C_CYAN)

    # Author / date
    sans(sl, "중간발표  ·  2026. 05. 14",
         Inches(1.0), Inches(3.75), Inches(7.8), Inches(0.35),
         sz=14, clr=C_MUTED, align=CENTER)
    sans(sl, "인공지능전공 4학년 이민우  |  Team Monaf",
         Inches(1.0), Inches(4.15), Inches(7.8), Inches(0.35),
         sz=13, clr=C_MUTED, align=CENTER)

    # Corner decorations
    rect(sl, Inches(8.8), Inches(4.8), Inches(0.9), Inches(0.04), fill=C_CYAN)
    rect(sl, Inches(8.8), Inches(4.85), Inches(0.04), Inches(0.6), fill=C_CYAN)
    rect(sl, Inches(0.2), Inches(4.8), Inches(0.9), Inches(0.04), fill=C_CYAN)
    rect(sl, Inches(1.06), Inches(4.85), Inches(0.04), Inches(0.6), fill=C_CYAN)


def s_flow_matching(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "핵심 개념: 왜 Flow Matching인가?",
        "기존 Diffusion 대비 ODE 기반 직선 경로 최적화 → 실시간 제어 가능")

    cards = [
        (0, "Diffusion", '"이산적 & 느림"',
         "Markov Chain\nSDE",
         [("수많은 스텝 필요", 3.576), ("연산 비용 과다", 3.934)],
         "고주파 제어 불가", C_MUTED, C_RED),
        (1, "CFM", '"연속 벡터 필드"',
         "dx/dt = v_θ(x, t)",
         [("ODE 기반 모델링", 3.350), ("수학적 깔끔함", 3.707)],
         "적은 스텝 가능성", C_ORANGE, C_ORANGE),
        (2, "Flow Matching", '"직선 경로 (채택 ✅)"',
         "L = ||v_θ - (x₁-x₀)||²",
         [("Optimal Transport", 3.592), ("단 5스텝 복원", 3.950)],
         "실시간 연속 제어 ✅", C_GREEN, C_GREEN),
    ]
    xs = [Inches(0.469), Inches(3.568), Inches(6.667)]
    cw, ch = Inches(2.865), Inches(4.250)

    for col_idx, ctitle, csub, ccode, bullets, footer, ec, fc in cards:
        x = xs[col_idx]
        card(sl, x, Inches(1.1), cw, ch, border_clr=ec, border_w=2.0, fill=C_CARD)
        # Top accent strip
        rect(sl, x, Inches(1.1), cw, Inches(0.07), fill=ec)

        mono(sl, ctitle,
             x + Inches(0.23), Inches(1.25), cw - Inches(0.46), Inches(0.50),
             sz=18, bold=True, clr=ec)

        sans(sl, csub,
             x + Inches(0.23), Inches(1.85), cw - Inches(0.46), Inches(0.30),
             sz=13, bold=True, clr=C_TEXT)

        # Code block
        rect(sl, x + Inches(0.23), Inches(2.30), cw - Inches(0.46), Inches(0.52), fill=C_CARD2)
        rect(sl, x + Inches(0.23), Inches(2.30), Inches(0.04), Inches(0.52), fill=ec)
        mono(sl, ccode,
             x + Inches(0.30), Inches(2.32), cw - Inches(0.53), Inches(0.48),
             sz=10, clr=ec)

        for txt, iy in bullets:
            sans(sl, "→ " + txt,
                 x + Inches(0.23), Inches(iy), cw - Inches(0.46), Inches(0.270),
                 sz=11, clr=C_TEXT)

        # Footer
        rect(sl, x + Inches(0.23), Inches(4.58), cw - Inches(0.46), Inches(0.45), fill=C_CARD2)
        sans(sl, footer,
             x + Inches(0.23), Inches(4.58), cw - Inches(0.46), Inches(0.45),
             sz=12, bold=True, clr=fc, align=CENTER)


def s02_background(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "연구 배경 및 문제 정의",
        "옴니휠 로봇 + 자연어 → 연속 제어  |  4가지 핵심 도전")

    # ── Left: What ───────────────────────────────────────────────────────────
    card(sl, Inches(0.469), Inches(1.10), Inches(3.977), Inches(4.20),
         border_clr=C_BORDER, fill=C_CARD)
    sans(sl, "무엇을 만드는가?",
         Inches(0.62), Inches(1.18), Inches(3.70), Inches(0.35),
         sz=14, bold=True, clr=C_CYAN)
    rect(sl, Inches(0.62), Inches(1.53), Inches(3.70), Inches(0.025), fill=C_BORDER)

    specs = [
        ("🤖", "3-DOF 옴니휠 로봇", "linear_x / linear_y / angular_z",       1.65),
        ("📷", "Fish-eye 8-Frame", "Wide-angle 카메라, window=0.8s",          2.25),
        ("💬", "자연어 명령 제어", "Instruction Pool (15 para/category)",       2.85),
        ("⚡", "4Hz + 50Hz 이중 루프", "4Hz VLA 재계획 / 50Hz 로컬 실행",      3.45),
    ]
    for icon, bold_t, muted_t, iy in specs:
        sans(sl, icon, Inches(0.62), Inches(iy), Inches(0.30), Inches(0.35), sz=16, clr=C_TEXT)
        sans(sl, bold_t, Inches(0.98), Inches(iy), Inches(3.30), Inches(0.30),
             sz=13, bold=True, clr=C_TEXT)
        sans(sl, muted_t, Inches(0.98), Inches(iy + 0.30), Inches(3.30), Inches(0.25),
             sz=10, clr=C_MUTED)

    # Photo placeholders
    card(sl, Inches(0.62), Inches(4.20), Inches(1.75), Inches(0.90),
         border_clr=C_BORDER, border_w=1.0, fill=C_CARD2)
    mono(sl, "(사진)\n로봇 정면", Inches(0.62), Inches(4.30), Inches(1.75), Inches(0.70),
         sz=9, clr=C_MUTED, align=CENTER)
    card(sl, Inches(2.50), Inches(4.20), Inches(1.85), Inches(0.90),
         border_clr=C_BORDER, border_w=1.0, fill=C_CARD2)
    mono(sl, "(사진)\n주행 환경", Inches(2.50), Inches(4.30), Inches(1.85), Inches(0.70),
         sz=9, clr=C_MUTED, align=CENTER)

    # ── Right: Challenges ────────────────────────────────────────────────────
    challenges = [
        (C_RED,    "① 연속 고주파 제어 필요",
         "이산 분류(class) → 궤적 생성 불가  →  Flow Matching 도입"),
        (C_ORANGE, "② 행동 일관성 확보",
         "단발 예측 → jitter  →  horizon=10 Action Chunking"),
        (C_CYAN,   "③ 언어 이해 일반화",
         "고정 문장 → 과적합  →  카테고리당 15개 Paraphrase Pool"),
        (C_PURPLE, "④ 실시간 배포 현실성",
         "GX10 서버 ↔ Jetson: 4Hz 예산 내 추론 완료 필요"),
    ]
    for i, (ec, title, detail) in enumerate(challenges):
        cy = Inches(1.10 + i * 1.06)
        card(sl, Inches(4.65), cy, Inches(5.0), Inches(0.95),
             border_clr=ec, border_w=1.5, fill=C_CARD, left_accent=ec)
        sans(sl, title, Inches(4.85), cy + Inches(0.10), Inches(4.70), Inches(0.38),
             sz=13, bold=True, clr=ec)
        sans(sl, detail, Inches(4.85), cy + Inches(0.48), Inches(4.70), Inches(0.35),
             sz=11, clr=C_MUTED)


def s03_design_focus(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "핵심 설계 포인트 & 기대 효과",
        "각 결정의 이유 + 예상 기대 효과")

    grid = [
        ("01", "연속 Flow Matching",
         "이산 Class 아닌 연속 경로\n단 5스텝 ODE 복원",
         "물리적으로 부드러운 연속 궤적", C_CYAN),
        ("02", "Action Chunking (h=10)",
         "미래 10스텝 동시 예측\n4Hz 재계획 / 50Hz 로컬 루프",
         "고주파 제어 & Jitter 감소", C_GREEN),
        ("03", "Instruction Pool",
         "카테고리당 15개 Paraphrase\n매 스텝 무작위 선택",
         "다양한 자연어에 강건한 일반화", C_ORANGE),
        ("04", "BF16 학습 & 배포",
         "FP32 대비 메모리 절반\nGradient 안정성 확보",
         "GX10·Jetson 네이티브 지원", C_PURPLE),
    ]
    xs = [Inches(0.469), Inches(5.117)]
    ys = [Inches(1.10),  Inches(3.22)]
    cw, ch = Inches(4.414), Inches(1.90)

    for i, (num, title, body, effect, ec) in enumerate(grid):
        col, row = i % 2, i // 2
        x, y = xs[col], ys[row]
        card(sl, x, y, cw, ch, border_clr=ec, border_w=1.5, fill=C_CARD)
        rect(sl, x, y, cw, Inches(0.06), fill=ec)
        mono(sl, num,
             x + Inches(0.18), y + Inches(0.18), Inches(0.80), Inches(0.52),
             sz=26, bold=True, clr=ec, wrap=False)
        sans(sl, title,
             x + Inches(1.10), y + Inches(0.18), Inches(3.15), Inches(0.40),
             sz=14, bold=True, clr=C_TEXT)
        sans(sl, body,
             x + Inches(1.10), y + Inches(0.65), Inches(3.15), Inches(0.60),
             sz=11, clr=C_MUTED)
        sans(sl, "→ " + effect,
             x + Inches(0.18), y + Inches(1.45), Inches(4.05), Inches(0.34),
             sz=11, bold=True, clr=ec)


def s04_approach(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "제안 방법: 세 가지 핵심 전환",
        "기존 방식의 한계를 돌파한 3가지 핵심 설계 변경")

    rows = [
        (1.10, "1", "이산 분류 (Class)", "→", "Flow Matching",
         "연속 액션 공간 경로 생성 · 5-step ODE · L=||v_θ-(x₁-x₀)||²",
         C_CYAN),
        (2.48, "2", "단발 예측 (1-step)", "→", "Action Chunking",
         "미래 10스텝 동시 예측 · 50Hz 로컬 루프 · 4Hz 재계획",
         C_GREEN),
        (3.86, "3", "고정 문장 (Fixed)", "→", "Instruction 다양화",
         "카테고리당 15개 Paraphrase · Success@1.0 +10%p 향상",
         C_ORANGE),
    ]
    for iy, num, from_t, arrow, to_t, detail, ec in rows:
        card(sl, Inches(0.469), Inches(iy), Inches(9.062), Inches(1.17),
             border_clr=ec, border_w=1.5, fill=C_CARD)
        rect(sl, Inches(0.469), Inches(iy), Inches(0.06), Inches(1.17), fill=ec)

        # Number badge
        rect(sl, Inches(0.70), Inches(iy + 0.30), Inches(0.50), Inches(0.50), fill=ec)
        mono(sl, num,
             Inches(0.70), Inches(iy + 0.30), Inches(0.50), Inches(0.50),
             sz=18, bold=True, clr=C_CARD, align=CENTER)

        sans(sl, from_t,
             Inches(1.40), Inches(iy + 0.32), Inches(2.60), Inches(0.38),
             sz=15, bold=False, clr=C_MUTED)

        mono(sl, "→",
             Inches(4.20), Inches(iy + 0.24), Inches(0.40), Inches(0.55),
             sz=26, bold=True, clr=ec, align=CENTER)

        sans(sl, to_t,
             Inches(4.75), Inches(iy + 0.18), Inches(4.50), Inches(0.40),
             sz=19, bold=True, clr=ec)

        sans(sl, detail,
             Inches(4.75), Inches(iy + 0.66), Inches(4.50), Inches(0.38),
             sz=11, clr=C_MUTED)


def s05_architecture(prs, buf_arch):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "MoNa-pi 아키텍처 파이프라인",
        "PaliGemma (SigLIP + Gemma-2B) + Action Expert + Flow Matching ODE")

    # Block diagram image
    sl.shapes.add_picture(buf_arch, Inches(0.35), Inches(1.07), Inches(5.75), Inches(4.45))

    # Right panel: component details
    card(sl, Inches(6.30), Inches(1.07), Inches(3.50), Inches(4.45),
         border_clr=C_BORDER, fill=C_CARD)

    sans(sl, "구성 요소 상세",
         Inches(6.45), Inches(1.18), Inches(3.20), Inches(0.32),
         sz=13, bold=True, clr=C_CYAN)
    rect(sl, Inches(6.45), Inches(1.50), Inches(3.20), Inches(0.025), fill=C_BORDER)

    components = [
        (C_CYAN,   "Backbone (PaliGemma 3B)",
         "paligemma-3b-pt-224\nSigLIP SO400M + Gemma-2B\n전체 파라미터 BF16 학습", 1.58),
        (C_ORANGE, "Action Expert",
         "4-layer Transformer\n8-head cross-attention\nhidden_dim = 256", 2.58),
        (C_GREEN,  "Flow Matching (CFM)",
         "L = ||v_θ − (x₁−x₀)||²\nx₀~N(0,I), x₁=GT chunk\nODE n_steps=5", 3.50),
        (C_PURPLE, "추론 파이프라인",
         "FastAPI (GX10)\n→ Action Buffer (Jetson)\n→ 50Hz 로컬 제어", 4.35),
    ]
    for ec, ctitle, cbody, iy in components:
        rect(sl, Inches(6.45), Inches(iy), Inches(0.05), Inches(0.70), fill=ec)
        sans(sl, ctitle,
             Inches(6.55), Inches(iy), Inches(3.10), Inches(0.28),
             sz=11, bold=True, clr=ec)
        sans(sl, cbody,
             Inches(6.55), Inches(iy + 0.28), Inches(3.10), Inches(0.55),
             sz=9, clr=C_MUTED)


def s_adaln(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "π0 정통 Action Expert: AdaLN-Zero 시간 컨디셔닝",
        "단순 덧셈 → scale/shift/gate 분리 · Zero-init으로 학습 안정성 보장")

    # BEFORE card
    card(sl, Inches(0.469), Inches(1.10), Inches(4.375), Inches(4.30),
         border_clr=C_RED, border_w=2.0, fill=C_CARD)
    rect(sl, Inches(0.469), Inches(1.10), Inches(4.375), Inches(0.06), fill=C_RED)

    mono(sl, "BEFORE  단순 덧셈",
         Inches(0.60), Inches(1.22), Inches(4.10), Inches(0.38),
         sz=14, bold=True, clr=C_RED)

    rect(sl, Inches(0.60), Inches(1.70), Inches(4.10), Inches(1.40), fill=C_CARD2)
    rect(sl, Inches(0.60), Inches(1.70), Inches(0.04), Inches(1.40), fill=C_RED)
    mono_lines(sl, [
        ("t_emb = time_mlp(t)        # (B,1,512)", False, C_TEXT),
        ("h = action_proj(x_t) + t_emb  # 단순 덧셈", False, C_TEXT),
        ("", False, C_TEXT),
        ("# query만 norm, key/value는 raw", False, C_MUTED),
        ("attn = self_attn(norm(h), h, h)", False, C_TEXT),
        ("h = h + attn               # gate 없음", False, C_TEXT),
        ("h = h + mlp(norm(h))", False, C_TEXT),
    ], Inches(0.65), Inches(1.72), Inches(4.05), Inches(1.36), sz=8.5)

    sans_lines(sl, [
        ("문제점:", True, C_RED),
        ("- timestep이 단순 offset으로만 작용", False, C_TEXT),
        ("- 각 레이어 조건 강도 조절 불가", False, C_TEXT),
        ("- 학습 초기 불안정 (초기값 편향)", False, C_TEXT),
    ], Inches(0.60), Inches(3.20), Inches(4.10), Inches(1.00), sz=11)

    rect(sl, Inches(0.60), Inches(4.30), Inches(4.10), Inches(0.65), fill=C_RED)
    sans(sl, "논문과 불일치 — DiT/π0 방식 미적용",
         Inches(0.60), Inches(4.30), Inches(4.10), Inches(0.65),
         sz=11, bold=True, clr=C_WHITE, align=CENTER)

    # Arrow
    mono(sl, "→",
         Inches(4.88), Inches(2.80), Inches(0.40), Inches(0.55),
         sz=30, bold=True, clr=C_ORANGE, align=CENTER)

    # AFTER card
    card(sl, Inches(5.156), Inches(1.10), Inches(4.375), Inches(4.30),
         border_clr=C_GREEN, border_w=2.0, fill=C_CARD)
    rect(sl, Inches(5.156), Inches(1.10), Inches(4.375), Inches(0.06), fill=C_GREEN)

    mono(sl, "AFTER  AdaLN-Zero  (π0 정통)",
         Inches(5.28), Inches(1.22), Inches(4.10), Inches(0.38),
         sz=14, bold=True, clr=C_GREEN)

    rect(sl, Inches(5.28), Inches(1.70), Inches(4.10), Inches(1.85), fill=C_CARD2)
    rect(sl, Inches(5.28), Inches(1.70), Inches(0.04), Inches(1.85), fill=C_GREEN)
    mono_lines(sl, [
        ("cond_emb = TimestepEmbedder(t) # (B,512)", False, C_TEXT),
        ("α1,β1,γ1,α2,β2,γ2 = AdaLNMod(cond_emb)", False, C_TEXT),
        ("", False, C_TEXT),
        ("# scale + shift + gate (self-attn)", False, C_MUTED),
        ("h += γ1 * self_attn(norm(h)*(1+α1)+β1)", False, C_TEXT),
        ("# VLM cross-attn (timestep 미적용)", False, C_MUTED),
        ("h += cross_attn(norm(h), vlm_cond)", False, C_TEXT),
        ("# scale + shift + gate (mlp)", False, C_MUTED),
        ("h += γ2 * mlp(norm(h)*(1+α2)+β2)", False, C_TEXT),
    ], Inches(5.33), Inches(1.72), Inches(4.05), Inches(1.81), sz=8.0)

    sans_lines(sl, [
        ("핵심 개선:", True, C_GREEN),
        ("- Zero-init → 학습 초기 identity 보장", False, C_TEXT),
        ("- 레이어별 독립적 scale/shift/gate", False, C_TEXT),
        ("- DiT(2023), π0(2024) 동일 방식", False, C_TEXT),
    ], Inches(5.28), Inches(3.65), Inches(4.10), Inches(0.90), sz=11)

    rect(sl, Inches(5.28), Inches(4.30), Inches(4.10), Inches(0.65), fill=C_GREEN)
    mono(sl, "검증 완료  verify_mona_expert.py: Max Error < 1e-6",
         Inches(5.28), Inches(4.30), Inches(4.10), Inches(0.65),
         sz=10, bold=True, clr=C_WHITE, align=CENTER)


def s06_dataset(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "데이터셋 및 수집 파이프라인",
        "HDF5 150 ep · 9카테고리 · 비동기 수집 (Camera 10Hz / Teleop 50Hz)")

    # ── Left column ───────────────────────────────────────────────────────────
    card(sl, Inches(0.469), Inches(1.10), Inches(4.30), Inches(4.40),
         border_clr=C_BORDER, fill=C_CARD)

    sans(sl, "데이터셋 구성",
         Inches(0.62), Inches(1.18), Inches(3.90), Inches(0.30),
         sz=13, bold=True, clr=C_CYAN)

    stats = [
        ("150", "에피소드 (HDF5)",  "Train 120 / Val 30 episode-level split", 1.60),
        ("9",   "이동 카테고리",     "center/left/right × straight/left/right", 2.20),
        ("15×", "Paraphrase Pool", "카테고리당 15개 · 매 step 무작위 선택",   2.80),
    ]
    for num, bold_t, muted_t, iy in stats:
        mono(sl, num,
             Inches(0.62), Inches(iy), Inches(0.65), Inches(0.40),
             sz=18, bold=True, clr=C_CYAN)
        sans(sl, bold_t,
             Inches(1.32), Inches(iy), Inches(3.20), Inches(0.28),
             sz=12, bold=True, clr=C_TEXT)
        sans(sl, muted_t,
             Inches(1.32), Inches(iy + 0.28), Inches(3.20), Inches(0.25),
             sz=9.5, clr=C_MUTED)

    rect(sl, Inches(0.62), Inches(3.38), Inches(3.95), Inches(0.025), fill=C_BORDER)

    sans(sl, "비동기 수집 파이프라인 (Jetson 16GB)",
         Inches(0.62), Inches(3.48), Inches(3.90), Inches(0.28),
         sz=11, bold=True, clr=C_ORANGE)

    rect(sl, Inches(0.62), Inches(3.80), Inches(3.95), Inches(1.00), fill=C_CARD2)
    rect(sl, Inches(0.62), Inches(3.80), Inches(0.04), Inches(1.00), fill=C_ORANGE)
    mono_lines(sl, [
        ("[Camera]  10Hz  →  frame queue (timestamp)", False, C_ORANGE),
        ("[Teleop]  50Hz  →  action queue (timestamp)", False, C_ORANGE),
        ("  ↓  timestamp 정렬 (±50ms nearest)", False, C_TEXT),
        ("Episode Manager  →  HDF5 일괄 저장", False, C_TEXT),
    ], Inches(0.68), Inches(3.82), Inches(3.89), Inches(0.96), sz=9.5)

    # ── Right column ──────────────────────────────────────────────────────────
    # Hyperparameters table
    card(sl, Inches(4.95), Inches(1.10), Inches(4.60), Inches(2.10),
         border_clr=C_BORDER, fill=C_CARD)
    sans(sl, "학습 하이퍼파라미터",
         Inches(5.08), Inches(1.18), Inches(4.20), Inches(0.30),
         sz=13, bold=True, clr=C_CYAN)

    hp = [
        ("window_size", "8 프레임",        C_TEXT),
        ("horizon",     "10 액션",         C_TEXT),
        ("batch_size",  "4",               C_TEXT),
        ("lr",          "1e-4 Cosine+warm", C_TEXT),
        ("dtype",       "BF16",            C_ORANGE),
        ("optimizer",   "AdamW",           C_TEXT),
    ]
    for i, (k, v, vc) in enumerate(hp):
        ry = Inches(1.58 + i * 0.24)
        if i % 2 == 0:
            rect(sl, Inches(4.95), ry, Inches(4.60), Inches(0.24), fill=C_CARD2)
        mono(sl, k, Inches(5.08), ry + Inches(0.03), Inches(2.00), Inches(0.20), sz=10, clr=C_MUTED)
        mono(sl, v, Inches(7.30), ry + Inches(0.03), Inches(2.00), Inches(0.20), sz=10, bold=True, clr=vc)

    # Category images from sample_imgs
    sans(sl, "카테고리별 실제 주행 이미지",
         Inches(4.95), Inches(3.30), Inches(4.60), Inches(0.28),
         sz=11, bold=True, clr=C_CYAN)

    img_specs = [
        ("center_straight_mid.jpg", "center", 4.97),
        ("left_left_mid.jpg",       "left_*", 6.48),
        ("right_right_mid.jpg",     "right_*", 7.99),
    ]
    for fname, label, ix in img_specs:
        img_path = SAMPLE_DIR / fname
        if img_path.exists():
            sl.shapes.add_picture(str(img_path),
                                  Inches(ix), Inches(3.62),
                                  Inches(1.42), Inches(1.72))
        else:
            card(sl, Inches(ix), Inches(3.62), Inches(1.42), Inches(1.72),
                 border_clr=C_BORDER, fill=C_CARD2)
            mono(sl, f"(사진)\n{label}", Inches(ix), Inches(3.90),
                 Inches(1.42), Inches(1.00), sz=9, clr=C_MUTED, align=CENTER)
        sans(sl, label, Inches(ix), Inches(5.33), Inches(1.42), Inches(0.22),
             sz=9, clr=C_MUTED, align=CENTER)


def s07_ablation(prs, buf_ab):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "Ablation Study — Offline 평가 (Val n=53)",
        "Instruction 기여도 검증 · FP16 배포 불가 확정")

    # Table
    hdr_defs = [("실험", Inches(0.469), Inches(1.50)),
                ("설명", Inches(1.967), Inches(1.70)),
                ("FPE ↓", Inches(3.667), Inches(1.17))]
    row_h = Inches(0.51)
    for htxt, hx, hw in hdr_defs:
        rect(sl, hx, Inches(1.10), hw, row_h, fill=C_CYAN)
        mono(sl, htxt, hx + Inches(0.05), Inches(1.10) + Inches(0.05),
             hw - Inches(0.10), row_h - Inches(0.10),
             sz=12, bold=True, clr=C_CARD, align=CENTER)

    data_rows = [
        (["E3  v2 best", "기준 모델",  "0.673"], None),
        (["E4  text-off", "text=0 벡터", "1.085"], 'red'),
        (["E6  FP16",    "BF16→FP16", "2.093"], 'red'),
    ]
    col_xs = [Inches(0.469), Inches(1.967), Inches(3.667)]
    col_ws = [Inches(1.498), Inches(1.700), Inches(1.170)]

    for i, (cells, hl) in enumerate(data_rows):
        ry = Inches(1.61 + i * 0.51)
        fill = C_RED if hl == 'red' else C_CARD2 if i % 2 == 0 else C_CARD
        if hl == 'red':
            rect(sl, Inches(0.469), ry, Inches(4.40), Inches(0.51), fill=fill)
        for txt, cx, cw in zip(cells, col_xs, col_ws):
            tc = C_WHITE if hl == 'red' else C_TEXT
            mono(sl, txt, cx + Inches(0.05), ry + Inches(0.05),
                 cw - Inches(0.10), Inches(0.41),
                 sz=11, bold=(hl == 'red'), clr=tc, align=CENTER)

    sl.shapes.add_picture(buf_ab, Inches(0.469), Inches(3.26), Inches(4.37), Inches(2.20))

    # Right interpretation card
    card(sl, Inches(5.10), Inches(1.10), Inches(4.43), Inches(4.36),
         border_clr=C_BORDER, fill=C_CARD)

    sans(sl, "핵심 해석",
         Inches(5.28), Inches(1.22), Inches(4.05), Inches(0.38),
         sz=15, bold=True, clr=C_CYAN)

    interps = [
        (C_GREEN,  "💡 Instruction 유효성",
         "text-off → FPE +61% 증가\n언어 명령이 궤적 생성에 필수적",  2.00),
        (C_RED,    "💡 FP16 배포 불가 확정",
         "FPE 3.1× 증가 · CL Success 0%\n정밀도 손실 → 성능 완전 붕괴", 3.10),
        (C_ORANGE, "💡 BF16 배포 확정",
         "GX10 & Jetson BF16 네이티브 지원\nFP32 대비 메모리 절반",     4.20),
    ]
    for ec, title, detail, iy in interps:
        rect(sl, Inches(5.28), Inches(iy), Inches(0.05), Inches(0.75), fill=ec)
        sans(sl, title,
             Inches(5.38), Inches(iy), Inches(3.95), Inches(0.32),
             sz=12, bold=True, clr=ec)
        sans(sl, detail,
             Inches(5.38), Inches(iy + 0.32), Inches(3.95), Inches(0.50),
             sz=10.5, clr=C_MUTED)


def s_metrics(prs, buf_traj, buf_sweep):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "평가 지표 정의 및 측정 방법",
        "FPE + Success@T · Bird's-eye 궤적 시각화 · Closed-Loop가 진짜 지표")

    # Left: metrics
    card(sl, Inches(0.469), Inches(1.10), Inches(4.10), Inches(1.40),
         border_clr=C_CYAN, border_w=1.5, fill=C_CARD)
    mono(sl, "FPE (Final Position Error)",
         Inches(0.62), Inches(1.18), Inches(3.80), Inches(0.35),
         sz=13, bold=True, clr=C_CYAN)
    rect(sl, Inches(0.62), Inches(1.58), Inches(3.80), Inches(0.40), fill=C_CARD2)
    mono(sl, "|| p_pred − p_gt ||₂",
         Inches(0.62), Inches(1.58), Inches(3.80), Inches(0.40),
         sz=12, clr=C_ORANGE, align=CENTER)
    sans(sl, "최종 위치 유클리드 거리 오차",
         Inches(0.62), Inches(2.06), Inches(3.80), Inches(0.30),
         sz=10, clr=C_MUTED)

    card(sl, Inches(0.469), Inches(2.58), Inches(4.10), Inches(1.40),
         border_clr=C_ORANGE, border_w=1.5, fill=C_CARD)
    mono(sl, "Success@T",
         Inches(0.62), Inches(2.66), Inches(3.80), Inches(0.35),
         sz=13, bold=True, clr=C_ORANGE)
    rect(sl, Inches(0.62), Inches(3.06), Inches(3.80), Inches(0.40), fill=C_CARD2)
    mono(sl, "FPE < T  (Primary: T=1.0)",
         Inches(0.62), Inches(3.06), Inches(3.80), Inches(0.40),
         sz=12, clr=C_CYAN, align=CENTER)
    sans(sl, "GT mean disp(9.42) 대비 10.6% 허용 오차",
         Inches(0.62), Inches(3.54), Inches(3.80), Inches(0.30),
         sz=10, clr=C_MUTED)

    sl.shapes.add_picture(buf_sweep, Inches(0.469), Inches(4.08), Inches(4.10), Inches(1.40))

    # Right: trajectory chart
    sans(sl, "Bird's-eye 궤적 시각화 (5개 시나리오, v3-A)",
         Inches(4.80), Inches(1.10), Inches(4.85), Inches(0.30),
         sz=12, bold=True, clr=C_CYAN)
    card(sl, Inches(4.80), Inches(1.44), Inches(4.85), Inches(3.25),
         border_clr=C_BORDER, fill=C_CARD)
    sl.shapes.add_picture(buf_traj,
                          Inches(4.86), Inches(1.50), Inches(4.73), Inches(3.13))

    card(sl, Inches(4.80), Inches(4.74), Inches(4.85), Inches(0.80),
         border_clr=C_GREEN, border_w=1.5, fill=C_CARD, left_accent=C_GREEN)
    sans(sl, "Closed-Loop(CL) vs Offline 차이",
         Inches(4.98), Inches(4.80), Inches(4.50), Inches(0.30),
         sz=12, bold=True, clr=C_GREEN)
    sans(sl, "Offline: 단발 스텝  →  CL: 오차 누적 → 진짜 성능 지표",
         Inches(4.98), Inches(5.10), Inches(4.50), Inches(0.30),
         sz=10, clr=C_MUTED)


def s08_results(prs, buf_cmp):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "주요 결과: 폐루프 시뮬레이션 (Val 30 ep)",
        "v3-A ★  Success@1.0 = 86.7%  |  Val Loss와 CL 성능 역전 현상")

    # Table
    col_defs = [("모델",         Inches(0.469), Inches(1.10)),
                ("Val Loss",     Inches(1.579), Inches(0.98)),
                ("FPE ↓",        Inches(2.559), Inches(0.86)),
                ("Success@1.0 ↑",Inches(3.419), Inches(1.83))]
    row_h = Inches(0.44)
    for htxt, hx, hw in col_defs:
        rect(sl, hx, Inches(1.10), hw, row_h, fill=C_CYAN)
        mono(sl, htxt, hx + Inches(0.04), Inches(1.10) + Inches(0.03),
             hw - Inches(0.08), row_h - Inches(0.06),
             sz=11, bold=True, clr=C_CARD, align=CENTER)

    data_rows = [
        (["Random",  "—",      "13.490", "—"    ], None),
        (["v2 best", "0.0619", "0.857",  "~79%" ], None),
        (["v3",      "0.0682", "0.740",  "76.7%"], None),
        (["v3-A ★",  "0.0714", "0.731",  "86.7%"], 'orange'),
    ]
    col_xs = [Inches(0.469), Inches(1.579), Inches(2.559), Inches(3.419)]
    col_ws = [Inches(1.110), Inches(0.980), Inches(0.860), Inches(1.830)]

    for i, (cells, hl) in enumerate(data_rows):
        ry = Inches(1.54 + i * 0.44)
        if hl == 'orange':
            rect(sl, Inches(0.469), ry, Inches(5.25), Inches(0.44), fill=C_ORANGE)
            tc = C_WHITE; tb = True
        else:
            fill = C_CARD2 if i % 2 == 0 else C_CARD
            rect(sl, Inches(0.469), ry, Inches(5.25), Inches(0.44), fill=fill)
            tc = C_TEXT; tb = False
        for txt, cx, cw in zip(cells, col_xs, col_ws):
            mono(sl, txt, cx + Inches(0.04), ry + Inches(0.04),
                 cw - Inches(0.08), Inches(0.36),
                 sz=11, bold=tb, clr=tc, align=CENTER)

    # Table ends at: 1.54 + 4*0.44 = 3.30 → chart from 3.38
    sl.shapes.add_picture(buf_cmp, Inches(0.469), Inches(3.38), Inches(4.80), Inches(2.10))

    # Right cards
    card(sl, Inches(5.55), Inches(1.10), Inches(3.99), Inches(1.32),
         border_clr=C_GREEN, border_w=1.5, fill=C_CARD, left_accent=C_GREEN)
    sans(sl, "v3-A 성능 도약",
         Inches(5.70), Inches(1.18), Inches(3.65), Inches(0.35),
         sz=14, bold=True, clr=C_GREEN)
    box, tf = _tb_raw(sl, Inches(5.70), Inches(1.53), Inches(3.65), Inches(0.72))
    mixed_para(tf, [("Instruction 다양화 적용 후  ", False, C_TEXT, FS),
                    ("Success +10%p", True, C_GREEN, FS),
                    (" (86.7%)", False, C_TEXT, FS)], sz=12, first=True)

    card(sl, Inches(5.55), Inches(2.55), Inches(3.99), Inches(2.20),
         border_clr=C_RED, border_w=1.5, fill=C_CARD, left_accent=C_RED)
    sans(sl, "⚠️  Val Loss ≠ CL 성능",
         Inches(5.70), Inches(2.63), Inches(3.65), Inches(0.35),
         sz=14, bold=True, clr=C_RED)
    box, tf = _tb_raw(sl, Inches(5.70), Inches(3.05), Inches(3.65), Inches(1.55))
    mixed_para(tf, [("v2 Loss 가장 낮음 → 하지만", False, C_TEXT, FS)], sz=12, first=True)
    mixed_para(tf, [("실제 CL 성능은 v3-A가 최고", True, C_TEXT, FS)], sz=12)
    mixed_para(tf, [("Episode split + Instruction 다양화가", False, C_MUTED, FS)], sz=11)
    mixed_para(tf, [("단순 Val Loss보다 실질 일반화 개선", True, C_ORANGE, FS)], sz=11)

    card(sl, Inches(5.55), Inches(4.85), Inches(3.99), Inches(0.68),
         border_clr=C_CYAN, fill=C_CARD, left_accent=C_CYAN)
    mono(sl, "→ Closed-Loop Success가 진짜 평가 지표",
         Inches(5.70), Inches(4.95), Inches(3.80), Inches(0.42),
         sz=11, bold=True, clr=C_CYAN)


def s09_categories(prs, buf_cat):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "카테고리별 성능 분석 (v3-A, n=30)",
        "center_left 0% — 모든 버전 공통 실패  |  시각적 모호성 분석")

    col_defs = [("카테고리",    Inches(0.469), Inches(1.52)),
                ("n",           Inches(1.989), Inches(0.35)),
                ("FPE ↓",       Inches(2.339), Inches(0.87)),
                ("Success@1.0", Inches(3.209), Inches(1.65))]
    row_h = Inches(0.355)
    for htxt, hx, hw in col_defs:
        rect(sl, hx, Inches(1.10), hw, row_h, fill=C_CYAN)
        mono(sl, htxt, hx + Inches(0.03), Inches(1.10) + Inches(0.02),
             hw - Inches(0.06), row_h - Inches(0.04),
             sz=9.5, bold=True, clr=C_CARD, align=CENTER)

    cat_rows = [
        ("center_straight", "4", "0.582", "25%",  None),
        ("center_left",     "3", "1.055", "0%",   'red'),
        ("center_right",    "3", "1.328", "33%",  None),
        ("left_straight",   "4", "0.782", "25%",  None),
        ("left_left",       "3", "0.451", "67%",  'green'),
        ("left_right",      "3", "0.747", "0%",   'red'),
        ("right_straight",  "4", "0.650", "50%",  None),
        ("right_left",      "3", "0.572", "33%",  None),
        ("right_right",     "3", "0.471", "67%",  'green'),
    ]
    col_xs = [Inches(0.469), Inches(1.989), Inches(2.339), Inches(3.209)]
    col_ws = [Inches(1.520), Inches(0.350), Inches(0.870), Inches(1.650)]
    total_w = Inches(4.39)

    # 9 rows × 0.355 + header 0.355 = 3.55 total → ends at 1.10+3.55 = 4.65
    for i, (cat, n, fpe, s10, hl) in enumerate(cat_rows):
        ry = Inches(1.455 + i * 0.355)
        if hl == 'red':
            rect(sl, Inches(0.469), ry, total_w, row_h, fill=C_RED)
            tc = C_WHITE; tb = True
        elif hl == 'green':
            rect(sl, Inches(0.469), ry, total_w, row_h, fill=C_GREEN)
            tc = C_CARD; tb = True
        else:
            fill = C_CARD2 if i % 2 == 0 else C_CARD
            rect(sl, Inches(0.469), ry, total_w, row_h, fill=fill)
            tc = C_TEXT; tb = False
        for txt, cx, cw in zip([cat, n, fpe, s10], col_xs, col_ws):
            mono(sl, txt, cx + Inches(0.03), ry + Inches(0.03),
                 cw - Inches(0.06), row_h - Inches(0.04),
                 sz=9.5, bold=tb, clr=tc, align=CENTER)

    # Right: analysis
    card(sl, Inches(5.10), Inches(1.10), Inches(4.43), Inches(3.80),
         border_clr=C_BORDER, fill=C_CARD)
    sans(sl, "center_left  0% 실패 원인 분석",
         Inches(5.28), Inches(1.22), Inches(4.05), Inches(0.35),
         sz=14, bold=True, clr=C_RED)

    issues = [
        (C_RED,    "시각적 모호성",
         "fish-eye 카메라에서 left/right\n구별이 어려운 구도 발생", 1.72),
        (C_ORANGE, "통계적 불확실성",
         "val n=3으로 샘플 수 부족\n평가 결과 신뢰도 한계",       2.52),
        (C_MUTED,  "Instruction 다양화 한계",
         "언어 표현 다양화만으로는\n시각 모호성 해결 불가능",      3.28),
    ]
    for ec, title, detail, iy in issues:
        rect(sl, Inches(5.28), Inches(iy), Inches(0.05), Inches(0.60), fill=ec)
        sans(sl, title,
             Inches(5.40), Inches(iy), Inches(3.85), Inches(0.28),
             sz=12, bold=True, clr=ec)
        sans(sl, detail,
             Inches(5.40), Inches(iy + 0.28), Inches(3.85), Inches(0.42),
             sz=10, clr=C_MUTED)

    # Solution block
    rect(sl, Inches(5.28), Inches(4.08), Inches(4.05), Inches(0.72), fill=C_CYAN)
    sans(sl, "해결책 (H6)",
         Inches(5.35), Inches(4.12), Inches(3.95), Inches(0.28),
         sz=12, bold=True, clr=C_CARD)
    sans(sl, "다양한 시작 위치 에피소드 추가 수집 → v4 재학습",
         Inches(5.35), Inches(4.40), Inches(3.95), Inches(0.32),
         sz=11, bold=True, clr=C_CARD)

    # Category chart (bottom) — table ends at ~4.65, chart from 4.70
    sl.shapes.add_picture(buf_cat,
                          Inches(0.469), Inches(4.72), Inches(9.06), Inches(0.82))


def s10_server(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "추론 서버 검증 — ASUS ASCENT GX10 (GB10 Superchip)",
        "Warm 평균 ~243ms  <  4Hz 예산 250ms  →  실시간 배포 가능 확인")

    # Left: specs
    card(sl, Inches(0.469), Inches(1.10), Inches(3.99), Inches(4.40),
         border_clr=C_BORDER, fill=C_CARD)
    sans(sl, "서버 스펙",
         Inches(0.62), Inches(1.22), Inches(3.65), Inches(0.32),
         sz=14, bold=True, clr=C_CYAN)
    rect(sl, Inches(0.62), Inches(1.54), Inches(3.65), Inches(0.025), fill=C_BORDER)

    specs = [
        ("GPU",   "GB10 Superchip",   C_CYAN),
        ("모델",   "v3-A BF16",        C_TEXT),
        ("VRAM",  "5.97 GB",           C_GREEN),
        ("CUDA",  "13.0",              C_TEXT),
        ("API",   "FastAPI",           C_TEXT),
        ("캐싱",   "비전 KV 캐싱 (TTL)", C_ORANGE),
    ]
    for i, (k, v, vc) in enumerate(specs):
        ry = Inches(1.62 + i * 0.48)
        sans(sl, k, Inches(0.62), ry, Inches(1.30), Inches(0.38), sz=12, clr=C_MUTED)
        mono(sl, v, Inches(2.10), ry, Inches(2.30), Inches(0.38), sz=12, bold=True, clr=vc)

    # Jetson compat note
    rect(sl, Inches(0.62), Inches(4.50), Inches(3.65), Inches(0.80), fill=C_CARD2)
    sans(sl, "Jetson AGX Orin 호환성",
         Inches(0.72), Inches(4.55), Inches(3.45), Inches(0.28),
         sz=11, bold=True, clr=C_GREEN)
    sans(sl, "5.97 GB → AGX Orin 16GB 여유 충분\nBF16 네이티브 지원 확인",
         Inches(0.72), Inches(4.83), Inches(3.45), Inches(0.40),
         sz=10, clr=C_MUTED)

    # Right: latency table
    lat_cols = [("시나리오", Inches(4.65), Inches(1.73)),
                ("Latency",  Inches(6.38), Inches(1.10)),
                ("상태",     Inches(7.48), Inches(2.05))]
    row_h = Inches(0.51)
    for htxt, hx, hw in lat_cols:
        rect(sl, hx, Inches(1.10), hw, row_h, fill=C_CYAN)
        mono(sl, htxt, hx + Inches(0.04), Inches(1.10) + Inches(0.04),
             hw - Inches(0.08), row_h - Inches(0.08),
             sz=12, bold=True, clr=C_CARD, align=CENTER)

    lat_rows = [
        (["cold start",  "885 ms",  "[x] KV 미채움 (~4배)"], 'red'),
        (["center_left", "238 ms",  "[ok] 예산 내"],          None),
        (["center_right","232 ms",  "[ok] 예산 내"],          None),
        (["left_str",    "248 ms",  "[ok] 예산 내"],          None),
        (["warm 평균",   "~243 ms", "[ok] 4Hz 이내"],         'green'),
    ]
    col_xs = [Inches(4.65), Inches(6.38), Inches(7.48)]
    col_ws = [Inches(1.73), Inches(1.10), Inches(2.05)]

    for i, (cells, hl) in enumerate(lat_rows):
        ry = Inches(1.61 + i * 0.51)
        if hl == 'red':
            rect(sl, Inches(4.65), ry, Inches(3.88), Inches(0.51), fill=C_RED)
            tc = C_WHITE; tb = True
        elif hl == 'green':
            rect(sl, Inches(4.65), ry, Inches(3.88), Inches(0.51), fill=C_GREEN)
            tc = C_CARD; tb = True
        else:
            fill = C_CARD2 if i % 2 == 0 else C_CARD
            rect(sl, Inches(4.65), ry, Inches(3.88), Inches(0.51), fill=fill)
            tc = C_TEXT; tb = False
        for txt, cx, cw in zip(cells, col_xs, col_ws):
            mono(sl, txt, cx + Inches(0.04), ry + Inches(0.05),
                 cw - Inches(0.08), Inches(0.41),
                 sz=11, bold=tb, clr=tc, align=CENTER)

    # Deployment diagram
    card(sl, Inches(4.65), Inches(4.27), Inches(3.88), Inches(1.23),
         border_clr=C_CYAN, border_w=1.5, fill=C_CARD, left_accent=C_CYAN)
    mono_lines(sl, [
        ("[GX10 서버]  →  /predict API  →  (B,10,3) chunk",  False, C_CYAN),
        ("       ↓",                                           False, C_TEXT),
        ("[Jetson AGX] → Action Buffer → 50Hz 로컬 제어",     False, C_TEXT),
        ("            → ROS2 → 실로봇 actuator",              False, C_MUTED),
    ], Inches(4.80), Inches(4.38), Inches(3.65), Inches(1.00), sz=9.5)


def s11_discussion(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "한계 및 토의",
        "정직한 한계 인정 + 원인 분석 + 해결 방향 제시")

    rows = [
        (C_RED,    "❌", "center_left  0%  실패",
         "시각적 모호성 및 통계적 불확실성 (val n=3)",
         "→ 다양한 시작 위치 에피소드 추가 수집 (H6 milestone)",
         1.10),
        (C_ORANGE, "⚠️", "Val Loss ≠ CL 성능 불일치",
         "단순 Loss보다 Downstream 일반화가 실질적 성능 지표",
         "→ 평가 기준: Closed-Loop Success 확정 사용",
         2.60),
        (C_CYAN,   "⏳", "실로봇 미배포 상태",
         "현재 시뮬레이션 폐루프(CL)만 검증 완료된 상태",
         "→ GX10 구동 완료 / Jetson 코드 준비 완료 / 배포 대기",
         4.10),
    ]
    for ec, icon, title, muted, action, iy in rows:
        card(sl, Inches(0.469), Inches(iy), Inches(9.06), Inches(1.30),
             border_clr=ec, border_w=1.5, fill=C_CARD, left_accent=ec)
        sans(sl, icon,
             Inches(0.65), Inches(iy + 0.38), Inches(0.55), Inches(0.50),
             sz=28, clr=ec)
        sans(sl, title,
             Inches(1.40), Inches(iy + 0.10), Inches(7.60), Inches(0.40),
             sz=16, bold=True, clr=ec)
        sans(sl, muted,
             Inches(1.40), Inches(iy + 0.52), Inches(7.60), Inches(0.30),
             sz=12, clr=C_MUTED)
        sans(sl, action,
             Inches(1.40), Inches(iy + 0.88), Inches(7.60), Inches(0.30),
             sz=12, bold=True, clr=C_TEXT)


def s12_conclusion(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "결론 및 향후 계획",
        "무엇을 했고 · 다음에 무엇을 할 것인가")

    # Left: contributions
    card(sl, Inches(0.469), Inches(1.00), Inches(4.45), Inches(3.10),
         border_clr=C_GREEN, border_w=1.5, fill=C_CARD)
    rect(sl, Inches(0.469), Inches(1.00), Inches(4.45), Inches(0.06), fill=C_GREEN)
    sans(sl, "주요 기여 (완료)",
         Inches(0.62), Inches(1.12), Inches(4.15), Inches(0.35),
         sz=15, bold=True, clr=C_GREEN)

    contribs = [
        ("✅", "VLA 파이프라인 구현",
         "PaliGemma + Action Expert + Flow Matching 통합", 1.60),
        ("✅", "데이터 전략 개선",
         "Episode split + Instruction Pool → Success +10%p",  2.10),
        ("✅", "배포 검증 완료",
         "GX10 243ms ✅  GPU 5.97GB ✅  4Hz 예산 내",          2.60),
        ("✅", "실용적 진단",
         "BF16 확정 · Threshold 명확화 · center_left 병목 발견", 3.10),
    ]
    for icon, bold_t, detail_t, iy in contribs:
        sans(sl, icon,
             Inches(0.62), Inches(iy), Inches(0.28), Inches(0.40),
             sz=13, clr=C_GREEN)
        box, tf = _tb_raw(sl, Inches(0.98), Inches(iy), Inches(3.75), Inches(0.42))
        mixed_para(tf, [(bold_t, True, C_TEXT, FS),
                        ("\n" + detail_t, False, C_MUTED, FS)], sz=11, first=True)

    # Right: future plans
    card(sl, Inches(5.10), Inches(1.00), Inches(4.45), Inches(3.10),
         border_clr=C_ORANGE, border_w=1.5, fill=C_CARD)
    rect(sl, Inches(5.10), Inches(1.00), Inches(4.45), Inches(0.06), fill=C_ORANGE)
    sans(sl, "향후 계획",
         Inches(5.25), Inches(1.12), Inches(4.15), Inches(0.35),
         sz=15, bold=True, clr=C_ORANGE)

    plans = [
        ("📌", "단기 (D10-D11)",
         "Jetson 실로봇 배포 · 실환경 주행 10+ ep · Success 측정", 1.60),
        ("📌", "중기 (H6)",
         "center_left 에피소드 추가 수집 · v4 재학습 및 평가",     2.14),
        ("📌", "장기 (논문 완성)",
         "실로봇 결과 업데이트 · MoNaVLA 비교 · 시연 영상",        2.68),
    ]
    for icon, bold_t, detail_t, iy in plans:
        sans(sl, icon,
             Inches(5.25), Inches(iy), Inches(0.28), Inches(0.40),
             sz=13, clr=C_ORANGE)
        box, tf = _tb_raw(sl, Inches(5.62), Inches(iy), Inches(3.75), Inches(0.52))
        mixed_para(tf, [(bold_t, True, C_TEXT, FS),
                        ("\n" + detail_t, False, C_MUTED, FS)], sz=11, first=True)

    # Bottom goal bar
    card(sl, Inches(0.469), Inches(4.18), Inches(9.06), Inches(1.25),
         border_clr=C_CYAN, border_w=1.5, fill=C_CARD, left_accent=C_CYAN)
    sans(sl, "최종 목표",
         Inches(0.65), Inches(4.24), Inches(1.50), Inches(0.28),
         sz=12, bold=True, clr=C_CYAN)
    sans(sl, "4Hz 재계획 + 50Hz 로컬 제어  →  실로봇 완전 자율 내비게이션 실현",
         Inches(0.65), Inches(4.52), Inches(8.70), Inches(0.38),
         sz=18, bold=True, clr=C_TEXT)
    sans(sl, "MoNa-π  ·  Mobile Navigation π0  ·  Flow Matching VLA",
         Inches(0.65), Inches(4.92), Inches(8.70), Inches(0.28),
         sz=11, clr=C_MUTED)


# ── Entry point ────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    os.makedirs("/home/minum/26CS/MoNa-pi/reports", exist_ok=True)

    print("[1/8] chart_architecture...")
    buf_arch  = chart_architecture()
    print("[2/8] chart_flow_matching_concept...")
    buf_flow  = chart_flow_matching_concept()
    print("[3/8] chart_ablation...")
    buf_ab    = chart_ablation()
    print("[4/8] chart_comparison...")
    buf_cmp   = chart_comparison()
    print("[5/8] chart_categories...")
    buf_cat   = chart_categories()
    print("[6/8] chart_trajectories...")
    buf_traj  = chart_trajectories()
    print("[7/8] chart_threshold_sweep...")
    buf_sweep = chart_threshold_sweep()

    print("[8/8] Building slides...")
    prs = new_prs()
    s01_title(prs)                          # Slide 01
    s_flow_matching(prs)                    # Slide 02 — Flow Matching 개념
    s02_background(prs)                     # Slide 03 — 연구 배경
    s03_design_focus(prs)                   # Slide 04 — 설계 포인트
    s04_approach(prs)                       # Slide 05 — 3가지 전환
    s05_architecture(prs, buf_arch)         # Slide 06 — 아키텍처 블록 다이어그램
    s_adaln(prs)                            # Slide 07 — AdaLN-Zero
    s06_dataset(prs)                        # Slide 08 — 데이터셋 (실제 이미지)
    s07_ablation(prs, buf_ab)               # Slide 09 — Ablation
    s_metrics(prs, buf_traj, buf_sweep)     # Slide 10 — 평가 지표
    s08_results(prs, buf_cmp)               # Slide 11 — 주요 결과
    s09_categories(prs, buf_cat)            # Slide 12 — 카테고리 분석
    s10_server(prs)                         # Slide 13 — 추론 서버
    s11_discussion(prs)                     # Slide 14 — 한계
    s12_conclusion(prs)                     # Slide 15 — 결론

    out = "/home/minum/26CS/MoNa-pi/reports/mona_pi_v2.pptx"
    prs.save(out)
    print(f"\nSaved: {out}  ({len(prs.slides)} slides)")
