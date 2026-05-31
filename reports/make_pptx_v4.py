#!/usr/bin/env python3
"""MoNa-pi 중간발표 PPTX v4 — 2026-05-28
7 Slides · ~5분 · 평가기준(창의성30/차별성25/기획25/실현성20) 정렬
"""

import io, os
from pathlib import Path
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SW = Inches(10.0)
SH = Inches(5.625)

ROBOT_FRONT = Path(__file__).parent / "robot_front.jpg"
ROBOT_TOP   = Path(__file__).parent / "robot_top.jpg"

FS = "맑은 고딕"
FC = "Consolas"

BG_SLIDE  = RGBColor(0xF7, 0xF9, 0xFC)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD2   = RGBColor(0xF0, 0xF5, 0xFF)
C_TEXT    = RGBColor(0x11, 0x18, 0x27)
C_MUTED   = RGBColor(0x6B, 0x72, 0x80)
C_NAVY    = RGBColor(0x1D, 0x4E, 0xD8)
C_GREEN   = RGBColor(0x16, 0xA3, 0x4A)
C_RED     = RGBColor(0xDC, 0x26, 0x26)
C_ORANGE  = RGBColor(0xEA, 0x58, 0x0C)
C_PURPLE  = RGBColor(0x71, 0x27, 0xCE)
C_AMBER   = RGBColor(0xD9, 0x77, 0x06)
C_BLUE_T  = RGBColor(0xEB, 0xF3, 0xFF)
C_GREEN_T = RGBColor(0xDC, 0xFC, 0xEA)
C_RED_T   = RGBColor(0xFE, 0xE2, 0xE2)
C_AMBER_T = RGBColor(0xFF, 0xF7, 0xED)
C_PURPLE_T= RGBColor(0xF3, 0xE8, 0xFF)
C_BORDER  = RGBColor(0xE5, 0xE7, 0xEB)
C_BDR2    = RGBColor(0x93, 0xC5, 0xFD)

BG_PLT    = '#FFFFFF'
NAVY_PLT  = '#1D4ED8'
GREEN_PLT = '#16A34A'
RED_PLT   = '#DC2626'
ORANGE_PLT= '#EA580C'
PURPLE_PLT= '#7127CE'
TEXT_PLT  = '#111827'
MUTED_PLT = '#9CA3AF'
GRID_PLT  = '#E5E7EB'

LEFT   = PP_ALIGN.LEFT
CENTER = PP_ALIGN.CENTER


# ── Primitives ────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(sl):
    f = sl.background.fill
    f.solid()
    f.fore_color.rgb = BG_SLIDE

def rect(sl, x, y, w, h, fill=None, line=None, lw=0):
    s = sl.shapes.add_shape(1, x, y, w, h)
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line and lw:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb(sl, x, y, w, h, wrap=True):
    box = sl.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    return box, tf

def sans(sl, text, x, y, w, h, sz=12, bold=False, clr=None, align=LEFT, wrap=True):
    if clr is None: clr = C_TEXT
    _, tf = _tb(sl, x, y, w, h, wrap)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = clr; r.font.name = FS

def mono(sl, text, x, y, w, h, sz=12, bold=False, clr=None, align=LEFT, wrap=True):
    if clr is None: clr = C_TEXT
    _, tf = _tb(sl, x, y, w, h, wrap)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = clr; r.font.name = FC

def card(sl, x, y, w, h, border_clr=None, border_w=1.0, fill=None, accent=None):
    if fill is None:       fill       = C_WHITE
    if border_clr is None: border_clr = C_BORDER
    rect(sl, x, y, w, h, fill=fill, line=border_clr, lw=border_w)
    if accent:
        rect(sl, x, y, Inches(0.055), h, fill=accent)

def hdr(sl, title, subtitle=None):
    rect(sl, Inches(0.469), Inches(0.240), Inches(0.055), Inches(0.680), fill=C_NAVY)
    sans(sl, title, Inches(0.563), Inches(0.260), Inches(9.062), Inches(0.420),
         sz=20, bold=True, clr=C_TEXT)
    if subtitle:
        sans(sl, subtitle, Inches(0.563), Inches(0.680), Inches(9.062), Inches(0.240),
             sz=10, clr=C_MUTED)
    rect(sl, Inches(0.469), Inches(1.010), Inches(9.062), Inches(0.010), fill=C_BORDER)

def savefig(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0); plt.close(fig)
    return buf

def _ax(ax):
    ax.set_facecolor(BG_PLT)
    for sp in ax.spines.values(): sp.set_color(GRID_PLT)
    ax.tick_params(colors=TEXT_PLT, labelsize=8)
    ax.yaxis.grid(True, color=GRID_PLT, linewidth=0.7, linestyle='--')
    ax.set_axisbelow(True)


# ── Charts ────────────────────────────────────────────────────────────────────

def chart_comparison():
    fig, axes = plt.subplots(1, 2, figsize=(5.2, 2.8), facecolor=BG_PLT)

    ax1 = axes[0]; _ax(ax1)
    models = ['Baseline\n(v2)', 'MoNa-π\n(v3-A)']
    fpe = [0.857, 0.731]
    bars = ax1.bar(models, fpe, color=[MUTED_PLT, NAVY_PLT], width=0.5, edgecolor='none', alpha=0.85)
    for b, v in zip(bars, fpe):
        ax1.text(b.get_x()+b.get_width()/2, v+0.01, f'{v:.3f}',
                 ha='center', va='bottom', color=TEXT_PLT, fontsize=9)
    ax1.set_title('FPE  (lower is better)', color=TEXT_PLT, fontsize=9, pad=4)
    ax1.set_ylim(0, 1.1)
    ax1.annotate('−14.7%', xy=(1, 0.731), xytext=(0.55, 0.62),
                 color=NAVY_PLT, fontsize=9, fontweight='bold',
                 arrowprops=dict(arrowstyle='->', color=NAVY_PLT, lw=1.1))

    ax2 = axes[1]; _ax(ax2)
    s10 = [79.0, 86.7]
    bars2 = ax2.bar(models, s10, color=[MUTED_PLT, ORANGE_PLT], width=0.5, edgecolor='none', alpha=0.85)
    for b, v in zip(bars2, s10):
        ax2.text(b.get_x()+b.get_width()/2, v+0.5, f'{v:.1f}%',
                 ha='center', va='bottom', color=TEXT_PLT, fontsize=9)
    ax2.set_title('Success@1.0  ↑', color=TEXT_PLT, fontsize=9, pad=4)
    ax2.set_ylim(0, 108)
    ax2.annotate('+9.7%p', xy=(1, 86.7), xytext=(0.5, 94),
                 color=ORANGE_PLT, fontsize=9, fontweight='bold',
                 arrowprops=dict(arrowstyle='->', color=ORANGE_PLT, lw=1.1))

    fig.tight_layout(pad=1.4)
    return savefig(fig)


def annotated_robot_hw():
    """robot_top.jpg with hardware component callout labels"""
    import matplotlib.image as mpimg
    img = mpimg.imread(str(ROBOT_TOP))
    H, W = img.shape[:2]  # 1024 × 576

    fig, ax = plt.subplots(figsize=(3.2, 4.8), facecolor=BG_PLT)
    ax.imshow(img, extent=[0, W, H, 0])
    ax.set_xlim(-W * 0.58, W * 1.60)
    ax.set_ylim(H * 1.06, -H * 0.06)
    ax.axis('off')

    # (label, arrow_tip_x, arrow_tip_y, text_x, text_y, color)
    callouts = [
        ("Fish-eye Camera ×2\n8-frame / Wide-angle",
         105, 330,  -W*0.28, 280, GREEN_PLT),
        ("LiDAR",
         285, 195,   W*1.32, 130, NAVY_PLT),
        ("Jetson Orin AGX\n16 GB  BF16",
         350, 455,   W*1.32, 420, ORANGE_PLT),
        ("Touch Display",
         285, 640,   W*1.32, 640, PURPLE_PLT),
        ("Omni Wheel\n3-DOF",
         440, 865,   W*1.32, 850, RED_PLT),
    ]
    for label, tx, ty, lx, ly, color in callouts:
        ax.annotate(
            label, xy=(tx, ty), xytext=(lx, ly),
            fontsize=7, fontweight='bold', color='white',
            ha='center', va='center',
            bbox=dict(boxstyle='round,pad=0.35', fc=color,
                      ec='white', lw=0.8, alpha=0.92),
            arrowprops=dict(arrowstyle='->', color=color, lw=1.3,
                            connectionstyle='arc3,rad=0.05'))

    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
    return savefig(fig)


def chart_gantt():
    fig, ax = plt.subplots(figsize=(5.6, 2.2), facecolor=BG_PLT)
    ax.set_facecolor(BG_PLT)
    for sp in ax.spines.values(): sp.set_color(GRID_PLT)

    phases = [
        ('Data Collection & Pipeline', 1,  3, GREEN_PLT),
        ('Model Training v2/v3',       3,  7, NAVY_PLT),
        ('Validation & Ablation',      7, 10, ORANGE_PLT),
        ('Serbot2 Deploy',            10, 14, PURPLE_PLT),
    ]
    for i, (label, s, e, col) in enumerate(phases):
        y = len(phases)-1-i
        ax.barh(y, e-s, left=s, height=0.55, color=col, alpha=0.80, edgecolor='none')
        ax.text((s+e)/2, y, f'D{s}–D{e}', ha='center', va='center',
                color='white', fontsize=8, fontweight='bold')
        ax.text(s-0.15, y, label, ha='right', va='center', color=TEXT_PLT, fontsize=8)

    ax.axvline(x=7, color=RED_PLT, lw=1.5, linestyle='--', alpha=0.8)
    ax.text(7.15, 3.35, 'Now(D7)', color=RED_PLT, fontsize=7.5, fontweight='bold')

    ax.set_yticks([])
    ax.set_xticks(range(1, 15))
    ax.set_xticklabels([f'D{i}' for i in range(1, 15)], color=MUTED_PLT, fontsize=7)
    ax.set_xlim(0, 15); ax.set_ylim(-0.4, 3.8)
    ax.xaxis.grid(True, color=GRID_PLT, linewidth=0.5, linestyle='--')
    ax.set_axisbelow(True)
    ax.tick_params(axis='y', length=0)
    ax.set_title('Project Schedule (D1–D14)', color=TEXT_PLT, fontsize=9, pad=4)
    fig.tight_layout(pad=0.6)
    return savefig(fig)


# ── Slides ────────────────────────────────────────────────────────────────────

def s01_title(prs):
    sl = blank(prs); set_bg(sl)
    rect(sl, 0, 0, SW, Inches(0.22), fill=C_NAVY)
    rect(sl, 0, SH-Inches(0.22), SW, Inches(0.22), fill=C_NAVY)

    sans(sl, "MoNa-π",
         Inches(0.60), Inches(0.48), Inches(5.80), Inches(1.50),
         sz=80, bold=True, clr=C_NAVY)

    card(sl, Inches(0.60), Inches(2.16), Inches(6.00), Inches(0.64),
         border_clr=C_BDR2, border_w=1.2, fill=C_WHITE)
    sans(sl, "Flow Matching 기반 고주파 모바일 내비게이션 VLA",
         Inches(0.76), Inches(2.27), Inches(5.70), Inches(0.42),
         sz=15, bold=True, clr=C_TEXT)

    rect(sl, Inches(0.60), Inches(2.96), Inches(2.0), Inches(0.018), fill=C_BORDER)
    sans(sl, "중간발표  ·  2026. 05. 14",
         Inches(0.60), Inches(3.04), Inches(5.0), Inches(0.30),
         sz=12, clr=C_MUTED)

    card(sl, Inches(0.60), Inches(3.46), Inches(5.20), Inches(1.92),
         border_clr=C_BORDER, fill=C_WHITE)
    sans(sl, "Team Monaf  ·  인공지능전공",
         Inches(0.78), Inches(3.56), Inches(4.80), Inches(0.30),
         sz=11, bold=True, clr=C_NAVY)
    rect(sl, Inches(0.78), Inches(3.88), Inches(4.80), Inches(0.015), fill=C_BORDER)

    members = [
        ("이민우", "모델 아키텍처 설계 · 학습 파이프라인",  3.96),
        ("정재연", "데이터셋 수집 파이프라인 구축",          4.40),
        ("오은석", "데이터셋 수집 및 인프라 관리",           4.84),
    ]
    for name, role, iy in members:
        sans(sl, name, Inches(0.78), Inches(iy), Inches(0.90), Inches(0.32),
             sz=12, bold=True, clr=C_TEXT)
        sans(sl, role, Inches(1.76), Inches(iy), Inches(3.80), Inches(0.32),
             sz=11, clr=C_MUTED)

    # Serbot2 hardware photo (right half)
    sl.shapes.add_picture(str(ROBOT_TOP),
                          Inches(6.30), Inches(0.28), Inches(3.30), Inches(5.10))


def s02_creativity(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "아이디어의 창의성 — 기존 한계 → MoNa-π 해결책",
        "이산 분류 VLA의 3가지 근본 한계를 Flow Matching + VLA 통합으로 극복  [창의성 30점]")

    # Left panel: 기존 한계
    card(sl, Inches(0.469), Inches(1.08), Inches(4.28), Inches(4.38),
         border_clr=C_RED, border_w=1.4, fill=C_WHITE)
    rect(sl, Inches(0.469), Inches(1.08), Inches(4.28), Inches(0.07), fill=C_RED)
    sans(sl, "기존 VLA의 한계",
         Inches(0.62), Inches(1.20), Inches(4.00), Inches(0.34),
         sz=14, bold=True, clr=C_RED)

    problems = [
        (C_RED,    C_RED_T,
         "① 이산 분류 제어",
         "동작을 클래스(왼쪽/직진/오른쪽)로 분류\n→ 연속적 궤적 생성 근본적으로 불가능"),
        (C_ORANGE, C_AMBER_T,
         "② 단발 예측 Jitter",
         "매 스텝 독립 예측 → 행동 불일치·진동\n→ 실로봇에서 불안정한 주행 패턴"),
        (C_PURPLE, C_PURPLE_T,
         "③ 고정 명령어 과적합",
         "특정 문장에만 반응 → 미등록 표현 실패\n→ 실사용 자연어 다양성 처리 불가"),
    ]
    for i, (ec, bg, title, detail) in enumerate(problems):
        iy = Inches(1.66 + i * 1.10)
        card(sl, Inches(0.62), iy, Inches(3.98), Inches(0.96),
             border_clr=ec, border_w=1.0, fill=bg, accent=ec)
        sans(sl, title, Inches(0.78), iy+Inches(0.07), Inches(3.65), Inches(0.28),
             sz=12, bold=True, clr=ec)
        sans(sl, detail, Inches(0.78), iy+Inches(0.38), Inches(3.65), Inches(0.52),
             sz=10, clr=C_MUTED)

    # Center arrow
    sans(sl, "→", Inches(4.75), Inches(2.60), Inches(0.48), Inches(0.50),
         sz=30, bold=True, clr=C_NAVY, align=CENTER)

    # Right panel: 해결책
    card(sl, Inches(5.25), Inches(1.08), Inches(4.28), Inches(4.38),
         border_clr=C_GREEN, border_w=1.4, fill=C_WHITE)
    rect(sl, Inches(5.25), Inches(1.08), Inches(4.28), Inches(0.07), fill=C_GREEN)
    sans(sl, "MoNa-π 해결책",
         Inches(5.40), Inches(1.20), Inches(4.00), Inches(0.34),
         sz=14, bold=True, clr=C_GREEN)

    solutions = [
        (C_GREEN,  C_GREEN_T,
         "① Flow Matching 연속 경로",
         "연속 벡터 필드로 직접 궤적 학습\n5스텝 ODE 복원 → 실시간 연속 제어"),
        (C_NAVY,   C_BLUE_T,
         "② Action Chunking (h=10)",
         "미래 10스텝 동시 예측 후 버퍼 실행\n4Hz 재계획 / 50Hz 로컬 제어 분리"),
        (C_AMBER,  C_AMBER_T,
         "③ Instruction Pool",
         "카테고리당 15개 Paraphrase 무작위 선택\n다양한 자연어 표현에 강건한 일반화"),
    ]
    for i, (ec, bg, title, detail) in enumerate(solutions):
        iy = Inches(1.66 + i * 1.10)
        card(sl, Inches(5.40), iy, Inches(3.98), Inches(0.96),
             border_clr=ec, border_w=1.0, fill=bg, accent=ec)
        sans(sl, title, Inches(5.56), iy+Inches(0.07), Inches(3.65), Inches(0.28),
             sz=12, bold=True, clr=ec)
        sans(sl, detail, Inches(5.56), iy+Inches(0.38), Inches(3.65), Inches(0.52),
             sz=10, clr=C_MUTED)

    # Bottom insight bar
    card(sl, Inches(0.469), Inches(5.04), Inches(9.06), Inches(0.44),
         border_clr=C_NAVY, border_w=1.2, fill=C_BLUE_T, accent=C_NAVY)
    sans(sl, "독창성: PaliGemma (언어이해) + Flow Matching (연속제어) + Serbot2 (실배포) — 세 요소의 통합이 핵심",
         Inches(0.62), Inches(5.08), Inches(8.82), Inches(0.36),
         sz=11, bold=True, clr=C_NAVY)


def s03_differentiation(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "기존 기술과의 차별성 — 3가지 접근법 비교",
        "이산 VLA · Diffusion Policy · MoNa-π  |  5가지 핵심 기준  [차별성 25점]")

    col_xs = [Inches(0.469), Inches(3.06), Inches(5.48), Inches(7.55)]
    col_ws = [Inches(2.59),  Inches(2.40), Inches(2.05), Inches(2.05)]
    headers = ["평가 기준", "이산 분류 VLA", "Diffusion Policy", "MoNa-π  ★"]
    h_fills = [RGBColor(0xF1,0xF5,0xF9), RGBColor(0xF8,0xFA,0xFC),
               RGBColor(0xF8,0xFA,0xFC), C_BLUE_T]
    h_clrs  = [C_TEXT, C_MUTED, C_MUTED, C_NAVY]

    HDR_H = Inches(0.52)
    for hx, hw, ht, hf, hc in zip(col_xs, col_ws, headers, h_fills, h_clrs):
        rect(sl, hx, Inches(1.08), hw-Inches(0.03), HDR_H, fill=hf)
        sans(sl, ht, hx+Inches(0.10), Inches(1.11), hw-Inches(0.20), HDR_H-Inches(0.08),
             sz=12, bold=True, clr=hc, align=CENTER)

    ROW_H = Inches(0.69)
    rows = [
        ("제어 방식",
         ("이산 클래스",           C_RED,    C_RED_T),
         ("연속 확률 분포",         C_ORANGE, C_AMBER_T),
         ("연속 벡터 필드  ✓",     C_GREEN,  C_GREEN_T)),
        ("추론 속도",
         ("빠름 (1 forward)",       C_MUTED,  C_WHITE),
         ("느림 (1000+ 스텝)",      C_RED,    C_RED_T),
         ("빠름 (5스텝 ODE)  ✓",   C_GREEN,  C_GREEN_T)),
        ("언어 이해",
         ("제한 (고정 명령)",        C_ORANGE, C_AMBER_T),
         ("없음 (순수 모방)",        C_RED,    C_RED_T),
         ("PaliGemma 3B  ✓",        C_GREEN,  C_GREEN_T)),
        ("일반화",
         ("낮음 (과적합 위험)",      C_RED,    C_RED_T),
         ("중간",                    C_MUTED,  C_WHITE),
         ("Instr. Pool 강건  ✓",    C_GREEN,  C_GREEN_T)),
        ("실로봇 배포",
         ("가능 (경량)",             C_MUTED,  C_WHITE),
         ("어려움 (메모리 과다)",    C_RED,    C_RED_T),
         ("Serbot2 BF16  ✓",        C_GREEN,  C_GREEN_T)),
    ]

    for ri, (criterion, v1, v2, v3) in enumerate(rows):
        ry = Inches(1.60 + ri * 0.72)
        row_fill = C_CARD2 if ri % 2 == 0 else C_WHITE
        rect(sl, col_xs[0], ry, col_ws[0]-Inches(0.03), ROW_H, fill=row_fill)
        sans(sl, criterion, col_xs[0]+Inches(0.10), ry+Inches(0.10),
             col_ws[0]-Inches(0.20), ROW_H-Inches(0.20),
             sz=12, bold=True, clr=C_TEXT)
        for ci, (val, vc, vbg) in enumerate([v1, v2, v3]):
            cx = col_xs[ci+1]; cw = col_ws[ci+1]
            cell_fill = C_BLUE_T if ci == 2 else vbg
            rect(sl, cx, ry, cw-Inches(0.03), ROW_H, fill=cell_fill)
            clr = C_NAVY if ci == 2 else vc
            sans(sl, val, cx+Inches(0.08), ry+Inches(0.10),
                 cw-Inches(0.16), ROW_H-Inches(0.20),
                 sz=10, bold=(ci == 2), clr=clr, align=CENTER)

    card(sl, Inches(0.469), Inches(5.21), Inches(9.06), Inches(0.28),
         border_clr=C_NAVY, border_w=1.0, fill=C_BLUE_T, accent=C_NAVY)
    sans(sl, "MoNa-π는 언어 이해 + 연속 제어 + 실배포를 동시 달성한 유일한 접근법",
         Inches(0.62), Inches(5.23), Inches(8.82), Inches(0.24),
         sz=10.5, bold=True, clr=C_NAVY)


def s04_planning(prs, buf_gantt):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "과제 기획 및 추진전략 — 체계적 접근",
        "D1–D14 단계별 일정 · 3인 역할 분담 · 정량 평가 기반 반복 개선  [과제기획 25점]")

    # Left: Gantt chart
    sl.shapes.add_picture(buf_gantt,
                          Inches(0.469), Inches(1.10), Inches(5.52), Inches(2.38))

    # Left bottom: methodology
    card(sl, Inches(0.469), Inches(3.60), Inches(5.52), Inches(1.84),
         border_clr=C_BORDER, fill=C_WHITE)
    sans(sl, "연구 방법론 — 논리적 문제 해결 순서",
         Inches(0.62), Inches(3.70), Inches(5.20), Inches(0.30),
         sz=12, bold=True, clr=C_NAVY)
    rect(sl, Inches(0.62), Inches(4.02), Inches(5.20), Inches(0.014), fill=C_BORDER)

    methods = [
        ("①", "데이터 수집",  "비동기 파이프라인 · HDF5 · 카테고리 균형"),
        ("②", "반복 학습",    "v2→v3→v3-A 에피소드 분할 · BF16 전환"),
        ("③", "정량 평가",    "FPE + CL Success@1.0 이중 지표 적용"),
        ("④", "배포 검증",    "GX10 latency · Serbot2 VRAM 호환 확인"),
    ]
    for i, (num, step, detail) in enumerate(methods):
        iy = Inches(4.10 + i * 0.32)
        sans(sl, num,  Inches(0.62), iy, Inches(0.22), Inches(0.28),
             sz=10, bold=True, clr=C_NAVY, align=CENTER)
        sans(sl, step, Inches(0.88), iy, Inches(1.00), Inches(0.28),
             sz=10, bold=True, clr=C_TEXT)
        sans(sl, detail, Inches(1.94), iy, Inches(3.90), Inches(0.28),
             sz=10, clr=C_MUTED)

    # Right: Role blocks
    sans(sl, "역할 분담",
         Inches(6.20), Inches(1.10), Inches(3.34), Inches(0.32),
         sz=13, bold=True, clr=C_NAVY)
    rect(sl, Inches(6.20), Inches(1.44), Inches(3.34), Inches(0.014), fill=C_BORDER)

    roles = [
        ("이민우", C_NAVY, C_BLUE_T,
         "모델 아키텍처 설계\n학습 파이프라인 구축\n평가 지표 분석 · GX10 검증"),
        ("정재연", C_GREEN, C_GREEN_T,
         "HDF5 수집 파이프라인 구현\n카테고리 데이터 균형 관리\n실험 재현성 보장"),
        ("오은석", C_ORANGE, C_AMBER_T,
         "데이터셋 실제 수집 실행\nSerbot2 인프라 관리\nROS2 통신 셋업"),
    ]
    ry = Inches(1.52)
    for name, ec, bg, tasks in roles:
        bh = Inches(1.14)
        card(sl, Inches(6.20), ry, Inches(3.34), bh,
             border_clr=ec, border_w=1.0, fill=bg, accent=ec)
        sans(sl, name, Inches(6.36), ry+Inches(0.07), Inches(3.00), Inches(0.28),
             sz=12, bold=True, clr=ec)
        sans(sl, tasks, Inches(6.36), ry+Inches(0.38), Inches(3.00), Inches(0.70),
             sz=9.5, clr=C_MUTED)
        ry += bh + Inches(0.06)


def s05_feasibility(prs, buf_hw):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "실현가능성 — 시제품 완성 & Serbot2 실배포 준비",
        "모델 학습 완료 · GX10 검증 완료 · Serbot2 (Jetson Orin AGX 16GB) 호환 확인  [실현가능성 20점]")

    # Left: annotated hardware photo
    sl.shapes.add_picture(buf_hw,
                          Inches(0.25), Inches(1.06), Inches(3.80), Inches(4.42))

    # Right col: camera info (top priority)
    card(sl, Inches(4.22), Inches(1.10), Inches(5.31), Inches(1.62),
         border_clr=C_GREEN, border_w=1.4, fill=C_WHITE)
    rect(sl, Inches(4.22), Inches(1.10), Inches(5.31), Inches(0.07), fill=C_GREEN)
    sans(sl, "카메라 입력 스펙",
         Inches(4.38), Inches(1.22), Inches(5.00), Inches(0.30),
         sz=13, bold=True, clr=C_GREEN)
    rect(sl, Inches(4.38), Inches(1.54), Inches(5.00), Inches(0.014), fill=C_BORDER)

    cam_specs = [
        ("센서",   "Fish-eye Wide-angle  ×2",    C_GREEN),
        ("입력",   "(B, 8, H, W, 3)  — 8-frame", C_NAVY),
        ("윈도우", "0.8 s  temporal window",      C_TEXT),
        ("해상도", "224 × 224  (SigLIP patch)",   C_TEXT),
    ]
    for i, (k, v, vc) in enumerate(cam_specs):
        ry = Inches(1.62 + i * 0.28)
        sans(sl, k, Inches(4.38), ry, Inches(1.00), Inches(0.26), sz=11, clr=C_MUTED)
        mono(sl, v, Inches(5.46), ry, Inches(3.88), Inches(0.26), sz=11, bold=True, clr=vc)

    # Deployment status (compact)
    card(sl, Inches(4.22), Inches(2.84), Inches(5.31), Inches(1.44),
         border_clr=C_NAVY, border_w=1.4, fill=C_WHITE)
    rect(sl, Inches(4.22), Inches(2.84), Inches(5.31), Inches(0.07), fill=C_NAVY)
    sans(sl, "배포 준비 상태",
         Inches(4.38), Inches(2.96), Inches(5.00), Inches(0.30),
         sz=13, bold=True, clr=C_NAVY)

    statuses = [
        ("[완료]  v3-A 학습 · Success@1.0  86.7%",      C_GREEN),
        ("[완료]  GX10 Warm 243ms  <  250ms 예산",       C_GREEN),
        ("[완료]  BF16 · Serbot2 16GB VRAM 여유 확인",   C_GREEN),
        ("[예정]  Serbot2 실배포  —  D10~D11",           C_ORANGE),
    ]
    for i, (txt, tc) in enumerate(statuses):
        sans(sl, txt, Inches(4.38), Inches(3.34 + i * 0.24),
             Inches(5.00), Inches(0.22), sz=10, bold=(tc == C_GREEN), clr=tc)

    # Pipeline card
    card(sl, Inches(4.22), Inches(4.40), Inches(5.31), Inches(0.84),
         border_clr=C_PURPLE, border_w=1.2, fill=C_PURPLE_T, accent=C_PURPLE)
    sans(sl, "추론 파이프라인",
         Inches(4.38), Inches(4.48), Inches(2.00), Inches(0.26),
         sz=10, bold=True, clr=C_PURPLE)
    mono(sl, "Camera  →  GX10 FastAPI  →  (B,10,3) chunk\n"
             "       →  Serbot2 Buffer  →  50Hz 실행",
         Inches(4.38), Inches(4.76), Inches(5.00), Inches(0.42),
         sz=10, clr=C_TEXT)


def s06_results(prs, buf_cmp):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "핵심 결과 — 성능 검증 완료",
        "v3-A: Success@1.0 86.7%  ·  FPE 0.731  ·  추론 Warm 243ms  ·  4Hz 예산 내")

    # Top metric cards
    metrics = [
        ("86.7%",  "Success@1.0",    "v3-A  (Threshold=1.0)",    C_ORANGE, C_AMBER_T),
        ("243 ms", "추론 레이턴시",   "Warm 평균  <  250ms 예산",  C_GREEN,  C_GREEN_T),
        ("5.97 GB","모델 VRAM",       "BF16 / Serbot2 16GB 여유", C_NAVY,   C_BLUE_T),
    ]
    for i, (val, label, sub, ec, bg) in enumerate(metrics):
        mx = Inches(0.469 + i * 3.02)
        card(sl, mx, Inches(1.10), Inches(2.86), Inches(1.22),
             border_clr=ec, border_w=1.4, fill=bg)
        sans(sl, val,   mx+Inches(0.14), Inches(1.18), Inches(2.58), Inches(0.60),
             sz=32, bold=True, clr=ec, align=CENTER)
        sans(sl, label, mx+Inches(0.14), Inches(1.78), Inches(2.58), Inches(0.26),
             sz=11, bold=True, clr=ec, align=CENTER)
        sans(sl, sub,   mx+Inches(0.14), Inches(2.04), Inches(2.58), Inches(0.24),
             sz=9, clr=C_MUTED, align=CENTER)

    # Bottom left: comparison chart
    sl.shapes.add_picture(buf_cmp,
                          Inches(0.469), Inches(2.44), Inches(5.00), Inches(3.00))

    # Bottom right: robot_front — real navigation scene
    sl.shapes.add_picture(str(ROBOT_FRONT),
                          Inches(5.58), Inches(2.44), Inches(3.95), Inches(2.72))
    card(sl, Inches(5.58), Inches(5.10), Inches(3.95), Inches(0.34),
         border_clr=C_NAVY, fill=C_NAVY)
    sans(sl, "Serbot2 실험 환경  ·  자율 내비게이션 시연",
         Inches(5.66), Inches(5.14), Inches(3.79), Inches(0.26),
         sz=9, bold=True, clr=C_WHITE, align=CENTER)


def s07_conclusion(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "결론 — 연구 의의 및 파급효과",
        "Flow Matching VLA 실증 완료  ·  모바일 로봇 자율 내비게이션 표준화 가능성 확인")

    # Left: contributions
    card(sl, Inches(0.469), Inches(1.10), Inches(4.40), Inches(3.08),
         border_clr=C_GREEN, border_w=1.4, fill=C_WHITE)
    rect(sl, Inches(0.469), Inches(1.10), Inches(4.40), Inches(0.07), fill=C_GREEN)
    sans(sl, "완료된 기여",
         Inches(0.62), Inches(1.22), Inches(4.10), Inches(0.32),
         sz=14, bold=True, clr=C_GREEN)

    contribs = [
        ("VLA 파이프라인 구현",
         "PaliGemma + Action Expert + Flow Matching 통합",     1.64),
        ("데이터 전략 개선",
         "Episode Split + Instruction Pool → +9.7%p 향상",    2.18),
        ("배포 검증 완료",
         "GX10 243ms ✓  Serbot2 5.97GB ✓  4Hz 예산 내",        2.72),
        ("체계적 진단",
         "FPE+CL 이중 지표 · BF16 확정 · 병목 원인 분석",      3.26),
    ]
    for bold_t, detail, iy in contribs:
        rect(sl, Inches(0.62), Inches(iy+0.02), Inches(0.22), Inches(0.22), fill=C_GREEN)
        sans(sl, "✓", Inches(0.62), Inches(iy+0.01), Inches(0.22), Inches(0.22),
             sz=9, bold=True, clr=C_WHITE, align=CENTER)
        sans(sl, bold_t, Inches(0.90), Inches(iy),       Inches(3.80), Inches(0.24),
             sz=11, bold=True, clr=C_TEXT)
        sans(sl, detail,  Inches(0.90), Inches(iy+0.25), Inches(3.80), Inches(0.24),
             sz=10, clr=C_MUTED)

    # Right top: 파급효과
    card(sl, Inches(5.06), Inches(1.10), Inches(4.47), Inches(1.72),
         border_clr=C_ORANGE, border_w=1.4, fill=C_WHITE)
    rect(sl, Inches(5.06), Inches(1.10), Inches(4.47), Inches(0.07), fill=C_ORANGE)
    sans(sl, "사업화 & 파급효과",
         Inches(5.22), Inches(1.22), Inches(4.18), Inches(0.32),
         sz=14, bold=True, clr=C_ORANGE)

    impacts = [
        ("물류·창고",  "언어 명령 기반 자율 이동 → 물류 자동화"),
        ("서비스 로봇", "의료·호텔·공항 안내 로봇 일반화 적용"),
        ("오픈소스",   "Flow Matching VLA 프레임워크 공개"),
    ]
    for i, (sector, detail) in enumerate(impacts):
        iy = Inches(1.54 + i * 0.42)
        sans(sl, sector,  Inches(5.22), iy, Inches(1.14), Inches(0.28),
             sz=10.5, bold=True, clr=C_ORANGE)
        sans(sl, detail,  Inches(6.42), iy, Inches(2.98), Inches(0.28),
             sz=10.5, clr=C_MUTED)

    # Right bottom: 향후 계획
    card(sl, Inches(5.06), Inches(2.94), Inches(4.47), Inches(1.24),
         border_clr=C_NAVY, border_w=1.4, fill=C_WHITE)
    rect(sl, Inches(5.06), Inches(2.94), Inches(4.47), Inches(0.07), fill=C_NAVY)
    sans(sl, "향후 계획",
         Inches(5.22), Inches(3.06), Inches(4.18), Inches(0.32),
         sz=14, bold=True, clr=C_NAVY)

    plans = [
        ("단기 D10-D11", "Serbot2 실배포 · 실환경 10+ ep 주행"),
        ("중기 H6",       "center_left 보강 · v4 재학습 평가"),
        ("장기",          "실로봇 결과 논문화 · MoNaVLA 비교"),
    ]
    for i, (term, plan) in enumerate(plans):
        iy = Inches(3.38 + i * 0.26)
        sans(sl, term + ":", Inches(5.22), iy, Inches(1.14), Inches(0.24),
             sz=9.5, bold=True, clr=C_NAVY)
        sans(sl, plan,        Inches(6.42), iy, Inches(2.98), Inches(0.24),
             sz=9.5, clr=C_MUTED)

    # Bottom goal bar
    card(sl, Inches(0.469), Inches(4.30), Inches(9.06), Inches(1.10),
         border_clr=C_NAVY, border_w=1.2, fill=C_BLUE_T, accent=C_NAVY)
    sans(sl, "최종 목표",
         Inches(0.65), Inches(4.38), Inches(1.40), Inches(0.26),
         sz=11, bold=True, clr=C_NAVY)
    sans(sl, "4Hz 재계획 + 50Hz 로컬 제어  →  Serbot2 완전 자율 내비게이션 실현",
         Inches(0.65), Inches(4.66), Inches(8.76), Inches(0.50),
         sz=18, bold=True, clr=C_TEXT)
    sans(sl, "MoNa-π  ·  Mobile Navigation π0  ·  Flow Matching VLA",
         Inches(0.65), Inches(5.18), Inches(8.76), Inches(0.20),
         sz=9.5, clr=C_MUTED)


# ── Entry ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("[1/3] Generating charts...")
    buf_gantt = chart_gantt()
    buf_cmp   = chart_comparison()
    print("[2/3] Annotating hardware photo...")
    buf_hw    = annotated_robot_hw()

    print("[3/3] Building 7 slides...")
    prs = new_prs()
    s01_title(prs)
    s02_creativity(prs)
    s03_differentiation(prs)
    s04_planning(prs, buf_gantt)
    s05_feasibility(prs, buf_hw)
    s06_results(prs, buf_cmp)
    s07_conclusion(prs)

    out = "/home/minum/26CS/MoNa-pi/reports/mona_pi_v4.pptx"
    prs.save(out)
    print(f"\nSaved: {out}  ({len(prs.slides)} slides)")