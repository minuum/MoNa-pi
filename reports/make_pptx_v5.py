#!/usr/bin/env python3
"""MoNa-pi 캡스톤 발표 PPTX v5 — 2026-05-28
9 Slides · ~5분 · 교수 컨펌 구조
목차→기획의도/차별성→과제추진→시스템아키텍처→실험환경→핵심결과→기여/결론→감사합니다
"""

import io, os
from pathlib import Path
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.image as mpimg

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SW = Inches(10.0)
SH = Inches(5.625)

ROBOT_FRONT = Path(__file__).parent / "robot_front.jpg"
ROBOT_TOP   = Path(__file__).parent / "robot_top.jpg"

FS = "맑은 고딕"
FC = "Consolas"

BG_SLIDE  = RGBColor(0xF7, 0xF9, 0xFC)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD2   = RGBColor(0xF0, 0xF5, 0xFF)
C_TEXT    = RGBColor(0x11, 0x18, 0x27)
C_MUTED   = RGBColor(0x6B, 0x72, 0x80)
C_NAVY    = RGBColor(0x1D, 0x4E, 0xD8)
C_GREEN   = RGBColor(0x16, 0xA3, 0x4A)
C_RED     = RGBColor(0xDC, 0x26, 0x26)
C_ORANGE  = RGBColor(0xEA, 0x58, 0x0C)
C_PURPLE  = RGBColor(0x71, 0x27, 0xCE)
C_TEAL    = RGBColor(0x0D, 0x94, 0x88)
C_BLUE_T  = RGBColor(0xEB, 0xF3, 0xFF)
C_GREEN_T = RGBColor(0xDC, 0xFC, 0xEA)
C_RED_T   = RGBColor(0xFE, 0xE2, 0xE2)
C_AMBER_T = RGBColor(0xFF, 0xF7, 0xED)
C_PURPLE_T= RGBColor(0xF3, 0xE8, 0xFF)
C_TEAL_T  = RGBColor(0xCC, 0xFB, 0xF1)
C_BORDER  = RGBColor(0xE5, 0xE7, 0xEB)
C_BDR2    = RGBColor(0x93, 0xC5, 0xFD)

BG_PLT    = '#FFFFFF'
NAVY_PLT  = '#1D4ED8'
GREEN_PLT = '#16A34A'
RED_PLT   = '#DC2626'
ORANGE_PLT= '#EA580C'
PURPLE_PLT= '#7127CE'
TEAL_PLT  = '#0D9488'
TEXT_PLT  = '#111827'
MUTED_PLT = '#9CA3AF'
GRID_PLT  = '#E5E7EB'

LEFT   = PP_ALIGN.LEFT
CENTER = PP_ALIGN.CENTER
RIGHT  = PP_ALIGN.RIGHT


# ── Primitives ────────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def set_bg(sl):
    f = sl.background.fill
    f.solid(); f.fore_color.rgb = BG_SLIDE

def rect(sl, x, y, w, h, fill=None, line=None, lw=0):
    s = sl.shapes.add_shape(1, x, y, w, h)
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line and lw:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb(sl, x, y, w, h, wrap=True):
    box = sl.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame; tf.word_wrap = wrap
    return box, tf

def sans(sl, text, x, y, w, h, sz=12, bold=False, clr=None, align=LEFT, wrap=True):
    if clr is None: clr = C_TEXT
    _, tf = _tb(sl, x, y, w, h, wrap)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = clr; r.font.name = FS

def mono(sl, text, x, y, w, h, sz=12, bold=False, clr=None, align=LEFT, wrap=True):
    if clr is None: clr = C_TEXT
    _, tf = _tb(sl, x, y, w, h, wrap)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.color.rgb = clr; r.font.name = FC

def card(sl, x, y, w, h, border_clr=None, border_w=1.0, fill=None, accent=None):
    if fill is None:       fill       = C_WHITE
    if border_clr is None: border_clr = C_BORDER
    rect(sl, x, y, w, h, fill=fill, line=border_clr, lw=border_w)
    if accent:
        rect(sl, x, y, Inches(0.055), h, fill=accent)

def hdr(sl, title, subtitle=None):
    rect(sl, Inches(0.469), Inches(0.240), Inches(0.055), Inches(0.680), fill=C_NAVY)
    sans(sl, title, Inches(0.563), Inches(0.260), Inches(9.062), Inches(0.420),
         sz=20, bold=True, clr=C_TEXT)
    if subtitle:
        sans(sl, subtitle, Inches(0.563), Inches(0.680), Inches(9.062), Inches(0.240),
             sz=10, clr=C_MUTED)
    rect(sl, Inches(0.469), Inches(1.010), Inches(9.062), Inches(0.010), fill=C_BORDER)

def savefig(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=fig.get_facecolor())
    buf.seek(0); plt.close(fig)
    return buf

def _ax(ax):
    ax.set_facecolor(BG_PLT)
    for sp in ax.spines.values(): sp.set_color(GRID_PLT)
    ax.tick_params(colors=TEXT_PLT, labelsize=8)
    ax.yaxis.grid(True, color=GRID_PLT, linewidth=0.7, linestyle='--')
    ax.set_axisbelow(True)


# ── Charts ────────────────────────────────────────────────────────────────────

def annotated_robot_hw():
    img = mpimg.imread(str(ROBOT_TOP))
    H, W = img.shape[:2]  # 1024 × 576

    fig, ax = plt.subplots(figsize=(3.2, 4.8), facecolor=BG_PLT)
    ax.imshow(img, extent=[0, W, H, 0])
    ax.set_xlim(-W * 0.58, W * 1.60)
    ax.set_ylim(H * 1.06, -H * 0.06)
    ax.axis('off')

    callouts = [
        ("Fish-eye Camera ×2\n8-frame / Wide-angle",
         105, 330, -W*0.28, 280, GREEN_PLT),
        ("LiDAR",
         285, 195,  W*1.32, 130, NAVY_PLT),
        ("Jetson Orin AGX\n16 GB  BF16",
         350, 455,  W*1.32, 420, ORANGE_PLT),
        ("Touch Display",
         285, 640,  W*1.32, 640, PURPLE_PLT),
        ("Omni Wheel\n3-DOF",
         440, 865,  W*1.32, 850, RED_PLT),
    ]
    for label, tx, ty, lx, ly, color in callouts:
        ax.annotate(
            label, xy=(tx, ty), xytext=(lx, ly),
            fontsize=7, fontweight='bold', color='white',
            ha='center', va='center',
            bbox=dict(boxstyle='round,pad=0.35', fc=color, ec='white', lw=0.8, alpha=0.92),
            arrowprops=dict(arrowstyle='->', color=color, lw=1.3,
                            connectionstyle='arc3,rad=0.05'))

    fig.subplots_adjust(left=0, right=1, top=1, bottom=0)
    return savefig(fig)


def chart_comparison():
    fig, axes = plt.subplots(1, 2, figsize=(5.2, 2.8), facecolor=BG_PLT)

    ax1 = axes[0]; _ax(ax1)
    models = ['Baseline\n(v2)', 'MoNa-π\n(v3-A)']
    fpe = [0.857, 0.731]
    bars = ax1.bar(models, fpe, color=[MUTED_PLT, NAVY_PLT], width=0.5, edgecolor='none', alpha=0.85)
    for b, v in zip(bars, fpe):
        ax1.text(b.get_x()+b.get_width()/2, v+0.01, f'{v:.3f}',
                 ha='center', va='bottom', color=TEXT_PLT, fontsize=9)
    ax1.set_title('FPE  (lower is better)', color=TEXT_PLT, fontsize=9, pad=4)
    ax1.set_ylim(0, 1.1)
    ax1.annotate('-14.7%', xy=(1, 0.731), xytext=(0.55, 0.62),
                 color=NAVY_PLT, fontsize=9, fontweight='bold',
                 arrowprops=dict(arrowstyle='->', color=NAVY_PLT, lw=1.1))

    ax2 = axes[1]; _ax(ax2)
    s10 = [79.0, 86.7]
    bars2 = ax2.bar(models, s10, color=[MUTED_PLT, ORANGE_PLT], width=0.5, edgecolor='none', alpha=0.85)
    for b, v in zip(bars2, s10):
        ax2.text(b.get_x()+b.get_width()/2, v+0.5, f'{v:.1f}%',
                 ha='center', va='bottom', color=TEXT_PLT, fontsize=9)
    ax2.set_title('Success@1.0  (higher is better)', color=TEXT_PLT, fontsize=9, pad=4)
    ax2.set_ylim(0, 108)
    ax2.annotate('+9.7%p', xy=(1, 86.7), xytext=(0.5, 94),
                 color=ORANGE_PLT, fontsize=9, fontweight='bold',
                 arrowprops=dict(arrowstyle='->', color=ORANGE_PLT, lw=1.1))

    fig.tight_layout(pad=1.4)
    return savefig(fig)


# ── Slides ────────────────────────────────────────────────────────────────────

def s01_title(prs):
    sl = blank(prs); set_bg(sl)
    rect(sl, 0, 0, SW, Inches(0.22), fill=C_NAVY)
    rect(sl, 0, SH-Inches(0.22), SW, Inches(0.22), fill=C_NAVY)

    sans(sl, "MoNa-π",
         Inches(0.60), Inches(0.48), Inches(5.80), Inches(1.50),
         sz=80, bold=True, clr=C_NAVY)

    card(sl, Inches(0.60), Inches(2.16), Inches(6.00), Inches(0.64),
         border_clr=C_BDR2, border_w=1.2, fill=C_WHITE)
    sans(sl, "Flow Matching 기반 고주파 모바일 내비게이션 VLA",
         Inches(0.76), Inches(2.27), Inches(5.70), Inches(0.42),
         sz=15, bold=True, clr=C_TEXT)

    rect(sl, Inches(0.60), Inches(2.96), Inches(2.0), Inches(0.018), fill=C_BORDER)
    sans(sl, "2026 캡스톤 경진대회  ·  2026. 06. 04",
         Inches(0.60), Inches(3.04), Inches(5.0), Inches(0.30),
         sz=12, clr=C_MUTED)

    card(sl, Inches(0.60), Inches(3.46), Inches(5.20), Inches(1.92),
         border_clr=C_BORDER, fill=C_WHITE)
    sans(sl, "Team MONA-F  ·  인공지능전공",
         Inches(0.78), Inches(3.56), Inches(4.80), Inches(0.30),
         sz=11, bold=True, clr=C_NAVY)
    rect(sl, Inches(0.78), Inches(3.88), Inches(4.80), Inches(0.015), fill=C_BORDER)

    members = [
        ("임현우 (팀장)", "모델 아키텍처 설계 · 학습 파이프라인",  3.96),
        ("정재연",        "데이터셋 수집 파이프라인 구축",          4.40),
        ("오은석",        "데이터셋 수집 및 인프라 관리",           4.84),
    ]
    for name, role, iy in members:
        sans(sl, name, Inches(0.78), Inches(iy), Inches(1.10), Inches(0.32),
             sz=12, bold=True, clr=C_TEXT)
        sans(sl, role, Inches(1.94), Inches(iy), Inches(3.65), Inches(0.32),
             sz=11, clr=C_MUTED)

    sl.shapes.add_picture(str(ROBOT_TOP),
                          Inches(6.30), Inches(0.28), Inches(3.30), Inches(5.10))


def s02_toc(prs):
    sl = blank(prs); set_bg(sl)
    rect(sl, Inches(0.469), Inches(0.240), Inches(0.055), Inches(0.580), fill=C_NAVY)
    sans(sl, "목차",
         Inches(0.563), Inches(0.260), Inches(9.062), Inches(0.380),
         sz=20, bold=True, clr=C_TEXT)
    rect(sl, Inches(0.469), Inches(0.86), Inches(9.062), Inches(0.010), fill=C_BORDER)

    sections = [
        ("01", "연구 배경 및 차별성",
         "도메인 이전 동기 · 기존 기술 한계 · MoNa-π 접근법",
         C_NAVY, C_BLUE_T),
        ("02", "개발 과정",
         "데이터 수집 → 모델 학습 → 검증 → Serbot2 배포",
         C_GREEN, C_GREEN_T),
        ("03", "시스템 아키텍처",
         "Serbot2 하드웨어 · 추론 파이프라인 · 이중 제어 루프",
         C_ORANGE, C_AMBER_T),
        ("04", "실험 환경 및 평가",
         "9-카테고리 내비게이션 · 평가 지표 · 데이터셋",
         C_TEAL, C_TEAL_T),
        ("05", "실험 결과",
         "Success@1.0  86.7%  ·  추론 243ms  ·  4Hz 달성",
         C_PURPLE, C_PURPLE_T),
        ("06", "결론 및 기대 효과",
         "도메인 이전 성공 · 최적화 파이프라인 · 향후 적용",
         C_RED, C_RED_T),
    ]

    # 2 rows × 3 cols
    xs = [Inches(0.469), Inches(3.57), Inches(6.67)]
    ys = [Inches(0.96),  Inches(3.12)]
    bw, bh = Inches(2.96), Inches(1.98)

    for i, (num, title, desc, ec, bg) in enumerate(sections):
        col, row = i % 3, i // 3
        x, y = xs[col], ys[row]
        card(sl, x, y, bw, bh, border_clr=ec, border_w=1.3, fill=bg)
        rect(sl, x, y, bw, Inches(0.055), fill=ec)

        # Number badge
        rect(sl, x+Inches(0.14), y+Inches(0.13), Inches(0.48), Inches(0.48), fill=ec)
        sans(sl, num, x+Inches(0.14), y+Inches(0.13), Inches(0.48), Inches(0.48),
             sz=16, bold=True, clr=C_WHITE, align=CENTER, wrap=False)

        sans(sl, title, x+Inches(0.72), y+Inches(0.16), bw-Inches(0.82), Inches(0.40),
             sz=13, bold=True, clr=ec)
        sans(sl, desc,  x+Inches(0.14), y+Inches(0.72), bw-Inches(0.24), Inches(0.90),
             sz=10, clr=C_MUTED)

        # Row connector arrow (between cols)
        if col < 2:
            sans(sl, "→",
                 x+bw+Inches(0.02), y+Inches(0.74), Inches(0.09), Inches(0.50),
                 sz=18, bold=True, clr=C_MUTED, align=CENTER, wrap=False)

    # Down arrow between rows (right side)
    sans(sl, "↓", Inches(9.50), Inches(2.00), Inches(0.12), Inches(1.10),
         sz=20, bold=True, clr=C_MUTED, align=CENTER, wrap=False)


def s03_motivation(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "연구 배경 및 차별성",
        "모바일 매니플레이터 VLA → 모바일 로봇 도메인 이전  |  기존 기술 한계 극복")

    # Left: 기획 의도
    card(sl, Inches(0.469), Inches(1.08), Inches(4.52), Inches(4.38),
         border_clr=C_NAVY, border_w=1.4, fill=C_WHITE)
    rect(sl, Inches(0.469), Inches(1.08), Inches(4.52), Inches(0.07), fill=C_NAVY)
    sans(sl, "연구 배경",
         Inches(0.62), Inches(1.20), Inches(4.25), Inches(0.34),
         sz=14, bold=True, clr=C_NAVY)

    # Domain transfer concept
    card(sl, Inches(0.62), Inches(1.64), Inches(4.25), Inches(0.88),
         border_clr=C_BORDER, fill=C_CARD2)
    sans(sl, "기존 VLA",
         Inches(0.78), Inches(1.72), Inches(1.60), Inches(0.26),
         sz=10, bold=True, clr=C_MUTED)
    sans(sl, "모바일 매니플레이터\n(로봇 팔) 중심",
         Inches(0.78), Inches(1.98), Inches(1.60), Inches(0.42),
         sz=10, clr=C_MUTED)
    sans(sl, "→",
         Inches(2.48), Inches(1.88), Inches(0.38), Inches(0.40),
         sz=22, bold=True, clr=C_NAVY, align=CENTER, wrap=False)
    sans(sl, "MoNa-π",
         Inches(2.92), Inches(1.72), Inches(1.72), Inches(0.26),
         sz=10, bold=True, clr=C_NAVY)
    sans(sl, "모바일 로봇\n도메인 이전 적용",
         Inches(2.92), Inches(1.98), Inches(1.72), Inches(0.42),
         sz=10, bold=True, clr=C_NAVY)

    problems = [
        (C_RED,    C_RED_T,    "이산 분류 제어의 한계",
         "클래스 분류로는 연속 궤적 생성 불가\n→ Flow Matching 연속 벡터 필드 도입"),
        (C_ORANGE, C_AMBER_T,  "낮은 일반화 성능",
         "고정 명령어 과적합 · 미등록 표현 실패\n→ Paraphrase Pool (15개/카테고리)"),
        (C_PURPLE, C_PURPLE_T, "느린 추론 (Diffusion Policy)",
         "1000+ 스텝 필요 → 실시간 제어 불가\n→ Flow Matching 5스텝 ODE 복원"),
    ]
    for i, (ec, bg, title, detail) in enumerate(problems):
        iy = Inches(2.68 + i * 0.84)
        card(sl, Inches(0.62), iy, Inches(4.25), Inches(0.74),
             border_clr=ec, border_w=1.0, fill=bg, accent=ec)
        sans(sl, title, Inches(0.78), iy+Inches(0.06), Inches(4.00), Inches(0.26),
             sz=11, bold=True, clr=ec)
        sans(sl, detail, Inches(0.78), iy+Inches(0.34), Inches(4.00), Inches(0.36),
             sz=9.5, clr=C_MUTED)

    # Right: 차별성
    card(sl, Inches(5.12), Inches(1.08), Inches(4.42), Inches(4.38),
         border_clr=C_GREEN, border_w=1.4, fill=C_WHITE)
    rect(sl, Inches(5.12), Inches(1.08), Inches(4.42), Inches(0.07), fill=C_GREEN)
    sans(sl, "MoNa-π 차별성",
         Inches(5.26), Inches(1.20), Inches(4.15), Inches(0.34),
         sz=14, bold=True, clr=C_GREEN)

    diffs = [
        (C_GREEN,  C_GREEN_T,  "① 연속 Flow Matching",
         "L = ||v_θ − (x₁−x₀)||²\n5스텝 ODE → 실시간 연속 궤적"),
        (C_NAVY,   C_BLUE_T,   "② Action Chunking  h=10",
         "10스텝 동시 예측 → 버퍼 실행\n4Hz 재계획 / 50Hz 로컬 제어"),
        (C_ORANGE, C_AMBER_T,  "③ Instruction Pool",
         "15개 Paraphrase 무작위 선택\n다양한 자연어 명령 강건 처리"),
        (C_PURPLE, C_PURPLE_T, "④ PaliGemma 3B 백본",
         "SigLIP + Gemma-2B BF16\n언어 이해 + 시각 통합 처리"),
    ]
    for i, (ec, bg, title, detail) in enumerate(diffs):
        iy = Inches(1.64 + i * 0.94)
        card(sl, Inches(5.26), iy, Inches(4.15), Inches(0.84),
             border_clr=ec, border_w=1.0, fill=bg, accent=ec)
        sans(sl, title, Inches(5.42), iy+Inches(0.07), Inches(3.92), Inches(0.28),
             sz=12, bold=True, clr=ec)
        sans(sl, detail, Inches(5.42), iy+Inches(0.38), Inches(3.92), Inches(0.40),
             sz=10, clr=C_MUTED)


def s04_process(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "개발 과정",
        "데이터 수집 → 모델 학습 → 검증 → Serbot2 배포  |  D1–D14")

    # 4-step horizontal flow
    steps = [
        ("01", "데이터 수집",
         "HDF5 파이프라인\n비동기 수집\n9 카테고리",
         C_GREEN, C_GREEN_T),
        ("02", "모델 학습",
         "v2 → v3 → v3-A\nBF16 정밀도\nEpisode Split",
         C_NAVY, C_BLUE_T),
        ("03", "검증 · 분석",
         "FPE + Success@1.0\nClosed-Loop 평가\n어블레이션",
         C_ORANGE, C_AMBER_T),
        ("04", "Serbot2 배포",
         "GX10 FastAPI\n243ms 검증\n실환경 주행",
         C_PURPLE, C_PURPLE_T),
    ]

    sw = Inches(2.08); sh = Inches(2.40)
    xs = [Inches(0.469 + i*(2.08+0.26)) for i in range(4)]
    y0 = Inches(1.30)

    for i, (num, title, body, ec, bg) in enumerate(steps):
        x = xs[i]
        card(sl, x, y0, sw, sh, border_clr=ec, border_w=1.4, fill=bg)
        rect(sl, x, y0, sw, Inches(0.07), fill=ec)

        # Circle number
        rect(sl, x+Inches(0.76), y0+Inches(0.12), Inches(0.56), Inches(0.56), fill=ec)
        sans(sl, num, x+Inches(0.76), y0+Inches(0.12), Inches(0.56), Inches(0.56),
             sz=16, bold=True, clr=C_WHITE, align=CENTER, wrap=False)

        sans(sl, title, x+Inches(0.12), y0+Inches(0.80), sw-Inches(0.24), Inches(0.34),
             sz=13, bold=True, clr=ec, align=CENTER)
        sans(sl, body,  x+Inches(0.12), y0+Inches(1.22), sw-Inches(0.24), Inches(1.00),
             sz=11, clr=C_TEXT, align=CENTER)

        if i < 3:
            sans(sl, "→",
                 x+sw+Inches(0.04), y0+Inches(0.98), Inches(0.18), Inches(0.44),
                 sz=22, bold=True, clr=C_MUTED, align=CENTER, wrap=False)

    # Current status bar
    card(sl, Inches(0.469), Inches(3.84), Inches(9.06), Inches(0.52),
         border_clr=C_GREEN, border_w=1.2, fill=C_GREEN_T, accent=C_GREEN)
    sans(sl, "현재 진행",
         Inches(0.62), Inches(3.92), Inches(1.20), Inches(0.28),
         sz=11, bold=True, clr=C_GREEN)
    sans(sl, "01~03 완료  ·  Success@1.0 86.7% 달성  ·  04 Serbot2 실배포 D10~D11 예정",
         Inches(1.90), Inches(3.92), Inches(7.50), Inches(0.28),
         sz=11, clr=C_TEXT)

    # Team roles compact
    card(sl, Inches(0.469), Inches(4.48), Inches(9.06), Inches(0.96),
         border_clr=C_BORDER, fill=C_WHITE)
    sans(sl, "역할 분담",
         Inches(0.62), Inches(4.56), Inches(1.20), Inches(0.28),
         sz=11, bold=True, clr=C_NAVY)
    rect(sl, Inches(1.90), Inches(4.58), Inches(0.016), Inches(0.80), fill=C_BORDER)

    roles = [
        ("임현우", "모델 아키텍처 · 학습 파이프라인 · 검증",  C_NAVY,   Inches(2.00)),
        ("정재현", "HDF5 수집 파이프라인 · 데이터 균형",       C_GREEN,  Inches(5.10)),
        ("오은석", "데이터셋 수집 · Serbot2 인프라 · ROS2",    C_ORANGE, Inches(7.56)),
    ]
    for name, role, ec, rx in roles:
        sans(sl, name, rx, Inches(4.56), Inches(0.82), Inches(0.28),
             sz=11, bold=True, clr=ec)
        sans(sl, role, rx, Inches(4.86), Inches(2.30), Inches(0.26),
             sz=9.5, clr=C_MUTED)


def s05_architecture(prs, buf_hw):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "시스템 아키텍처",
        "Serbot2 하드웨어 · PaliGemma 3B + Action Expert · 이중 제어 루프")

    # Left: annotated hardware
    sl.shapes.add_picture(buf_hw,
                          Inches(0.25), Inches(1.06), Inches(3.80), Inches(4.42))

    # Right: pipeline
    rx = Inches(4.22)
    rw = Inches(5.31)

    # Input layer
    card(sl, rx, Inches(1.10), rw, Inches(0.68),
         border_clr=C_GREEN, border_w=1.2, fill=C_GREEN_T, accent=C_GREEN)
    sans(sl, "입력  —  Fish-eye Camera  (8-frame · 0.8s window)",
         rx+Inches(0.18), Inches(1.20), rw-Inches(0.24), Inches(0.28),
         sz=11, bold=True, clr=C_GREEN)
    mono(sl, "(B, 8, 224, 224, 3)  +  자연어 명령",
         rx+Inches(0.18), Inches(1.48), rw-Inches(0.24), Inches(0.24),
         sz=10, clr=C_MUTED)

    sans(sl, "↓", rx+rw/2-Inches(0.10), Inches(1.78), Inches(0.20), Inches(0.28),
         sz=16, bold=True, clr=C_MUTED, align=CENTER, wrap=False)

    # GX10 inference
    card(sl, rx, Inches(2.10), rw, Inches(1.36),
         border_clr=C_NAVY, border_w=1.2, fill=C_BLUE_T, accent=C_NAVY)
    sans(sl, "GX10 서버  —  VLA 추론  (4Hz)",
         rx+Inches(0.18), Inches(2.20), rw-Inches(0.24), Inches(0.28),
         sz=12, bold=True, clr=C_NAVY)
    rect(sl, rx+Inches(0.18), Inches(2.50), rw-Inches(0.24), Inches(0.014), fill=C_BORDER)

    components = [
        ("SigLIP SO400M  →  이미지 패치 임베딩",     2.56),
        ("Gemma-2B  →  언어+시각 컨텍스트 (2048)",   2.84),
        ("Action Expert (4-layer Xattn)  →  궤적",  3.12),
        ("Flow Matching ODE (5스텝)  →  (B,10,3)", 3.40),
    ]
    for txt, iy in components:
        mono(sl, txt, rx+Inches(0.18), Inches(iy), rw-Inches(0.24), Inches(0.26),
             sz=9.5, clr=C_TEXT)

    sans(sl, "↓", rx+rw/2-Inches(0.10), Inches(3.48), Inches(0.20), Inches(0.28),
         sz=16, bold=True, clr=C_MUTED, align=CENTER, wrap=False)

    # Serbot2 execution
    card(sl, rx, Inches(3.80), rw, Inches(0.68),
         border_clr=C_ORANGE, border_w=1.2, fill=C_AMBER_T, accent=C_ORANGE)
    sans(sl, "Serbot2  —  로컬 실행  (50Hz)",
         rx+Inches(0.18), Inches(3.90), rw-Inches(0.24), Inches(0.28),
         sz=12, bold=True, clr=C_ORANGE)
    mono(sl, "Action Buffer  →  ROS2  →  Omni Wheel 3-DOF",
         rx+Inches(0.18), Inches(4.18), rw-Inches(0.24), Inches(0.24),
         sz=10, clr=C_MUTED)

    # Key insight
    card(sl, rx, Inches(4.60), rw, Inches(0.76),
         border_clr=C_PURPLE, border_w=1.2, fill=C_PURPLE_T, accent=C_PURPLE)
    sans(sl, "Action Chunking 핵심",
         rx+Inches(0.18), Inches(4.68), Inches(1.80), Inches(0.24),
         sz=10, bold=True, clr=C_PURPLE)
    sans(sl, "GX10이 250ms 추론하는 동안 Serbot2는 이전 청크로 계속 50Hz 실행 → 끊김 없는 주행",
         rx+Inches(0.18), Inches(4.92), rw-Inches(0.24), Inches(0.38),
         sz=10, clr=C_TEXT)


def s06_experiment(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "실험 환경 및 평가",
        "9-카테고리 내비게이션 · Serbot2 실내 주행 · FPE + Closed-Loop Success 평가")

    # Left: robot front photo (navigation scene)
    sl.shapes.add_picture(str(ROBOT_FRONT),
                          Inches(0.469), Inches(1.10), Inches(4.50), Inches(3.60))
    card(sl, Inches(0.469), Inches(4.68), Inches(4.50), Inches(0.38),
         border_clr=C_NAVY, fill=C_NAVY)
    sans(sl, "Serbot2 실험 환경  ·  실내 자율 내비게이션",
         Inches(0.55), Inches(4.74), Inches(4.35), Inches(0.26),
         sz=9.5, bold=True, clr=C_WHITE, align=CENTER)

    # Right top: 9 categories grid
    card(sl, Inches(5.10), Inches(1.10), Inches(4.43), Inches(2.20),
         border_clr=C_BORDER, fill=C_WHITE)
    sans(sl, "9-카테고리 내비게이션 (3-DOF)",
         Inches(5.26), Inches(1.20), Inches(4.10), Inches(0.30),
         sz=12, bold=True, clr=C_NAVY)
    rect(sl, Inches(5.26), Inches(1.52), Inches(4.10), Inches(0.014), fill=C_BORDER)

    cats = [
        ("center_straight", "center_left", "center_right"),
        ("left_straight",   "left_left",   "left_right"),
        ("right_straight",  "right_left",  "right_right"),
    ]
    for ri, row in enumerate(cats):
        for ci, cat in enumerate(row):
            cx = Inches(5.26 + ci * 1.36)
            cy = Inches(1.60 + ri * 0.52)
            is_fail = cat in ("center_left", "left_right")
            bg = C_RED_T if is_fail else C_GREEN_T
            ec = C_RED   if is_fail else C_GREEN
            rect(sl, cx, cy, Inches(1.30), Inches(0.44), fill=bg)
            sans(sl, cat.replace("_", "\n"),
                 cx+Inches(0.04), cy+Inches(0.02), Inches(1.22), Inches(0.40),
                 sz=8, clr=ec, align=CENTER)

    # Right bottom: stats
    card(sl, Inches(5.10), Inches(3.42), Inches(4.43), Inches(1.64),
         border_clr=C_BORDER, fill=C_WHITE)
    sans(sl, "데이터셋 및 평가",
         Inches(5.26), Inches(3.52), Inches(4.10), Inches(0.30),
         sz=12, bold=True, clr=C_NAVY)
    rect(sl, Inches(5.26), Inches(3.84), Inches(4.10), Inches(0.014), fill=C_BORDER)

    stats = [
        ("액션 공간", "linear_x · linear_y · angular_z  (3-DOF)"),
        ("학습 에피소드", "카테고리당 15~30개  ·  HDF5 저장"),
        ("평가 지표",  "FPE (오프라인)  +  Success@1.0 (폐루프)"),
        ("성공 기준",  "목표 위치 1.0m 이내 도달"),
    ]
    for i, (k, v) in enumerate(stats):
        iy = Inches(3.92 + i * 0.27)
        sans(sl, k, Inches(5.26), iy, Inches(1.30), Inches(0.24),
             sz=10, bold=True, clr=C_MUTED)
        sans(sl, v, Inches(6.62), iy, Inches(2.80), Inches(0.24),
             sz=10, clr=C_TEXT)


def s07_results(prs, buf_cmp):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "실험 결과",
        "v3-A: Success@1.0 86.7%  ·  추론 243ms  ·  4Hz 실시간 제어 달성")

    # Top 3 metric cards
    metrics = [
        ("86.7%",  "Success@1.0",     "v3-A  Threshold=1.0",     C_ORANGE, C_AMBER_T),
        ("243 ms", "추론 레이턴시",    "4Hz 예산(250ms) 이내",      C_GREEN,  C_GREEN_T),
        ("5.97 GB","모델 VRAM",        "BF16 / Serbot2 16GB 여유", C_NAVY,   C_BLUE_T),
    ]
    for i, (val, label, sub, ec, bg) in enumerate(metrics):
        mx = Inches(0.469 + i * 3.02)
        card(sl, mx, Inches(1.10), Inches(2.86), Inches(1.20),
             border_clr=ec, border_w=1.4, fill=bg)
        sans(sl, val,   mx+Inches(0.14), Inches(1.18), Inches(2.58), Inches(0.60),
             sz=32, bold=True, clr=ec, align=CENTER)
        sans(sl, label, mx+Inches(0.14), Inches(1.78), Inches(2.58), Inches(0.26),
             sz=11, bold=True, clr=ec, align=CENTER)
        sans(sl, sub,   mx+Inches(0.14), Inches(2.04), Inches(2.58), Inches(0.22),
             sz=9, clr=C_MUTED, align=CENTER)

    # Bottom left: comparison chart
    sl.shapes.add_picture(buf_cmp,
                          Inches(0.469), Inches(2.44), Inches(5.00), Inches(3.00))

    # Bottom right: key findings
    card(sl, Inches(5.58), Inches(2.44), Inches(3.95), Inches(3.00),
         border_clr=C_BORDER, fill=C_WHITE)
    sans(sl, "주요 발견",
         Inches(5.74), Inches(2.54), Inches(3.62), Inches(0.30),
         sz=12, bold=True, clr=C_NAVY)
    rect(sl, Inches(5.74), Inches(2.86), Inches(3.62), Inches(0.014), fill=C_BORDER)

    findings = [
        (C_GREEN,  "[개선] Episode Split + Instr. Pool",
                   "→ +9.7%p Success 향상"),
        (C_RED,    "[한계] center_left  0% 실패",
                   "→ 시각 모호성, 추가 데이터 필요"),
        (C_ORANGE, "[확정] BF16 필수",
                   "→ FP16 시 FPE 3.1× 증가, 완전 실패"),
        (C_NAVY,   "[달성] 4Hz 실시간",
                   "→ Warm 243ms < 250ms 예산"),
    ]
    for i, (ec, title, detail) in enumerate(findings):
        iy = Inches(2.96 + i * 0.60)
        rect(sl, Inches(5.74), iy, Inches(0.04), Inches(0.50), fill=ec)
        sans(sl, title,  Inches(5.84), iy+Inches(0.02), Inches(3.40), Inches(0.24),
             sz=10.5, bold=True, clr=ec)
        sans(sl, detail, Inches(5.84), iy+Inches(0.26), Inches(3.40), Inches(0.22),
             sz=10, clr=C_MUTED)


def s08_conclusion(prs):
    sl = blank(prs); set_bg(sl)
    hdr(sl, "결론 및 기대 효과",
        "도메인 이전 성공 · 최적화 파이프라인 제시 · 향후 적용 및 파급효과")

    # Two main contributions (large, prominent)
    card(sl, Inches(0.469), Inches(1.10), Inches(9.06), Inches(1.36),
         border_clr=C_NAVY, border_w=1.5, fill=C_BLUE_T, accent=C_NAVY)
    sans(sl, "기대 효과 ①  —  도메인 이전",
         Inches(0.62), Inches(1.20), Inches(3.00), Inches(0.28),
         sz=11, bold=True, clr=C_NAVY)
    sans(sl, "모바일 매니플레이터(로봇 팔) 중심의 VLA를 모바일 로봇 내비게이션으로 도메인 이전 성공",
         Inches(0.62), Inches(1.52), Inches(8.78), Inches(0.52),
         sz=16, bold=True, clr=C_TEXT)

    card(sl, Inches(0.469), Inches(2.58), Inches(9.06), Inches(1.36),
         border_clr=C_GREEN, border_w=1.5, fill=C_GREEN_T, accent=C_GREEN)
    sans(sl, "기대 효과 ②  —  파이프라인 기여",
         Inches(0.62), Inches(2.68), Inches(3.00), Inches(0.28),
         sz=11, bold=True, clr=C_GREEN)
    sans(sl, "적은 데이터셋으로 최적화 가능한 Flow Matching VLA 파이프라인·프레임워크 제시",
         Inches(0.62), Inches(3.00), Inches(8.78), Inches(0.52),
         sz=16, bold=True, clr=C_TEXT)

    # Bottom: future work (compact)
    card(sl, Inches(0.469), Inches(4.06), Inches(9.06), Inches(1.32),
         border_clr=C_BORDER, fill=C_WHITE)
    sans(sl, "추후 계획",
         Inches(0.62), Inches(4.16), Inches(1.50), Inches(0.28),
         sz=12, bold=True, clr=C_MUTED)
    rect(sl, Inches(2.20), Inches(4.18), Inches(0.016), Inches(1.00), fill=C_BORDER)

    future = [
        ("단기",  "D10-D11  Serbot2 실배포 · 실환경 주행 10+ ep · CL Success 측정",  C_ORANGE, Inches(2.30)),
        ("중기",  "H6  center_left 에피소드 보강 · v4 재학습",                        C_NAVY,   Inches(2.30)),
        ("장기",  "실로봇 결과 논문화 · MoNaVLA 비교 · 시연 영상 제작",               C_PURPLE, Inches(2.30)),
    ]
    for i, (term, plan, ec, _) in enumerate(future):
        iy = Inches(4.20 + i * 0.28)
        sans(sl, term + ":",   Inches(2.30), iy, Inches(0.60), Inches(0.26),
             sz=10, bold=True, clr=ec)
        sans(sl, plan, Inches(2.98), iy, Inches(6.40), Inches(0.26),
             sz=10, clr=C_MUTED)


def s09_thanks(prs):
    sl = blank(prs); set_bg(sl)
    rect(sl, 0, 0, SW, Inches(0.22), fill=C_NAVY)
    rect(sl, 0, SH-Inches(0.22), SW, Inches(0.22), fill=C_NAVY)

    # Big thank you
    sans(sl, "감사합니다",
         Inches(0.50), Inches(0.80), Inches(6.50), Inches(1.80),
         sz=72, bold=True, clr=C_NAVY)

    rect(sl, Inches(0.60), Inches(2.72), Inches(3.80), Inches(0.020), fill=C_BORDER)

    sans(sl, "MoNa-π  ·  Mobile Navigation π0",
         Inches(0.60), Inches(2.82), Inches(5.80), Inches(0.38),
         sz=16, bold=True, clr=C_TEXT)
    sans(sl, "Flow Matching 기반 고주파 모바일 내비게이션 VLA",
         Inches(0.60), Inches(3.24), Inches(5.80), Inches(0.30),
         sz=13, clr=C_MUTED)

    # Team
    card(sl, Inches(0.60), Inches(3.70), Inches(5.20), Inches(1.60),
         border_clr=C_BORDER, fill=C_WHITE)
    sans(sl, "Team MONA-F",
         Inches(0.78), Inches(3.80), Inches(4.80), Inches(0.30),
         sz=12, bold=True, clr=C_NAVY)
    rect(sl, Inches(0.78), Inches(4.12), Inches(4.80), Inches(0.015), fill=C_BORDER)

    members = [
        ("임현우 (팀장)", "모델 아키텍처 · 학습 · 검증",      4.20),
        ("정재연",        "데이터셋 수집 파이프라인 구축",     4.50),
        ("오은석",        "데이터셋 수집 · Serbot2 인프라",   4.80),
    ]
    for name, role, iy in members:
        sans(sl, name, Inches(0.78), Inches(iy), Inches(1.10), Inches(0.28),
             sz=11, bold=True, clr=C_TEXT)
        sans(sl, role, Inches(1.94), Inches(iy), Inches(3.65), Inches(0.28),
             sz=10.5, clr=C_MUTED)

    # Robot front photo right side
    sl.shapes.add_picture(str(ROBOT_FRONT),
                          Inches(6.20), Inches(0.50), Inches(3.55), Inches(4.84))

    # Caption on photo
    card(sl, Inches(6.20), Inches(5.09), Inches(3.55), Inches(0.30),
         border_clr=C_NAVY, fill=C_NAVY)
    sans(sl, "Serbot2  ·  자율 내비게이션 시연",
         Inches(6.28), Inches(5.13), Inches(3.39), Inches(0.22),
         sz=8.5, bold=True, clr=C_WHITE, align=CENTER)


# ── Entry ─────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("[1/3] Annotating hardware photo...")
    buf_hw  = annotated_robot_hw()
    print("[2/3] Generating comparison chart...")
    buf_cmp = chart_comparison()

    print("[3/3] Building 9 slides...")
    prs = new_prs()
    s01_title(prs)            # 01 표지
    s02_toc(prs)              # 02 목차
    s03_motivation(prs)       # 03 기획의도 + 차별성
    s04_process(prs)          # 04 과제 추진
    s05_architecture(prs, buf_hw)  # 05 시스템 아키텍처
    s06_experiment(prs)       # 06 실험 환경
    s07_results(prs, buf_cmp) # 07 핵심 결과
    s08_conclusion(prs)       # 08 기여 및 결론
    s09_thanks(prs)           # 09 감사합니다

    out = "/home/minum/26CS/MoNa-pi/reports/mona_pi_v5.pptx"
    prs.save(out)
    print(f"\nSaved: {out}  ({len(prs.slides)} slides)")